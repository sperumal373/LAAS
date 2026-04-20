from pptx import Presentation
from pptx.util import Pt

prs = Presentation(r'C:\Users\Administrator\Desktop\Wipro Template.pptx')

for i, layout in enumerate(prs.slide_layouts):
    print(f'\nLayout {i} [{layout.name}]:')
    for ph in layout.placeholders:
        print(f'  ph idx={ph.placeholder_format.idx} name="{ph.name}"')

# Sample first 5 existing slides
print('\n\n=== FIRST 5 EXISTING SLIDES ===')
for i, slide in enumerate(prs.slides[:5]):
    print(f'\nSlide {i+1}:')
    for ph in slide.placeholders:
        txt = ph.text[:80].replace('\n','|') if ph.has_text_frame else '[no text]'
        print(f'  ph idx={ph.placeholder_format.idx} name="{ph.name}" text="{txt}"')
    print(f'  layout: {slide.slide_layout.name}')

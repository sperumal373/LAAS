"""Analyze PPT template structure in detail."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'C:\caas-dashboard\LaaS_Portal_Presentation_v63.pptx')

for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name if slide.slide_layout else 'N/A'
    print(f'=== Slide {i+1} (Layout: {layout}) ===')
    has_image = False
    for j, shape in enumerate(slide.shapes):
        stype = str(shape.shape_type)
        left_in = round(shape.left / 914400, 2)
        top_in = round(shape.top / 914400, 2)
        w_in = round(shape.width / 914400, 2)
        h_in = round(shape.height / 914400, 2)
        print(f'  Shape {j}: type={stype}, name="{shape.name}", pos=({left_in},{top_in}), size=({w_in}x{h_in})')
        if shape.has_text_frame:
            tf = shape.text_frame
            for pi, para in enumerate(tf.paragraphs):
                text = para.text.strip()
                if text and pi < 4:
                    font = para.runs[0].font if para.runs else None
                    fsize = str(Pt(font.size.pt).pt) if font and font.size else 'inherit'
                    fname = font.name if font and font.name else 'inherit'
                    fbold = font.bold if font else None
                    try:
                        fcolor = str(font.color.rgb) if font and font.color and font.color.rgb else 'inherit'
                    except:
                        fcolor = 'inherit'
                    print(f'    p[{pi}]: "{text[:70]}" | {fname} {fsize}pt bold={fbold} color={fcolor}')
        try:
            if hasattr(shape, 'image') and shape.image:
                has_image = True
                print(f'    IMAGE: {shape.image.content_type}')
        except:
            pass
    print(f'  >> has_image={has_image}')
    print()

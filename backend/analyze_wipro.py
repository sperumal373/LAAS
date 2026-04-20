"""Analyze Wipro Template layouts in detail."""
from pptx import Presentation

prs = Presentation(r"C:\Users\Administrator\Desktop\Wipro Template.pptx")

print(f"Slide width: {prs.slide_width}  height: {prs.slide_height}")
print(f"Total layouts: {len(prs.slide_layouts)}")
print(f"Total slides: {len(prs.slides)}")

for i, layout in enumerate(prs.slide_layouts):
    phs = list(layout.placeholders)
    print(f"\n--- Layout {i}: '{layout.name}' ({len(phs)} placeholders) ---")
    for ph in phs:
        txt = ""
        if ph.has_text_frame and ph.text_frame.text.strip():
            txt = f" text='{ph.text_frame.text[:60]}'"
        print(f"  PH{ph.placeholder_format.idx}: {ph.name} | L={ph.left} T={ph.top} W={ph.width} H={ph.height}{txt}")

print("\n\n=== FIRST 5 SLIDES ===")
for si in range(min(5, len(prs.slides))):
    slide = list(prs.slides)[si]
    print(f"\nSlide {si+1}: layout='{slide.slide_layout.name}' | {len(slide.shapes)} shapes")
    for sh in slide.shapes:
        txt = ""
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = f" | '{sh.text_frame.text[:80]}'"
        print(f"  {sh.name} | L={sh.left} T={sh.top} W={sh.width} H={sh.height}{txt}")

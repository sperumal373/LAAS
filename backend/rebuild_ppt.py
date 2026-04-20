"""
Rebuild LaaS Portal PPT with fresh screenshots on ALL content slides.
- Slides 1-4: Keep as-is (cover, agenda, platform overview, tech stack)
- Slides 5-17: Already styled (update screenshot images only)
- Slides 18-22: Currently plain text → convert to styled layout with screenshots
- Slides 23-24: Already styled (update screenshot images)
- Slides 25-27: Currently plain text → convert to styled layout with screenshots
- Slide 28: Keep as-is (closing cover)
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

SRC  = r"C:\caas-dashboard\LaaS_Portal_Presentation_v63.pptx"
OUT  = r"C:\caas-dashboard\LaaS_Portal_Presentation_v64.pptx"
SDIR = r"C:\caas-dashboard\screenshots"

# ── Styling constants (extracted from existing styled slides 5-17) ──
BLUE_BAR   = RGBColor(0x1F, 0x39, 0x64)
TITLE_CLR  = RGBColor(0x1F, 0x39, 0x64)
SUBTITLE_CLR = RGBColor(0x59, 0x59, 0x59)
BULLET_CLR = RGBColor(0x00, 0x70, 0xC0)

# Positions (EMU) — from styled slides
BAR_L, BAR_T, BAR_W, BAR_H = 347472, 182880, 11494008, 50292
TITLE_L, TITLE_T, TITLE_W, TITLE_H = 347472, 246888, 11494008, 548640
SUB_L, SUB_T, SUB_W, SUB_H = 347472, 795528, 11494008, 292608
RECT_L, RECT_T, RECT_W, RECT_H = 320040, 1143000, 7918704, 4478274
IMG_L, IMG_T, IMG_W, IMG_H = 347472, 1170432, 7863840, 4423410
BULLET_L, BULLET_T, BULLET_W, BULLET_H = 8439912, 1170432, 3566160, 5047488

# ── Slide-to-screenshot mapping ──
# key = slide index (0-based), value = (screenshot_file, title, subtitle, bullets)
SLIDE_MAP = {
    # --- Already styled (slides 5-17, index 4-16) → just update image ---
    4:  ("01_overview.png", None, None, None),
    5:  ("02_vmware_vms.png", None, None, None),
    6:  ("03_snapshots.png", None, None, None),
    7:  ("04_networks.png", None, None, None),
    8:  ("05_capacity.png", None, None, None),
    9:  ("06_projects.png", None, None, None),
    10: ("07_chargeback.png", None, None, None),
    11: ("08_requests.png", None, None, None),
    12: ("09_ipam.png", None, None, None),
    13: ("10_assets.png", None, None, None),
    14: ("11_openshift.png", None, None, None),
    15: ("12_nutanix.png", None, None, None),
    16: ("13_ansible.png", None, None, None),

    # --- Plain text → convert to styled (slides 18-22, index 17-21) ---
    17: ("14_aws.png",
         "Amazon Web Services (AWS)",
         "EC2, S3, RDS, VPC management and Cost Explorer integration",
         [
             "☁️  EC2 instance discovery across all configured AWS accounts",
             "⚡  Start, Stop and Reboot EC2 instances directly from the portal",
             "📦  S3 bucket listing with region and creation date",
             "🗄️  RDS database engine, status and endpoint visibility",
             "🌐  VPC inventory with subnets, route tables and security groups",
             "💰  Cost Explorer: per-service daily spend with interactive charts",
             "🔑  IAM integration via STS assume-role for multi-account access",
             "📊  Per-account tab navigation with instance count badges",
         ]),
    18: ("15_hyperv.png",
         "Microsoft Hyper-V",
         "Standalone Hyper-V host management via WinRM PowerShell remoting",
         [
             "🪟  Connect to standalone Hyper-V hosts — no SCVMM required",
             "💻  VM inventory with CPU, memory, state and uptime",
             "⚡  Power actions: Start, Stop, Restart, Checkpoint",
             "📊  Per-host CPU and RAM usage bar charts",
             "🔌  Agentless — uses WinRM/PowerShell remoting only",
             "🗂️  Multi-host tab navigation with status badges",
         ]),
    19: ("16_storage.png",
         "Enterprise Storage Management",
         "Pure Storage, Dell EMC, HPE and NetApp — all via native REST APIs",
         [
             "💾  Pure Storage FlashArray & FlashBlade: volumes, hosts, capacity",
             "🗄️  Dell EMC PowerStore: volumes, host groups, storage pools",
             "📦  HPE StoreServ: CPGs, virtual volumes, port inventory",
             "🌊  NetApp ONTAP: aggregates, volumes, LIFs, cluster health",
             "📊  Per-array capacity bar charts with threshold alerts",
             "🔗  Native REST API integration — no gateway or agent needed",
             "📈  Multi-vendor normalised view with unified KPI cards",
         ]),
    20: ("17_backup.png",
         "Backup & Data Protection",
         "Rubrik Security Cloud, Cohesity DataProtect and Veeam Backup & Replication",
         [
             "🛡️  Rubrik: SLA domains, on-demand snapshots, live mount, restore",
             "🗄️  Cohesity: protection groups, run history, compliance status",
             "💿  Veeam: backup jobs, repository capacity, session logs",
             "📊  Per-platform SLA compliance dashboards",
             "⏱️  Backup job trend charts: duration, size, success rate",
             "🔄  Instant recovery and file-level restore from portal",
             "📋  Unified backup inventory across all three vendors",
         ]),
    21: ("18_cmdb.png",
         "CMDB — Configuration Management Database",
         "ServiceNow-aligned CI registry with auto-discovery across all platforms",
         [
             "🗄️  1,800+ CIs auto-discovered from VMware, AWS, OCP, Nutanix, etc.",
             "🏷️  Tagging auto-populated from vCenter tags per VM",
             "🏢  Department set to SDx-COE across all CI types",
             "🖥️  ESXi host IP and OS version (ESXi 8.0.x) populated",
             "📊  Columns: Server, IP, Status, Technology, OS, Manager, Owner",
             "🌍  Environment, Timezone, Region, Department, Tagging fields",
             "📥  Export to CSV and PDF with one click",
             "✏️  Inline editing — click any cell to update CI attributes",
         ]),

    # --- Already styled (slides 23-24, index 22-23) → just update image ---
    22: ("19_ad_dns.png", None, None, None),
    # Slide 24 is special (dual image layout) — handle separately
    23: None,  # skip, handled specially

    # --- Plain text → convert to styled (slides 25-27, index 24-26) ---
    24: ("22_insights.png",
         "Insights & Analytics Engine",
         "PostgreSQL 16 daily snapshots — Health Scorecard, Change Detection, Capacity & Cost",
         [
             "🩺  Health Scorecard: real-time RAG status across all platforms",
             "🔄  Change Detection: new, removed, modified CIs day over day",
             "📊  Capacity & Cost: per-platform resource utilisation and spend",
             "📋  Executive Summary: auto-generated digest for leadership",
             "📈  Trend arrows and sparklines for key metrics",
             "🗓️  Daily snapshots stored in PostgreSQL 16 for historical queries",
             "🔍  Drill-down from summary cards to detailed tables",
         ]),
    25: ("23_history.png",
         "Historical Trending",
         "Interactive charts for 7, 14 and 30-day windows across all collected metrics",
         [
             "📈  VM count, host count, datastore usage trend lines",
             "🖥️  Per-vCenter historical comparison charts",
             "🗓️  Selectable time windows: 7 / 14 / 30 days",
             "📊  Overlay charts: CPU vs Memory utilisation over time",
             "💾  Datastore growth trend for capacity planning",
             "🔴  OCP node count and pod trends",
             "📉  Anomaly highlights when metrics deviate from baseline",
         ]),
    26: ("24_forecast.png",
         "Capacity Forecasting",
         "Linear regression predictions for resource planning across all platforms",
         [
             "🔮  Linear regression forecast for CPU, memory, storage",
             "📅  Projected exhaustion date per resource pool",
             "📊  Confidence bands on forecast charts",
             "🖥️  Per-vCenter and per-cluster forecast breakdown",
             "⚠️  Threshold alerts when forecasted usage > 80%",
             "📈  30, 60 and 90-day forward projections",
             "📋  Exportable forecast report for capacity reviews",
         ]),
}


def _clear_shapes(slide):
    """Remove all shapes from a slide."""
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        if sp.tag.endswith('}sp') or sp.tag.endswith('}pic') or sp.tag.endswith('}grpSp'):
            sp_tree.remove(sp)


def _add_blue_bar(slide):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, BAR_L, BAR_T, BAR_W, BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_BAR
    bar.line.fill.background()
    return bar


def _add_title(slide, text):
    txbox = slide.shapes.add_textbox(TITLE_L, TITLE_T, TITLE_W, TITLE_H)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TITLE_CLR
    return txbox


def _add_subtitle(slide, text):
    txbox = slide.shapes.add_textbox(SUB_L, SUB_T, SUB_W, SUB_H)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = SUBTITLE_CLR
    return txbox


def _add_screenshot_rect(slide):
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, RECT_L, RECT_T, RECT_W, RECT_H)
    rect.fill.background()
    rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    rect.line.width = Pt(0.75)
    return rect


def _add_screenshot(slide, img_path):
    pic = slide.shapes.add_picture(img_path, IMG_L, IMG_T, IMG_W, IMG_H)
    return pic


def _add_bullet_panel(slide, bullets):
    txbox = slide.shapes.add_textbox(BULLET_L, BULLET_T, BULLET_W, BULLET_H)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet_text
        p.font.size = Pt(11)
        p.font.color.rgb = BULLET_CLR
        p.space_after = Pt(6)

        # Add empty line between bullets for readability
        if i < len(bullets) - 1:
            spacer = tf.add_paragraph()
            spacer.text = ""
            spacer.space_before = Pt(2)
            spacer.space_after = Pt(2)

    return txbox


def _replace_image_in_styled_slide(slide, img_path):
    """Find the existing Picture shape and replace it with new image."""
    # Remove old picture(s)
    sp_tree = slide.shapes._spTree
    pics_to_remove = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics_to_remove.append(sh._element)

    for pic_el in pics_to_remove:
        sp_tree.remove(pic_el)

    # Add new picture at same position
    pic = slide.shapes.add_picture(img_path, IMG_L, IMG_T, IMG_W, IMG_H)
    return pic


def _convert_plain_to_styled(slide, img_path, title, subtitle, bullets):
    """Convert a plain-text slide to styled layout with screenshot + bullets."""
    _clear_shapes(slide)
    _add_blue_bar(slide)
    _add_title(slide, title)
    _add_subtitle(slide, subtitle)
    _add_screenshot_rect(slide)
    _add_screenshot(slide, img_path)
    _add_bullet_panel(slide, bullets)


def _handle_slide_24(slide):
    """Slide 24 (Users/Audit) has dual-image layout. Update both images."""
    sp_tree = slide.shapes._spTree
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Remove old pictures
    for pic in pics:
        sp_tree.remove(pic._element)

    # Slide 24 original has two images side by side:
    # Left image (Audit Log): pos=(347472, 1170432) size=(5577840, 3137535)
    # Right image (User Mgmt): pos=(6089904, 1170432) size=(5577840, 3137535)
    audit_img = os.path.join(SDIR, "21_audit.png")
    users_img = os.path.join(SDIR, "20_users.png")

    if os.path.exists(audit_img):
        slide.shapes.add_picture(audit_img, 347472, 1170432, 5577840, 3137535)
    if os.path.exists(users_img):
        slide.shapes.add_picture(users_img, 6089904, 1170432, 5577840, 3137535)


def main():
    print("Loading presentation...")
    prs = Presentation(SRC)
    slides = list(prs.slides)
    print(f"  Total slides: {len(slides)}")

    for idx, slide in enumerate(slides):
        if idx not in SLIDE_MAP:
            print(f"  Slide {idx+1}: SKIP (not mapped)")
            continue

        entry = SLIDE_MAP[idx]
        if entry is None:
            # Special handling for slide 24
            if idx == 23:
                print(f"  Slide {idx+1}: Updating dual-image layout (Users & Audit)")
                _handle_slide_24(slide)
            continue

        img_file, title, subtitle, bullets = entry
        img_path = os.path.join(SDIR, img_file)

        if not os.path.exists(img_path):
            print(f"  Slide {idx+1}: WARNING - {img_file} not found, skipping")
            continue

        if title is None:
            # Already styled slide → just replace the image
            print(f"  Slide {idx+1}: Replacing screenshot with {img_file}")
            _replace_image_in_styled_slide(slide, img_path)
        else:
            # Plain text slide → convert to styled format
            print(f"  Slide {idx+1}: Converting to styled layout with {img_file}")
            _convert_plain_to_styled(slide, img_path, title, subtitle, bullets)

    print(f"\nSaving to: {OUT}")
    prs.save(OUT)
    size_mb = os.path.getsize(OUT) / (1024 * 1024)
    print(f"  File size: {size_mb:.1f} MB")
    print("Done!")


if __name__ == "__main__":
    main()

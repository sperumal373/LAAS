"""
Build a completely new LaaS Portal presentation using the Wipro Template.
Uses Wipro branded layouts for cover slides and blank content layout for custom slides.
All 30 slides rebuilt with professional Wipro styling.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

WIPRO_TPL = r"C:\Users\Administrator\Desktop\Wipro Template.pptx"
EXISTING  = r"C:\caas-dashboard\LaaS_Portal_Presentation_v64.pptx"
SDIR      = r"C:\caas-dashboard\screenshots"
OUT       = r"C:\Users\Administrator\Desktop\LaaS_Portal_Presentation_v65_Wipro.pptx"

# ── Colors ──
WIPRO_BLUE  = RGBColor(0x3F, 0x1D, 0x72)  # Wipro purple-blue
NAVY        = RGBColor(0x1F, 0x39, 0x64)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG     = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x59, 0x59, 0x59)
BODY_TEXT    = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)

# ── Positions ──
BAR_L, BAR_T, BAR_W, BAR_H = 347472, 182880, 11494008, 50292
TITLE_L, TITLE_T, TITLE_W, TITLE_H = 347472, 246888, 11494008, 548640
SUB_L, SUB_T, SUB_W, SUB_H = 347472, 795528, 11494008, 292608
RECT_L, RECT_T, RECT_W, RECT_H = 320040, 1143000, 7918704, 4478274
IMG_L, IMG_T, IMG_W, IMG_H = 347472, 1170432, 7863840, 4423410
BULLET_L, BULLET_T, BULLET_W, BULLET_H = 8439912, 1170432, 3566160, 5047488


# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════

def _clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag in ('sp', 'pic', 'grpSp', 'cxnSp'):
            sp_tree.remove(sp)


def _add_bar(slide, color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, BAR_L, BAR_T, BAR_W, BAR_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_title(slide, text, color=NAVY):
    tx = slide.shapes.add_textbox(TITLE_L, TITLE_T, TITLE_W, TITLE_H)
    tf = tx.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = color


def _add_subtitle(slide, text):
    tx = slide.shapes.add_textbox(SUB_L, SUB_T, SUB_W, SUB_H)
    tf = tx.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(12); r.font.color.rgb = DARK_GRAY


def _add_header(slide, title, subtitle, bar_color=ACCENT_BLUE):
    _add_bar(slide, bar_color)
    _add_title(slide, title)
    _add_subtitle(slide, subtitle)


def _add_screenshot(slide, img_file):
    img_path = os.path.join(SDIR, img_file)
    if not os.path.exists(img_path):
        print(f"    WARNING: {img_file} not found")
        return
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, RECT_L, RECT_T, RECT_W, RECT_H)
    rect.fill.background()
    rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); rect.line.width = Pt(0.75)
    slide.shapes.add_picture(img_path, IMG_L, IMG_T, IMG_W, IMG_H)


def _add_bullets(slide, bullets, color=ACCENT_BLUE):
    tx = slide.shapes.add_textbox(BULLET_L, BULLET_T, BULLET_W, BULLET_H)
    tf = tx.text_frame; tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = text
        r.font.size = Pt(11); r.font.color.rgb = color
        p.space_after = Pt(4)
        if i < len(bullets) - 1:
            sp = tf.add_paragraph(); sp.text = ""; sp.space_after = Pt(2)


def _add_tile(slide, left, top, width, hdr_h, body_h, hdr_color, title, bullets):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, hdr_h)
    h.fill.solid(); h.fill.fore_color.rgb = hdr_color; h.line.fill.background()
    ht = slide.shapes.add_textbox(left + Emu(91440), top + Emu(27432), width - Emu(182880), hdr_h - Emu(27432))
    r = ht.text_frame.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE

    bt = top + hdr_h
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, bt, width, body_h)
    b.fill.solid(); b.fill.fore_color.rgb = GRAY_BG; b.line.fill.background()

    y = bt + Emu(73152)
    for txt in bullets:
        tx = slide.shapes.add_textbox(left + Emu(109728), y, width - Emu(182880), Emu(274320))
        r = tx.text_frame.paragraphs[0].add_run(); r.text = txt
        r.font.size = Pt(11); r.font.color.rgb = NAVY
        y += Emu(274320)


def _add_icon_card(slide, x, y, w, h, bg, icon, title, desc, tsz=18):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg; rect.line.fill.background()
    try: rect.adjustments[0] = 0.06
    except: pass

    tx = slide.shapes.add_textbox(x + Emu(137160), y + Emu(91440), w - Emu(274320), Emu(457200))
    tf = tx.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = f"{icon}  {title}"
    r.font.size = Pt(tsz); r.font.bold = True; r.font.color.rgb = WHITE

    if desc:
        tx2 = slide.shapes.add_textbox(x + Emu(137160), y + Emu(548640), w - Emu(274320), h - Emu(640080))
        tf2 = tx2.text_frame; tf2.word_wrap = True
        r2 = tf2.paragraphs[0].add_run(); r2.text = desc
        r2.font.size = Pt(10.5); r2.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)


def _screenshot_slide(slide, title, subtitle, img_file, bullets):
    _add_header(slide, title, subtitle)
    _add_screenshot(slide, img_file)
    _add_bullets(slide, bullets)


# ══════════════════════════════════════════════════════════════
# SLIDE DEFINITIONS
# ══════════════════════════════════════════════════════════════

SCREENSHOT_SLIDES = [
    ("Multi-Platform Dashboard — Overview",
     "Home page: KPIs, platform health rings and IPAM summary",
     "01_overview.png",
     ["🖥️  VMware vCenter KPI cards: VMs, Hosts, Alerts, Storage",
      "🔴  Red Hat OpenShift cluster health ring",
      "🟩  Nutanix AHV cluster & VM count",
      "☁️  AWS EC2 health status and instance count",
      "🪟  Hyper-V host connection status and VM count",
      "💾  Storage array health across all vendors",
      "🛡️  Backup compliance: Rubrik, Cohesity, Veeam",
      "📡  IPAM subnet utilisation summary"]),
    ("VMware vSphere — VM Management",
     "Full VM inventory with per-vCenter tabs, power actions and search",
     "02_vmware_vms.png",
     ["🗂️  Per-vCenter tab navigation with live status",
      "🔍  Search and filter VMs by name, IP, OS, host",
      "⚡  Power actions: Start, Stop, Restart, Suspend",
      "📋  Full VM details: CPU, RAM, disk, IP, tools",
      "📊  Guest OS distribution chart",
      "🔄  Clone, migrate, reconfigure from portal"]),
    ("VM Snapshot Management",
     "Create, restore and delete snapshots with retention visibility",
     "03_snapshots.png",
     ["📸  List all snapshots across all vCenters",
      "✅  Create snapshot with description from portal",
      "↩️  Restore VM to any snapshot point-in-time",
      "🗑️  Delete individual or all snapshots",
      "⏱️  Snapshot age and size tracking",
      "⚠️  Retention alerts for stale snapshots"]),
    ("Network Inventory",
     "Port group and network inventory across all vCenters",
     "04_networks.png",
     ["🌐  All port groups listed per vCenter",
      "🏷️  Network name, type (Standard/DVS) visible",
      "📡  Associated VLAN IDs and switch info",
      "📊  VM count per network",
      "🔍  Search across all networks"]),
    ("Capacity Planning & Resource Monitoring",
     "ESXi host and datastore capacity across vCenters + OCP + Nutanix",
     "05_capacity.png",
     ["🗄️  Per-ESXi host: CPU %, Memory %, connection status",
      "💾  Per-datastore: capacity, used, free GB",
      "📊  Colour-coded usage bars with thresholds",
      "🔴  OCP node CPU/memory utilisation",
      "🟩  Nutanix cluster resource overview"]),
    ("Project Utilisation Tracking",
     "Per-project VM resource consumption by vCenter",
     "06_projects.png",
     ["📊  Bar charts of vCPU and memory per project",
      "🏷️  VMs tagged by project code in vSphere",
      "👤  Project owner attribution",
      "📈  Trend analysis per project"]),
    ("Chargeback & Cost Management",
     "Cost rate configuration per platform — VMware, OpenShift, Nutanix",
     "07_chargeback.png",
     ["💳  Per-platform pricing tabs: VMware / OCP / Nutanix",
      "⚙️  Admin-configurable rates: ₹/vCPU·hr, ₹/GB·hr",
      "📊  Monthly cost projection per project",
      "📋  Exportable chargeback reports"]),
    ("VM Request & Approval Workflow",
     "Self-service VM provisioning with manager approval queue",
     "08_requests.png",
     ["📋  Request form: platform, config, project, lease",
      "✅  Pending approval queue for managers",
      "🔍  Search and filter requests by status",
      "📊  Request trend analytics"]),
    ("IPAM — IP Address Management",
     "Subnet tracking and per-subnet IP usage visibility",
     "09_ipam.png",
     ["📡  List all subnets with prefix, VLAN, usage %",
      "🔍  Drill into any subnet to see all IPs",
      "📊  Per-IP status: used, free, reserved",
      "📥  Export subnet data to CSV"]),
    ("Asset Inventory Management",
     "Physical and logical asset tracking with ping and power actions",
     "10_assets.png",
     ["🖥️  Full asset list: name, type, IP, location, status",
      "🏓  Ping selected assets for reachability check",
      "⚡  Power actions on physical assets (IPMI/iLO)",
      "📥  Export asset inventory to CSV/Excel"]),
    ("Red Hat OpenShift Container Platform",
     "Multi-cluster OCP overview — nodes, operators, namespaces, workloads",
     "11_openshift.png",
     ["🔴  Multi-cluster tab navigation",
      "📊  Cluster health: nodes, CPU %, memory %",
      "📦  Operator status: installed, degraded, available",
      "🐳  Namespace and pod inventory"]),
    ("Nutanix AHV Virtualisation",
     "Prism Central integration — VMs, clusters, images",
     "12_nutanix.png",
     ["🟩  Multi Prism Central tab navigation",
      "📊  AHV VM inventory: name, power state, CPU/RAM",
      "🗄️  Cluster resource overview with capacity bars",
      "🖼️  Image management and upload"]),
    ("Ansible Automation Platform (AAP)",
     "Job template and playbook execution via AAP API",
     "13_ansible.png",
     ["🤖  Job template listing from AAP",
      "▶️  Launch job templates with extra vars",
      "📊  Job status: pending, running, successful, failed",
      "📋  Job output log viewer"]),
    ("Amazon Web Services (AWS)",
     "EC2, S3, RDS, VPC management and Cost Explorer integration",
     "14_aws.png",
     ["☁️  EC2 instance discovery across all configured AWS accounts",
      "⚡  Start, Stop and Reboot EC2 instances from portal",
      "📦  S3 bucket listing with region and creation date",
      "🗄️  RDS database engine, status and endpoint visibility",
      "🌐  VPC inventory with subnets and security groups",
      "💰  Cost Explorer: per-service daily spend charts",
      "🔑  IAM integration via STS assume-role",
      "📊  Per-account tab navigation with instance count"]),
    ("Microsoft Hyper-V",
     "Standalone Hyper-V host management via WinRM PowerShell remoting",
     "15_hyperv.png",
     ["🪟  Connect to standalone Hyper-V hosts — no SCVMM required",
      "💻  VM inventory with CPU, memory, state and uptime",
      "⚡  Power actions: Start, Stop, Restart, Checkpoint",
      "📊  Per-host CPU and RAM usage bar charts",
      "🔌  Agentless — uses WinRM/PowerShell remoting only",
      "🗂️  Multi-host tab navigation with status badges"]),
    ("Enterprise Storage Management",
     "Pure Storage, Dell EMC, HPE and NetApp — all via native REST APIs",
     "16_storage.png",
     ["💾  Pure Storage FlashArray & FlashBlade: volumes, hosts, capacity",
      "🗄️  Dell EMC PowerStore: volumes, host groups, storage pools",
      "📦  HPE StoreServ: CPGs, virtual volumes, port inventory",
      "🌊  NetApp ONTAP: aggregates, volumes, LIFs, cluster health",
      "📊  Per-array capacity bar charts with threshold alerts",
      "🔗  Native REST API — no gateway or agent needed"]),
    ("Backup & Data Protection",
     "Rubrik Security Cloud, Cohesity DataProtect and Veeam Backup & Replication",
     "17_backup.png",
     ["🛡️  Rubrik: SLA domains, on-demand snapshots, live mount, restore",
      "🗄️  Cohesity: protection groups, run history, compliance status",
      "💿  Veeam: backup jobs, repository capacity, session logs",
      "📊  Per-platform SLA compliance dashboards",
      "⏱️  Backup job trend charts: duration, size, success rate",
      "🔄  Instant recovery and file-level restore from portal"]),
    ("CMDB — Configuration Management Database",
     "ServiceNow-aligned CI registry with auto-discovery across all platforms",
     "18_cmdb.png",
     ["🗄️  1,800+ CIs auto-discovered from VMware, AWS, OCP, Nutanix",
      "🏷️  Tagging auto-populated from vCenter tags per VM",
      "🏢  Department set to SDx-COE across all CI types",
      "🖥️  ESXi host IP and OS version (ESXi 8.0.x) populated",
      "📊  Columns: Server, IP, Status, Technology, OS, Manager, Owner",
      "🌍  Environment, Timezone, Region, Department, Tagging",
      "📥  Export to CSV and PDF with one click",
      "✏️  Inline editing — click any cell to update CI attributes"]),
    ("Active Directory & DNS Management",
     "Full AD user/group/OU management and DNS zone/record control",
     "19_ad_dns.png",
     ["👤  AD user list: search, enable/disable, unlock",
      "🔑  Reset AD user password from portal",
      "👥  Group management: add/remove members",
      "🏢  OU browser with tree navigation",
      "🌐  DNS zone listing with record count",
      "📝  Add/edit/delete DNS A, CNAME, PTR records"]),
]

# Slide 24 (Users/Audit) is special dual-layout - handled separately

SCREENSHOT_SLIDES_PART2 = [
    ("Insights & Analytics Engine",
     "PostgreSQL 16 daily snapshots — Health Scorecard, Change Detection, Capacity & Cost",
     "22_insights.png",
     ["🩺  Health Scorecard: real-time RAG status across all platforms",
      "🔄  Change Detection: new, removed, modified CIs day over day",
      "📊  Capacity & Cost: per-platform resource utilisation and spend",
      "📋  Executive Summary: auto-generated digest for leadership",
      "📈  Trend arrows and sparklines for key metrics",
      "🗓️  Daily snapshots stored in PostgreSQL 16"]),
    ("Historical Trending",
     "Interactive charts for 7, 14 and 30-day windows across all collected metrics",
     "23_history.png",
     ["📈  VM count, host count, datastore usage trend lines",
      "🖥️  Per-vCenter historical comparison charts",
      "🗓️  Selectable time windows: 7 / 14 / 30 days",
      "📊  Overlay charts: CPU vs Memory utilisation over time",
      "💾  Datastore growth trend for capacity planning"]),
    ("Capacity Forecasting",
     "Linear regression predictions for resource planning across all platforms",
     "24_forecast.png",
     ["🔮  Linear regression forecast for CPU, memory, storage",
      "📅  Projected exhaustion date per resource pool",
      "📊  Confidence bands on forecast charts",
      "🖥️  Per-vCenter and per-cluster forecast breakdown",
      "⚠️  Threshold alerts when forecasted usage > 80%",
      "📈  30, 60 and 90-day forward projections"]),
]


# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def main():
    print("Loading Wipro Template...")
    prs = Presentation(WIPRO_TPL)

    # Remove all existing template slides (130 demo slides)
    print("  Removing template demo slides...")
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId = sldId.get('r:id')
        sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except (KeyError, Exception):
                pass
    print(f"  Cleared. Remaining: {len(prs.slides)} slides")

    # Find layouts
    cover_layout = None
    content_layout = None
    for sl in prs.slide_layouts:
        if sl.name == "Cover_WiproBlue":
            cover_layout = sl
        if sl.name == "1_Text Content":
            content_layout = sl
    assert cover_layout, "Cover_WiproBlue layout not found!"
    assert content_layout, "1_Text Content layout not found!"

    # ─── SLIDE 1: Cover ───
    print("\n[1] Cover slide")
    s = prs.slides.add_slide(cover_layout)
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            ph.text = "LaaS Portal"
            for r in ph.text_frame.paragraphs[0].runs:
                r.font.size = Pt(40); r.font.bold = True
        elif idx == 1:  # Subtitle
            ph.text = "Lab as a Service — Wipro CaaS Platform"
        elif idx == 10:  # Date
            ph.text = "Confidential  |  Q2 2026  |  v6.5"

    # ─── SLIDE 2: Why LaaS ───
    print("[2] Why LaaS slide")
    s = prs.slides.add_slide(content_layout)
    _clear_slide(s)
    _add_header(s, "Why LaaS? — The Challenge",
                "Key pain points that drove the need for a unified lab infrastructure portal")
    
    TW = Emu(3657600); TH = Emu(1920240); TGX = Emu(182880); TGY = Emu(137160)
    SX = Emu(347472); SY = Emu(1188720)
    why_tiles = [
        (RGBColor(0xDC, 0x26, 0x26), "🔴", "Fragmented Tooling",
         "Engineers juggled 8+ separate consoles — vCenter, OCP, Nutanix Prism, AWS, Hyper-V, storage UIs, backup portals and AD tools — with no single pane of glass."),
        (RGBColor(0xEA, 0x58, 0x0C), "🟠", "No Visibility",
         "Management had zero real-time visibility into lab utilisation, capacity trends or cost allocation. Monthly manual spreadsheet collation."),
        (RGBColor(0xCA, 0x8A, 0x04), "🟡", "Slow Provisioning",
         "VM provisioning took 2-5 days through email chains and manual steps. No self-service, no approval workflow, no audit trail."),
        (RGBColor(0x16, 0x6A, 0x34), "🟢", "Siloed Knowledge",
         "Platform expertise locked in individual engineers. New team members took weeks to become productive across all platforms."),
        (RGBColor(0x25, 0x63, 0xEB), "🔵", "Compliance Gaps",
         "No centralised CMDB, no automated change detection, no snapshot compliance tracking. Audit readiness was a recurring scramble."),
        (RGBColor(0x7C, 0x3A, 0xED), "🟣", "Cost Opacity",
         "No chargeback mechanism. Labs consumed resources with no accountability, leading to uncontrolled sprawl and wasted capacity."),
    ]
    for idx, (bg, icon, title, desc) in enumerate(why_tiles):
        r, c = idx // 3, idx % 3
        _add_icon_card(s, SX + c*(TW+TGX), SY + r*(TH+TGY), TW, TH, bg, icon, title, desc, tsz=18)

    # ─── SLIDE 3: Agenda ───
    print("[3] Agenda slide")
    s = prs.slides.add_slide(content_layout)
    _clear_slide(s)
    tx = s.shapes.add_textbox(548640, 274320, 7772400, 457200)
    r = tx.text_frame.paragraphs[0].add_run(); r.text = "Agenda"
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    
    tx2 = s.shapes.add_textbox(548640, 777240, 7772400, 274320)
    r2 = tx2.text_frame.paragraphs[0].add_run()
    r2.text = "Capabilities covered in this presentation — LaaS Portal v6.5"
    r2.font.size = Pt(12); r2.font.color.rgb = DARK_GRAY

    agenda_left = [
        "01  Platform Overview & Architecture",
        "02  Multi-Platform Dashboard — Overview",
        "03  VMware vSphere — VM Management",
        "04  VM Snapshot Management",
        "05  Network Inventory",
        "06  Capacity Planning & Resource Monitoring",
        "07  Project Utilisation Tracking",
        "08  Chargeback & Cost Management",
        "09  VM Request & Approval Workflow",
        "10  IPAM — IP Address Management",
        "11  Asset Inventory Management",
        "12  Red Hat OpenShift (OCP)",
        "13  Nutanix AHV Virtualisation",
        "14  Ansible Automation Platform (AAP)",
    ]
    agenda_right = [
        "15  Amazon Web Services (AWS)",
        "16  Microsoft Hyper-V",
        "17  Enterprise Storage Management",
        "18  Backup & Data Protection",
        "19  CMDB — Configuration Management DB",
        "20  Active Directory & DNS Management",
        "21  User Management & RBAC",
        "22  Insights & Analytics Engine",
        "23  Historical Trending",
        "24  Capacity Forecasting",
        "25  Achievements & Business Outcomes",
    ]
    for col_items, x_pos in [(agenda_left, 548640), (agenda_right, 4754880)]:
        tx = s.shapes.add_textbox(x_pos, 1188720, 3931920, 4754880)
        tf = tx.text_frame; tf.word_wrap = True
        for i, item in enumerate(col_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = item
            r.font.size = Pt(12); r.font.color.rgb = NAVY
            p.space_after = Pt(6)

    # ─── SLIDE 4: Platform Overview (with tiles) ───
    print("[4] Platform Overview slide")
    s = prs.slides.add_slide(content_layout)
    _clear_slide(s)
    _add_header(s, "LaaS Portal — Lab as a Service",
                "A unified self-service portal for multi-platform lab infrastructure")

    # Row 1: VMware, OpenShift, Nutanix
    # Row 2: Ansible, AD&DNS, IPAM&Assets
    # Row 3: AWS, Hyper-V, Storage, Backup (4 columns)
    TILE_W3 = 3566160; GAP3 = Emu(219456)
    cols3 = [347472, 4114800, 7882128]
    HDR_H = 384048; BODY_H = 1100000
    row1_t = 1143000; row2_t = 2700000

    tiles_r1 = [
        (RGBColor(0x00, 0x70, 0xC0), "VMware vSphere",
         ["▪ Full VM lifecycle management", "▪ ESXi host & datastore monitoring", "▪ Cloning, migration, reconfiguration", "▪ Snapshot & network management"]),
        (RGBColor(0x9E, 0x1B, 0x32), "Red Hat OpenShift",
         ["▪ OCP cluster overview & health", "▪ Node & operator visibility", "▪ Container workload monitoring", "▪ Multi-cluster support"]),
        (RGBColor(0x00, 0x70, 0x40), "Nutanix AHV",
         ["▪ Prism Central integration", "▪ AHV VM inventory & metrics", "▪ Cluster resource visibility", "▪ Image management"]),
    ]
    tiles_r2 = [
        (RGBColor(0xFF, 0x6B, 0x00), "Ansible AAP",
         ["▪ Automation job templates", "▪ Playbook execution tracking", "▪ Workflow status visibility", "▪ AAP API integration"]),
        (RGBColor(0x5B, 0x2D, 0x8E), "AD & DNS",
         ["▪ Active Directory user management", "▪ Group & OU management", "▪ DNS zone & record management", "▪ Cache flush & PTR support"]),
        (RGBColor(0x06, 0xB6, 0xD4), "IPAM & Assets",
         ["▪ IP subnet tracking & allocation", "▪ Asset inventory (physical/virtual)", "▪ Ping & power actions on assets", "▪ Export to CSV/Excel"]),
    ]

    for i, (clr, title, bullets) in enumerate(tiles_r1):
        _add_tile(s, cols3[i], row1_t, TILE_W3, HDR_H, BODY_H, clr, title, bullets)
    for i, (clr, title, bullets) in enumerate(tiles_r2):
        _add_tile(s, cols3[i], row2_t, TILE_W3, HDR_H, BODY_H, clr, title, bullets)

    # Row 3: 4 smaller tiles
    row3_t = 4250000
    TILE_W4 = 2697480; GAP4 = 137160
    cols4 = [347472, 347472 + TILE_W4 + GAP4, 347472 + 2*(TILE_W4 + GAP4), 347472 + 3*(TILE_W4 + GAP4)]
    BODY_H4 = 900000

    tiles_r3 = [
        (RGBColor(0xFF, 0x99, 0x00), "Amazon Web Services",
         ["▪ EC2, S3, RDS, VPC management", "▪ Cost Explorer integration", "▪ IAM & SSO support"]),
        (RGBColor(0x00, 0x78, 0xD4), "Microsoft Hyper-V",
         ["▪ WinRM PowerShell remoting", "▪ VM power & checkpoints", "▪ Per-host CPU/RAM bars"]),
        (RGBColor(0x2D, 0x6B, 0x9E), "Enterprise Storage",
         ["▪ Pure, Dell, HPE, NetApp", "▪ Volumes, LUNs, capacity", "▪ Native REST APIs"]),
        (RGBColor(0x16, 0x7F, 0x6E), "Backup & DR",
         ["▪ Rubrik, Cohesity, Veeam", "▪ SLA & job management", "▪ Instant recovery & restore"]),
    ]
    for i, (clr, title, bullets) in enumerate(tiles_r3):
        _add_tile(s, cols4[i], row3_t, TILE_W4, HDR_H, BODY_H4, clr, title, bullets)

    # ─── SLIDE 5: Technology Stack ───
    print("[5] Technology Stack slide")
    s = prs.slides.add_slide(content_layout)
    _clear_slide(s)
    _add_header(s, "Technology Stack",
                "The components that power the LaaS Portal")

    tech_r1 = [
        (RGBColor(0x00, 0x70, 0xC0), "Frontend",
         ["• React 18 + Vite", "• TailwindCSS", "• Chart.js / Recharts", "• Playwright (testing)"]),
        (RGBColor(0xFF, 0x6B, 0x00), "Backend",
         ["• Python FastAPI", "• PostgreSQL 16", "• Pydantic validation", "• JWT Auth (LDAP/Local)"]),
        (RGBColor(0x00, 0x70, 0xC0), "VMware Layer",
         ["• pyVmomi SDK", "• vSphere 7/8 REST API", "• vCenter multi-VC", "• ESXi 7/8 hosts"]),
    ]
    tech_r2 = [
        (RGBColor(0x9E, 0x1B, 0x32), "Platforms",
         ["• OCP via REST API", "• Nutanix Prism Central", "• Ansible AAP API", "• ldap3 (AD)"]),
        (RGBColor(0x00, 0x70, 0x40), "Networking / IPAM",
         ["• IPAM REST integration", "• DNS (dnspython)", "• AD/LDAP (ldap3)", "• SMTP notifications"]),
        (RGBColor(0x5B, 0x2D, 0x8E), "Auth & Security",
         ["• AD LDAP bind auth", "• Local admin fallback", "• 4-tier RBAC model", "• Session token (sessionStorage)"]),
    ]

    for i, (clr, title, bullets) in enumerate(tech_r1):
        _add_tile(s, cols3[i], row1_t, TILE_W3, 347472, BODY_H, clr, title, bullets)
    for i, (clr, title, bullets) in enumerate(tech_r2):
        _add_tile(s, cols3[i], row2_t, TILE_W3, 347472, BODY_H, clr, title, bullets)

    # Row 3: New in v6.5 tiles
    new_tiles = [
        (RGBColor(0x8B, 0x5C, 0xF6), "New Platforms (v6.5)",
         ["▪ AWS SDK (boto3) — EC2, S3, RDS, VPC", "▪ WinRM — Hyper-V remoting", "▪ Storage REST — Pure, Dell, HPE, NetApp"]),
        (RGBColor(0x06, 0x96, 0x68), "New Data Layer (v6.5)",
         ["▪ PostgreSQL 16 — daily snapshot DB", "▪ Insights, CMDB, IPAM persistence", "▪ ServiceNow Table API — CMDB push"]),
        (RGBColor(0xD9, 0x4F, 0x00), "New Backup & DR (v6.5)",
         ["▪ Rubrik Security Cloud API", "▪ Cohesity DataProtect API", "▪ Veeam Backup & Replication API"]),
    ]
    COL3_W_NEW = (11494008 - 2 * 182880) // 3
    cols3_new = [347472, 347472 + COL3_W_NEW + 182880, 347472 + 2*(COL3_W_NEW + 182880)]
    for i, (clr, title, bullets) in enumerate(new_tiles):
        _add_tile(s, cols3_new[i], 4250000, COL3_W_NEW, 347472, BODY_H4, clr, title, bullets)

    # ─── SLIDES 6-24: Screenshot slides ───
    print("[6-24] Screenshot slides...")
    for i, (title, subtitle, img, bullets) in enumerate(SCREENSHOT_SLIDES):
        s = prs.slides.add_slide(content_layout)
        _clear_slide(s)
        _screenshot_slide(s, title, subtitle, img, bullets)
        print(f"  [{i+6}] {title}")

    # ─── SLIDE 25: Users & Audit (dual-image) ───
    print("[25] Users & Audit slide")
    s = prs.slides.add_slide(content_layout)
    _clear_slide(s)
    _add_header(s, "User Management & Role-Based Access Control",
                "Portal user management and RBAC model")
    
    # Two screenshots side by side
    audit_img = os.path.join(SDIR, "21_audit.png")
    users_img = os.path.join(SDIR, "20_users.png")
    if os.path.exists(audit_img):
        rect1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 329184, 1152144, 5614416, 3174111)
        rect1.fill.background(); rect1.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); rect1.line.width = Pt(0.75)
        s.shapes.add_picture(audit_img, 347472, 1170432, 5577840, 3137535)
    if os.path.exists(users_img):
        rect2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 6071616, 1152144, 5614416, 3174111)
        rect2.fill.background(); rect2.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); rect2.line.width = Pt(0.75)
        s.shapes.add_picture(users_img, 6089904, 1170432, 5577840, 3137535)
    
    # Labels
    for text, x in [("Audit Log", 347472), ("User Management", 6089904)]:
        tx = s.shapes.add_textbox(x, 4381119, 5577840, 256032)
        r = tx.text_frame.paragraphs[0].add_run(); r.text = text
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY

    # RBAC role cards
    roles = [
        (347472, "Super Admin", "Full access: all platforms, users, pricing, delete", RGBColor(0xDC,0x26,0x26)),
        (2633472, "Admin", "Manage VMs, approve requests, configure settings", RGBColor(0x25,0x63,0xEB)),
        (4919472, "Operator", "VM lifecycle ops, snapshot, clone, migrate", RGBColor(0x16,0x6A,0x34)),
        (7205472, "Requester", "Submit VM requests; view own resources only", RGBColor(0xCA,0x8A,0x04)),
        (9491472, "Viewer", "Read-only access to all dashboards and reports", RGBColor(0x64,0x74,0x8B)),
    ]
    for x, title, desc, clr in roles:
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, 4692015, 2240280, 310896)
        rect.fill.solid(); rect.fill.fore_color.rgb = clr; rect.line.fill.background()
        tx = s.shapes.add_textbox(x + 45720, 4728591, 2194560, 256032)
        r = tx.text_frame.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
        tx2 = s.shapes.add_textbox(x + 45720, 5021199, 2194560, 365760)
        r2 = tx2.text_frame.paragraphs[0].add_run(); r2.text = desc
        r2.font.size = Pt(9); r2.font.color.rgb = BODY_TEXT

    # ─── SLIDES 26-28: More screenshot slides ───
    print("[26-28] Insights, History, Forecast slides...")
    for i, (title, subtitle, img, bullets) in enumerate(SCREENSHOT_SLIDES_PART2):
        s = prs.slides.add_slide(content_layout)
        _clear_slide(s)
        _screenshot_slide(s, title, subtitle, img, bullets)
        print(f"  [{i+26}] {title}")

    # ─── SLIDE 29: Achievements & Business Outcomes ───
    print("[29] Achievements slide")
    s = prs.slides.add_slide(content_layout)
    _clear_slide(s)
    _add_header(s, "Achievements & Business Outcomes",
                "Measurable impact delivered by the LaaS Portal across CISS-COE infrastructure operations")

    AW = Emu(2743200); AH = Emu(1920240); AGX = Emu(137160); AGY = Emu(137160)
    ASX = Emu(347472); ASY = Emu(1188720)
    ach = [
        (RGBColor(0x05,0x96,0x69), "🖥️", "Single Pane of Glass",
         "8 platforms unified — VMware, OCP, Nutanix, AWS, Hyper-V, Storage, Backup, AD/DNS — all in one portal."),
        (RGBColor(0x25,0x63,0xEB), "⚡", "90% Faster Provisioning",
         "VM provisioning reduced from 2-5 days to under 30 minutes with self-service and automated approvals."),
        (RGBColor(0xDC,0x26,0x26), "📊", "1,800+ CIs Tracked",
         "Auto-discovered CMDB with tagging, department, environment, timezone, region — export to CSV/PDF."),
        (RGBColor(0x7C,0x3A,0xED), "💰", "Cost Transparency",
         "Per-project chargeback across VMware, OCP, Nutanix. AWS Cost Explorer integrated. Full accountability."),
        (RGBColor(0xEA,0x58,0x0C), "🩺", "Proactive Health",
         "Real-time RAG scorecard across all platforms. Daily snapshots, change detection and capacity forecasting."),
        (RGBColor(0x0E,0x7C,0x90), "🔒", "4-Tier RBAC",
         "Super Admin → Admin → Operator → Requester/Viewer. AD LDAP + local fallback. Full audit trail."),
        (RGBColor(0xB4,0x5D,0x09), "📈", "Data-Driven Insights",
         "PostgreSQL 16 daily snapshots power trend analysis, linear regression forecasting and executive summaries."),
        (RGBColor(0x16,0x6A,0x34), "🤝", "Team Enablement",
         "New engineers productive in hours, not weeks. Self-service reduces dependency on specialists by 70%."),
    ]
    for idx, (bg, icon, title, desc) in enumerate(ach):
        r, c = idx // 4, idx % 4
        _add_icon_card(s, ASX + c*(AW+AGX), ASY + r*(AH+AGY), AW, AH, bg, icon, title, desc, tsz=16)

    # ─── SLIDE 30: Thank You (Cover) ───
    print("[30] Thank You slide")
    s = prs.slides.add_slide(cover_layout)
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "Thank You"
            for r in ph.text_frame.paragraphs[0].runs:
                r.font.size = Pt(40); r.font.bold = True
        elif idx == 1:
            ph.text = "LaaS Portal — Lab as a Service  |  Wipro CaaS Platform"
        elif idx == 10:
            ph.text = "Confidential  |  Wipro Limited  |  2026"

    # ─── SAVE ───
    print(f"\nSaving to: {OUT}")
    prs.save(OUT)
    size_mb = os.path.getsize(OUT) / (1024 * 1024)
    print(f"  File size: {size_mb:.1f} MB")
    print(f"  Total slides: {len(prs.slides)}")
    print("Done!")


if __name__ == "__main__":
    main()

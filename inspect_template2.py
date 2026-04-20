"""Inspect Wipro template - logo positions, slide master shapes, and layout details."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

prs = Presentation(r'C:\Users\Administrator\Desktop\Wipro Template.pptx')

print("=== SLIDE MASTERS ===")
for mi, master in enumerate(prs.slide_masters):
    print(f"\nMaster {mi}:")
    for shape in master.shapes:
        print(f"  Shape: name='{shape.name}' type={shape.shape_type} "
              f"left={round(shape.left/914400,2)}in top={round(shape.top/914400,2)}in "
              f"w={round(shape.width/914400,2)}in h={round(shape.height/914400,2)}in")
        if hasattr(shape, 'text') and shape.text:
            print(f"    text='{shape.text[:60]}'")

print("\n=== SLIDE LAYOUTS (shapes) ===")
for li, layout in enumerate(prs.slide_layouts):
    print(f"\nLayout {li} [{layout.name}]:")
    for shape in layout.shapes:
        print(f"  Shape: name='{shape.name}' type={shape.shape_type} "
              f"left={round(shape.left/914400,2) if shape.left else 'N/A'}in "
              f"top={round(shape.top/914400,2) if shape.top else 'N/A'}in "
              f"w={round(shape.width/914400,2) if shape.width else 'N/A'}in "
              f"h={round(shape.height/914400,2) if shape.height else 'N/A'}in")
        if hasattr(shape, 'text') and shape.text:
            print(f"    text='{shape.text[:80]}'")

"""
Magic Migrate PPT Builder — VMware to Hyper-V Focus
10 slides using Wipro brand colours
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
import copy, os

# ── Brand Palette ─────────────────────────────────────────────────────────────
WIPRO_PURPLE   = RGBColor(0x5B, 0x21, 0xB6)   # deep purple
WIPRO_VIOLET   = RGBColor(0x7C, 0x3A, 0xED)   # mid purple
WIPRO_INDIGO   = RGBColor(0x43, 0x38, 0xCA)   # indigo accent
BG_DARK        = RGBColor(0x0D, 0x10, 0x17)   # near-black bg
CARD_BG        = RGBColor(0x11, 0x18, 0x27)   # card dark
CARD_BORDER    = RGBColor(0x1F, 0x29, 0x37)   # border
TEXT_WHITE     = RGBColor(0xF9, 0xFA, 0xFB)   # primary text
TEXT_MUTE      = RGBColor(0x6B, 0x72, 0x80)   # muted
ACCENT_BLUE    = RGBColor(0x3B, 0x82, 0xF6)   # Hyper-V blue
ACCENT_GREEN   = RGBColor(0x10, 0xB9, 0x81)   # pass/success
ACCENT_ORANGE  = RGBColor(0xF9, 0x73, 0x16)   # warning
ACCENT_RED     = RGBColor(0xEF, 0x44, 0x44)   # fail/risk

SCREENSHOT_DIR = r"C:\caas-dashboard\screenshots"
WIPRO_LOGO     = r"C:\caas-dashboard\wipro_logo.png"

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    blank = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank)

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=14, bold=False,
                color=TEXT_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def add_logo(slide):
    if os.path.isfile(WIPRO_LOGO):
        slide.shapes.add_picture(WIPRO_LOGO, Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.45))

def slide_bg(slide):
    add_rect(slide, 0, 0, 13.333, 7.5, BG_DARK)

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.333, 1.05, RGBColor(0x09, 0x0D, 0x14))
    add_rect(slide, 0, 1.05, 13.333, 0.04, WIPRO_VIOLET)
    add_textbox(slide, 0.35, 0.12, 10.5, 0.55, title,
                font_size=22, bold=True, color=TEXT_WHITE)
    if subtitle:
        add_textbox(slide, 0.35, 0.62, 10.5, 0.35, subtitle,
                    font_size=11.5, color=TEXT_MUTE)
    add_logo(slide)

def footer(slide, page_num):
    add_rect(slide, 0, 7.15, 13.333, 0.35, RGBColor(0x09, 0x0D, 0x14))
    add_textbox(slide, 0.3, 7.17, 8, 0.28,
                "Magic Migrate  |  VMware → Microsoft Hyper-V  |  Wipro CaaS Platform  |  Confidential",
                font_size=8, color=TEXT_MUTE)
    add_textbox(slide, 12.5, 7.17, 0.7, 0.28, str(page_num),
                font_size=8, color=TEXT_MUTE, align=PP_ALIGN.RIGHT)

def add_pill(slide, x, y, w, h, text, bg, fg=TEXT_WHITE, font_size=10):
    r = add_rect(slide, x, y, w, h, bg)
    r.line.fill.background()
    add_textbox(slide, x + 0.05, y + 0.02, w - 0.1, h - 0.04,
                text, font_size=font_size, color=fg, align=PP_ALIGN.CENTER, bold=True)

def add_screenshot(slide, path, x, y, w, h):
    if os.path.isfile(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        # grey placeholder
        r = add_rect(slide, x, y, w, h, CARD_BORDER)
        add_textbox(slide, x + 0.1, y + h/2 - 0.2, w - 0.2, 0.4,
                    f"[Screenshot: {os.path.basename(path)}]",
                    font_size=9, color=TEXT_MUTE, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, accent=WIPRO_VIOLET):
    r = add_rect(slide, x, y, w, h, CARD_BG)
    # left accent bar
    add_rect(slide, x, y, 0.04, h, accent)
    return r

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide1_cover(prs):
    """Cover — Magic Migrate branding"""
    sl = blank_slide(prs)
    slide_bg(sl)

    # Full-width purple gradient band
    add_rect(sl, 0, 0, 13.333, 7.5, RGBColor(0x0D, 0x10, 0x17))
    # left vivid stripe
    add_rect(sl, 0, 0, 0.55, 7.5, WIPRO_VIOLET)
    add_rect(sl, 0.55, 0, 0.08, 7.5, WIPRO_PURPLE)

    # big title
    add_textbox(sl, 1.1, 1.4, 9, 1.1, "✨ Magic Migrate",
                font_size=48, bold=True, color=TEXT_WHITE)
    add_textbox(sl, 1.1, 2.55, 9, 0.55,
                "Cross-Hypervisor VM Migration  |  VMware → Microsoft Hyper-V",
                font_size=20, bold=False, color=ACCENT_BLUE)
    add_textbox(sl, 1.1, 3.2, 9, 0.45,
                "Seamless, zero-downtime VM migrations with automated pre-flight, resource mapping & post-migration validation",
                font_size=13, color=TEXT_MUTE)

    # 4 target pills
    targets = [
        ("OpenShift", ACCENT_RED),
        ("Nutanix AHV", ACCENT_GREEN),
        ("Hyper-V", ACCENT_BLUE),
        ("HPE VME", RGBColor(0x00, 0x96, 0xD6)),
    ]
    for i, (lbl, col) in enumerate(targets):
        add_pill(sl, 1.1 + i * 2.05, 3.85, 1.85, 0.42, lbl, col, font_size=11)

    # bottom meta
    add_textbox(sl, 1.1, 6.5, 8, 0.35,
                "Wipro CaaS Platform  |  LaaS Portal v6.5  |  May 2026  |  Confidential",
                font_size=10, color=TEXT_MUTE)

    # wipro logo top right
    if os.path.isfile(WIPRO_LOGO):
        sl.shapes.add_picture(WIPRO_LOGO, Inches(11.5), Inches(0.25), Inches(1.6), Inches(0.55))


def slide2_problem(prs):
    """Why Migrate? — The VMware Pain"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Why Migrate? — The VMware Challenge",
               "Business drivers pushing organisations from VMware to Hyper-V")
    footer(sl, 2)

    pain_points = [
        ("💰", "Licensing Shock",
         "VMware by Broadcom doubled+ per-core pricing post-acquisition. "
         "Enterprise labs face 300-500% cost increases on renewal."),
        ("🔒", "Vendor Lock-in",
         "Proprietary vSphere APIs, VMFS datastores, and vCenter dependency "
         "create deep lock-in. No exit path without re-platforming."),
        ("📊", "Hyper-V Value",
         "Windows Server 2022 Datacenter includes unlimited Hyper-V VMs. "
         "Built-in WinRM, PowerShell, and Azure integration at no extra cost."),
        ("⚡", "Performance Parity",
         "Modern Hyper-V matches ESXi on NUMA, SR-IOV, NVMe, live migration "
         "(Live Migration) and clustering (Failover Cluster)."),
        ("🔄", "Azure Synergy",
         "Hyper-V shares Azure Stack HCI architecture. VMs migrate cleanly "
         "to Azure — same VHDX format, same networking model."),
        ("🛡️", "Security Posture",
         "Shielded VMs, TPM 2.0, Host Guardian Service, and BitLocker "
         "give Hyper-V enterprise-grade security without add-on cost."),
    ]

    col_w, col_h, gap = 3.9, 1.65, 0.18
    for i, (icon, title, body) in enumerate(pain_points):
        col = i % 3
        row = i // 3
        x = 0.3 + col * (col_w + gap)
        y = 1.3 + row * (col_h + gap)
        card(sl, x, y, col_w, col_h,
             accent=[ACCENT_RED, ACCENT_ORANGE, ACCENT_BLUE,
                     ACCENT_GREEN, WIPRO_VIOLET, ACCENT_BLUE][i])
        add_textbox(sl, x + 0.18, y + 0.1, col_w - 0.25, 0.35,
                    f"{icon}  {title}", font_size=12, bold=True, color=TEXT_WHITE)
        add_textbox(sl, x + 0.18, y + 0.45, col_w - 0.3, col_h - 0.55,
                    body, font_size=10, color=TEXT_MUTE)


def slide3_overview(prs):
    """Magic Migrate — Feature Overview"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Magic Migrate — Feature Overview",
               "5-step wizard with automated pre-flight, mapping, execution and validation")
    footer(sl, 3)

    steps = [
        ("1", "Source VMs", "Select vCenter & pick VMs to migrate\nFilter by cluster, host, datastore, OS", WIPRO_VIOLET),
        ("2", "Target Platform", "Choose Hyper-V host as destination\nMap to WS 2019/2022 cluster", ACCENT_BLUE),
        ("3", "Pre-flight", "Automated 12-point compatibility check\nCPU, RAM, disk, network, OS, tools", ACCENT_GREEN),
        ("4", "Resource Mapping", "vSwitch → vSwitch  |  VMFS → CSV\nCPU/RAM pinning on Hyper-V host", ACCENT_ORANGE),
        ("5", "Review & Execute", "Confirm plan, track live events\nPost-migration Ansible playbooks", ACCENT_RED),
    ]

    box_w = 2.3
    for i, (num, title, body, col) in enumerate(steps):
        x = 0.3 + i * (box_w + 0.12)
        add_rect(sl, x, 1.3, box_w, 3.5, CARD_BG)
        add_rect(sl, x, 1.3, box_w, 0.05, col)
        # number circle
        add_rect(sl, x + box_w/2 - 0.28, 1.45, 0.56, 0.56, col)
        add_textbox(sl, x + box_w/2 - 0.25, 1.47, 0.5, 0.52,
                    num, font_size=18, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, x + 0.1, 2.15, box_w - 0.2, 0.4,
                    title, font_size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, x + 0.1, 2.6, box_w - 0.2, 1.9,
                    body, font_size=10, color=TEXT_MUTE, align=PP_ALIGN.CENTER)
        # arrow (except last)
        if i < 4:
            add_textbox(sl, x + box_w + 0.0, 2.75, 0.15, 0.35,
                        "→", font_size=14, color=WIPRO_VIOLET, align=PP_ALIGN.CENTER)

    # bottom feature pills
    pills = ["Move Groups", "Batch Migration", "Live Event Log", "Post-Migration Tasks",
             "Preflight Report", "Plan Management", "VHDX Conversion", "AAP Integration"]
    for i, p in enumerate(pills):
        add_pill(sl, 0.3 + i * 1.58, 5.05, 1.45, 0.38, p,
                 RGBColor(0x1F, 0x29, 0x37), TEXT_MUTE, font_size=9)

    add_textbox(sl, 0.3, 5.6, 12.5, 0.3,
                "📋  Plan-based workflow: Create Plan → Add Move Groups → Run Pre-flight → Execute → Validate",
                font_size=11, color=ACCENT_GREEN, bold=True)


def slide4_wizard_screenshot(prs):
    """Migration Wizard — Live Screenshot"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Migration Wizard — Step-by-Step UI",
               "Guided 5-step wizard: Source selection → Target → Pre-flight → Mapping → Execute")
    footer(sl, 4)

    # Large screenshot (use vmware or overview screenshot as representative)
    img = os.path.join(SCREENSHOT_DIR, "02_vmware_vms.png")
    add_screenshot(sl, img, 0.3, 1.2, 7.7, 5.0)

    # Right side callouts
    callouts = [
        (WIPRO_VIOLET,  "Step 1 — Source",
         "Select vCenter instance\nFilter & multi-select VMs\nShows CPU/RAM/disk live"),
        (ACCENT_BLUE,   "Step 2 — Target: Hyper-V",
         "Pick Hyper-V host cluster\nSelect destination CSV\nMap vSwitch adapters"),
        (ACCENT_GREEN,  "Step 3 — Pre-flight",
         "12 automated checks\nOS compatibility, tools\nNetwork & storage"),
        (ACCENT_ORANGE, "Steps 4-5 — Execute",
         "Named migration plan\nMove group batching\nLive event streaming"),
    ]
    for i, (col, title, body) in enumerate(callouts):
        y = 1.2 + i * 1.28
        card(sl, 8.2, y, 4.85, 1.18, accent=col)
        add_textbox(sl, 8.38, y + 0.08, 4.55, 0.32,
                    title, font_size=11, bold=True, color=TEXT_WHITE)
        add_textbox(sl, 8.38, y + 0.4, 4.55, 0.72,
                    body, font_size=9.5, color=TEXT_MUTE)


def slide5_preflight(prs):
    """Pre-flight Checks — Automated Compatibility"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Automated Pre-flight Compatibility Checks",
               "12-point assessment before any migration executes — blocks unsafe migrations automatically")
    footer(sl, 5)

    checks = [
        ("✅", "OS Compatibility",     "Windows Server / RHEL on Hyper-V guest",  ACCENT_GREEN),
        ("✅", "CPU Architecture",      "x86-64, no ESXi-specific CPU features",   ACCENT_GREEN),
        ("✅", "RAM Headroom",           "Destination host has sufficient free RAM", ACCENT_GREEN),
        ("✅", "Disk Format",            "VMDK → VHDX convertible (no RDM)",        ACCENT_GREEN),
        ("✅", "Network Adapter",        "VMXNET3 mapped to Hyper-V synthetic NIC", ACCENT_GREEN),
        ("✅", "VMware Tools",           "Tools removed / replaced by Hyper-V IC",  ACCENT_GREEN),
        ("⚠️", "VM Snapshots",           "Snapshots must be consolidated first",    ACCENT_ORANGE),
        ("⚠️", "Shared VMDK",           "Multi-writer disks need manual review",   ACCENT_ORANGE),
        ("⚠️", "USB / PCI Passthrough", "Pass-through devices need re-config",     ACCENT_ORANGE),
        ("✅", "Storage Cluster",        "CSV path available on Hyper-V host",      ACCENT_GREEN),
        ("✅", "Live Migration Port",    "WinRM 5985/5986 reachable",               ACCENT_GREEN),
        ("✅", "Datastore Free Space",   ">120% of VM disk size available",         ACCENT_GREEN),
    ]

    col_w = 4.0
    for i, (icon, title, detail, col) in enumerate(checks):
        r = i % 3
        c = i // 3
        x = 0.3 + c * (col_w + 0.22)
        y = 1.3 + r * 0.9
        add_rect(sl, x, y, col_w, 0.78, CARD_BG)
        add_rect(sl, x, y, 0.04, 0.78, col)
        add_textbox(sl, x + 0.15, y + 0.05, 0.4, 0.3,
                    icon, font_size=13, color=col)
        add_textbox(sl, x + 0.55, y + 0.05, col_w - 0.65, 0.32,
                    title, font_size=11, bold=True, color=TEXT_WHITE)
        add_textbox(sl, x + 0.55, y + 0.38, col_w - 0.65, 0.35,
                    detail, font_size=9, color=TEXT_MUTE)

    # legend
    add_pill(sl, 0.3,  6.98, 1.2, 0.3, "✅ Pass", ACCENT_GREEN, font_size=9)
    add_pill(sl, 1.6,  6.98, 1.5, 0.3, "⚠️ Warning", ACCENT_ORANGE, font_size=9)
    add_textbox(sl, 3.3, 7.0, 8, 0.28,
                "Warnings do not block migration — they generate a preflight report for engineer review",
                font_size=9, color=TEXT_MUTE)


def slide6_resource_mapping(prs):
    """Resource Mapping — VMware to Hyper-V"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Resource Mapping — VMware → Hyper-V",
               "Automatic and manual mapping of vSphere resources to Hyper-V equivalents")
    footer(sl, 6)

    mappings = [
        ("VMware vSphere",     "Hyper-V Equivalent",        "Notes"),
        ("ESXi Host",          "Hyper-V Host (WS 2022)",    "WinRM PowerShell remoting"),
        ("vCenter (VC)",       "System Center VMM / SCVMM", "Or direct Hyper-V host mgmt"),
        ("Datastore (VMFS)",   "CSV (Cluster Shared Vol.)", "SMB 3.0 / iSCSI / FC"),
        ("VMDK Disk",          "VHDX Disk",                 "VHDX = same perf, dynamic"),
        ("vSwitch / dvSwitch", "Hyper-V Virtual Switch",    "External / Internal / Private"),
        ("VMXNET3 NIC",        "Hyper-V Synthetic NIC",     "No driver install needed"),
        ("VM Snapshot",        "Hyper-V Checkpoint",        "Production or Standard"),
        ("VMware Tools",       "Hyper-V Integration Svcs",  "Auto-installed via ISO"),
        ("vMotion (live)",     "Hyper-V Live Migration",    "Zero-downtime over cluster"),
        ("Storage vMotion",    "Storage Live Migration",    "Disk moved while VM runs"),
        ("VM Hardware v19+",   "VM Gen 2 (UEFI)",           "Secure Boot supported"),
    ]

    col_widths = [3.5, 3.5, 5.3]
    headers = mappings[0]
    rows = mappings[1:]

    # Header row
    x_positions = [0.3, 3.9, 7.5]
    for c, (hdr, w, x) in enumerate(zip(headers, col_widths, x_positions)):
        add_rect(sl, x, 1.25, w, 0.38, WIPRO_PURPLE)
        add_textbox(sl, x + 0.1, 1.28, w - 0.15, 0.32,
                    hdr, font_size=10.5, bold=True, color=TEXT_WHITE)

    row_colors = [CARD_BG, RGBColor(0x0F, 0x17, 0x24)]
    for ri, (vm_res, hv_res, notes) in enumerate(rows):
        y = 1.65 + ri * 0.47
        bg = row_colors[ri % 2]
        add_rect(sl, 0.3, y, 12.8, 0.44, bg)
        # VMware col — orange
        add_textbox(sl, 0.4, y + 0.07, 3.35, 0.3,
                    vm_res, font_size=10, color=ACCENT_ORANGE, bold=True)
        # Arrow
        add_textbox(sl, 3.65, y + 0.07, 0.3, 0.3,
                    "→", font_size=11, color=WIPRO_VIOLET, bold=True)
        # Hyper-V col — blue
        add_textbox(sl, 3.95, y + 0.07, 3.45, 0.3,
                    hv_res, font_size=10, color=ACCENT_BLUE, bold=True)
        # Notes col
        add_textbox(sl, 7.55, y + 0.07, 5.2, 0.3,
                    notes, font_size=9.5, color=TEXT_MUTE)


def slide7_hyperv_screenshot(prs):
    """Hyper-V Dashboard Screenshot"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Microsoft Hyper-V — Live Dashboard",
               "Real-time host monitoring, VM inventory & management via WinRM PowerShell remoting")
    footer(sl, 7)

    img = os.path.join(SCREENSHOT_DIR, "15_hyperv.png")
    add_screenshot(sl, img, 0.3, 1.2, 8.0, 5.1)

    # right callouts
    features = [
        (ACCENT_BLUE,   "Host Dashboard",
         "Per-host CPU / RAM bars\nVM count, uptime, status\nAuto-refresh via WinRM"),
        (ACCENT_GREEN,  "VM Management",
         "Power On / Off / Restart\nCheckpoint create & restore\nConnect to console"),
        (ACCENT_ORANGE, "Storage (CSV)",
         "Cluster Shared Volumes\nFree space monitoring\nVHDX path visibility"),
        (WIPRO_VIOLET,  "Migration Target",
         "Selected as Hyper-V target\nin Magic Migrate Step 2\nAuto-discovered hosts"),
    ]
    for i, (col, title, body) in enumerate(features):
        y = 1.2 + i * 1.28
        card(sl, 8.5, y, 4.6, 1.18, accent=col)
        add_textbox(sl, 8.68, y + 0.08, 4.3, 0.32,
                    title, font_size=11, bold=True, color=TEXT_WHITE)
        add_textbox(sl, 8.68, y + 0.42, 4.3, 0.7,
                    body, font_size=9.5, color=TEXT_MUTE)


def slide8_move_groups(prs):
    """Move Groups & Batch Migration"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Move Groups — Batch Migration & Post-Tasks",
               "Group VMs by app tier, migrate in waves, then trigger Ansible playbooks automatically")
    footer(sl, 8)

    # Left — move group concept
    add_rect(sl, 0.3, 1.25, 5.8, 5.45, CARD_BG)
    add_rect(sl, 0.3, 1.25, 5.8, 0.04, WIPRO_VIOLET)
    add_textbox(sl, 0.5, 1.35, 5.4, 0.38,
                "Move Group Architecture", font_size=13, bold=True, color=TEXT_WHITE)

    waves = [
        ("Wave 1 — Infrastructure",  ["DNS Server", "AD Controller", "NTP Server"],      ACCENT_RED),
        ("Wave 2 — Middleware",       ["Web Server", "App Server", "Cache Server"],       ACCENT_ORANGE),
        ("Wave 3 — Application",      ["CRM App", "ERP App", "Reporting"],                ACCENT_BLUE),
        ("Wave 4 — Data",             ["DB Primary", "DB Replica", "Archive NFS"],        ACCENT_GREEN),
    ]
    for wi, (wave_label, vms, col) in enumerate(waves):
        wy = 1.8 + wi * 1.18
        add_rect(sl, 0.45, wy, 5.5, 1.05, RGBColor(0x0F, 0x17, 0x24))
        add_rect(sl, 0.45, wy, 0.04, 1.05, col)
        add_textbox(sl, 0.6, wy + 0.05, 5.0, 0.28,
                    wave_label, font_size=10, bold=True, color=col)
        for vi, vm in enumerate(vms):
            add_textbox(sl, 0.65 + vi * 1.78, wy + 0.38, 1.65, 0.28,
                        f"🖥 {vm}", font_size=8.5, color=TEXT_MUTE)

    # Right — features
    add_rect(sl, 6.4, 1.25, 6.65, 5.45, CARD_BG)
    add_rect(sl, 6.4, 1.25, 6.65, 0.04, ACCENT_BLUE)
    add_textbox(sl, 6.6, 1.35, 6.2, 0.38,
                "Key Capabilities", font_size=13, bold=True, color=TEXT_WHITE)

    feats = [
        ("📦", "Named Move Groups",
         "Create named batches — e.g. 'Web Tier', 'DB Tier'.\nExecute independently or sequentially."),
        ("🔀", "Per-Group Target Mapping",
         "Each group can target a different Hyper-V host\nor CSV volume for workload isolation."),
        ("🤖", "Post-Migration Playbooks",
         "Trigger AAP / Ansible playbooks after each group\ncompletes — configure, validate, notify."),
        ("📊", "Live Event Log",
         "Real-time streaming of migration events:\nconvert, copy, register, power-on status."),
        ("🔁", "Retry & Resume",
         "Failed VMs retried automatically.\nPartial group completion tracked per VM."),
    ]
    for fi, (icon, title, body) in enumerate(feats):
        fy = 1.8 + fi * 1.0
        add_textbox(sl, 6.55, fy, 0.4, 0.38, icon, font_size=14, color=WIPRO_VIOLET)
        add_textbox(sl, 7.0, fy, 5.8, 0.3,
                    title, font_size=11, bold=True, color=TEXT_WHITE)
        add_textbox(sl, 7.0, fy + 0.32, 5.8, 0.55,
                    body, font_size=9.5, color=TEXT_MUTE)


def slide9_plan_management(prs):
    """Migration Plan Management & Execution"""
    sl = blank_slide(prs)
    slide_bg(sl)
    header_bar(sl, "Migration Plan Management & Execution Tracking",
               "Named plans, status lifecycle, live events and historical reporting")
    footer(sl, 9)

    # Plan lifecycle
    add_textbox(sl, 0.3, 1.22, 8, 0.32,
                "Plan Lifecycle", font_size=12, bold=True, color=TEXT_MUTE)

    statuses = [
        ("Draft",       CARD_BORDER,    "Plan created,\nnot yet run"),
        ("Pre-flight",  ACCENT_ORANGE,  "Running 12-point\nchecks now"),
        ("Ready",       ACCENT_BLUE,    "All checks\npassed"),
        ("Executing",   WIPRO_VIOLET,   "Migration\nin progress"),
        ("Completed",   ACCENT_GREEN,   "All VMs live\non Hyper-V"),
        ("Failed",      ACCENT_RED,     "Error — see\nevent log"),
    ]
    for i, (label, col, desc) in enumerate(statuses):
        x = 0.3 + i * 2.12
        add_rect(sl, x, 1.6, 1.88, 1.1, CARD_BG)
        add_rect(sl, x, 1.6, 1.88, 0.05, col)
        add_pill(sl, x + 0.12, 1.75, 1.64, 0.34, label, col, font_size=10)
        add_textbox(sl, x + 0.1, 2.17, 1.7, 0.48,
                    desc, font_size=9, color=TEXT_MUTE, align=PP_ALIGN.CENTER)

    # Reports section
    add_textbox(sl, 0.3, 2.88, 8, 0.32,
                "Reports & Audit", font_size=12, bold=True, color=TEXT_MUTE)

    report_cards = [
        (ACCENT_GREEN,  "Migration Summary",
         "Per-plan: VMs migrated, success rate,\ntotal time, data transferred, errors"),
        (ACCENT_BLUE,   "Pre-flight Report",
         "Detailed per-VM check results,\nwarnings with remediation steps"),
        (ACCENT_ORANGE, "Event Timeline",
         "Chronological event log per plan:\nconvert → copy → register → boot"),
        (WIPRO_VIOLET,  "Post-Task Audit",
         "Ansible job results per move group:\nplaybook, status, output snippet"),
    ]
    for i, (col, title, body) in enumerate(report_cards):
        x = 0.3 + i * 3.25
        card(sl, x, 3.25, 3.05, 1.5, accent=col)
        add_textbox(sl, x + 0.2, 3.35, 2.75, 0.35,
                    title, font_size=11, bold=True, color=TEXT_WHITE)
        add_textbox(sl, x + 0.2, 3.72, 2.75, 0.95,
                    body, font_size=9.5, color=TEXT_MUTE)

    # Screenshot — report dashboard
    img = os.path.join(SCREENSHOT_DIR, "01_overview.png")
    add_screenshot(sl, img, 0.3, 4.95, 12.7, 1.9)


def slide10_summary(prs):
    """Summary — Why Magic Migrate + Hyper-V"""
    sl = blank_slide(prs)
    slide_bg(sl)
    # Full header
    add_rect(sl, 0, 0, 13.333, 1.05, RGBColor(0x09, 0x0D, 0x14))
    add_rect(sl, 0, 1.05, 13.333, 0.04, WIPRO_VIOLET)
    add_textbox(sl, 0.35, 0.12, 10.5, 0.55,
                "Why Magic Migrate + Microsoft Hyper-V?",
                font_size=22, bold=True, color=TEXT_WHITE)
    add_textbox(sl, 0.35, 0.63, 10.5, 0.35,
                "Summary of capabilities and business value",
                font_size=11.5, color=TEXT_MUTE)
    if os.path.isfile(WIPRO_LOGO):
        sl.shapes.add_picture(WIPRO_LOGO, Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.45))
    footer(sl, 10)

    # 3 big value columns
    cols = [
        ("🚀", "Zero-Downtime Migration", WIPRO_VIOLET, [
            "5-step guided wizard — no CLI, no scripts",
            "Pre-flight blocks unsafe migrations",
            "Batch via Move Groups — app-tier aware",
            "Live event streaming — full visibility",
            "Automated post-migration Ansible tasks",
        ]),
        ("💰", "Cost & Licensing Savings", ACCENT_GREEN, [
            "Hyper-V free with Windows Server DC license",
            "Eliminate VMware per-core renewal costs",
            "CSVs on existing Windows storage — no vSAN",
            "Integrated with Azure Stack HCI roadmap",
            "No 3rd-party migration tool licences needed",
        ]),
        ("🛡️", "Enterprise-Grade Confidence", ACCENT_BLUE, [
            "12-point automated compatibility report",
            "VHDX = Azure-native disk format",
            "Shielded VMs + TPM 2.0 + Secure Boot",
            "Full audit trail in CaaS platform",
            "Rollback: re-migrate back to vSphere",
        ]),
    ]
    for ci, (icon, title, col, bullets) in enumerate(cols):
        x = 0.28 + ci * 4.34
        add_rect(sl, x, 1.2, 4.1, 5.55, CARD_BG)
        add_rect(sl, x, 1.2, 4.1, 0.05, col)
        add_textbox(sl, x + 0.15, 1.3, 3.8, 0.38,
                    f"{icon}  {title}", font_size=13, bold=True, color=col)
        for bi, bullet in enumerate(bullets):
            add_textbox(sl, x + 0.2, 1.82 + bi * 0.88, 3.7, 0.38,
                        f"▶  {bullet}", font_size=10.5, color=TEXT_WHITE)

    # Bottom CTA
    add_rect(sl, 0.28, 6.88, 12.77, 0.42, RGBColor(0x1E, 0x1B, 0x4B))
    add_textbox(sl, 0.5, 6.91, 12.3, 0.3,
                "📧  Contact: CaaS Platform Team  |  🌐  Portal: https://caas.wipro.internal  |  📄  Full feature docs available on request",
                font_size=9.5, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    prs = new_prs()

    print("Building slides...")
    slide1_cover(prs)           ; print("  ✅ Slide 1 — Cover")
    slide2_problem(prs)         ; print("  ✅ Slide 2 — Why Migrate")
    slide3_overview(prs)        ; print("  ✅ Slide 3 — Feature Overview")
    slide4_wizard_screenshot(prs); print("  ✅ Slide 4 — Wizard Screenshot")
    slide5_preflight(prs)       ; print("  ✅ Slide 5 — Pre-flight Checks")
    slide6_resource_mapping(prs); print("  ✅ Slide 6 — Resource Mapping")
    slide7_hyperv_screenshot(prs); print("  ✅ Slide 7 — Hyper-V Dashboard")
    slide8_move_groups(prs)     ; print("  ✅ Slide 8 — Move Groups")
    slide9_plan_management(prs) ; print("  ✅ Slide 9 — Plan Management")
    slide10_summary(prs)        ; print("  ✅ Slide 10 — Summary")

    out = r"C:\caas-dashboard\Magic_Migrate_VMware_to_HyperV.pptx"
    prs.save(out)
    print(f"\n✅ Saved: {out}")

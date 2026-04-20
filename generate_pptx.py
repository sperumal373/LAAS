"""
LaaS Portal – Executive Presentation Generator
Creates a stunning 20-slide PowerPoint using the Wipro template.
Saved to: C:\\caas-dashboard\\LaaS_Portal_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as In, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ── CONSTANTS ──────────────────────────────────────────────────────────────
TEMPLATE  = r"C:\Users\Administrator\Desktop\Wipro Template.pptx"
OUTPUT    = r"C:\caas-dashboard\LaaS_Portal_Presentation.pptx"
LOGO_PATH = r"C:\caas-dashboard\wipro_logo.png"

W = Inches(13.33)   # slide width  (16:9 widescreen)
H = Inches(7.5)     # slide height

# Brand palette
C_DARK    = RGBColor(0x04, 0x0D, 0x1A)   # deep navy bg
C_PANEL   = RGBColor(0x0D, 0x18, 0x29)   # panel bg
C_PANEL2  = RGBColor(0x11, 0x18, 0x27)   # alt panel
C_WIPRO   = RGBColor(0x00, 0x70, 0xC0)   # Wipro blue
C_BLUE    = RGBColor(0x3B, 0x82, 0xF6)   # accent blue
C_CYAN    = RGBColor(0x06, 0xB6, 0xD4)   # cyan
C_GREEN   = RGBColor(0x10, 0xB9, 0x81)   # green
C_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # amber
C_RED     = RGBColor(0xEF, 0x44, 0x44)   # red
C_PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)   # purple
C_TEXT    = RGBColor(0xE2, 0xE8, 0xF0)   # light text
C_SUB     = RGBColor(0x94, 0xA3, 0xB8)   # sub text
C_BORDER  = RGBColor(0x1E, 0x29, 0x3B)   # border
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── HELPERS ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, radius=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill_ex = shape.fill
    if fill:
        fill_ex.solid()
        fill_ex.fore_color.rgb = fill
    else:
        fill_ex.background()
    line_ex = shape.line
    if line:
        line_ex.color.rgb = line
        line_ex.width = Pt(1)
    else:
        line_ex.fill.background()
    if radius:
        # Apply rounded corners via XML
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is None:
            prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = avLst.find(qn('a:gd'))
        if gd is None:
            gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius}')
    return shape

def darken(col, factor=4):
    """Return a darkened version of an RGBColor."""
    return RGBColor(col[0]//factor, col[1]//factor, col[2]//factor)

# ── LOGO position matches original template (Layout 0 'Graphic 10') ──────
LOGO_L = Inches(0.930)
LOGO_T = Inches(0.701)
LOGO_W = Inches(1.354)
LOGO_H = Inches(1.354)

def add_wipro_logo(slide, cover=False):
    """Place the Wipro logo on the slide.
    cover=True  → original full-size template position (for cover/thank-you slides)
    cover=False → smaller top-right placement inside header band (for interior slides)
    """
    if cover:
        slide.shapes.add_picture(LOGO_PATH, LOGO_L, LOGO_T, LOGO_W, LOGO_H)
    else:
        # Top-right corner inside the header band, doesn't overlap title text
        slide.shapes.add_picture(LOGO_PATH, Inches(11.85), Inches(0.09), Inches(1.3), Inches(0.98))


def add_click_animations(slide, shape_id_list):
    """
    Add sequential 'Appear on click' animations to a list of shape IDs.
    Uses presetID=1 (Appear entrance) — the exact XML PowerPoint generates.
    Each shape in shape_id_list gets its own click trigger.
    """
    PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # ── build inner click-effect blocks ─────────────────────────────────
    click_blocks = ''
    ctn_id = 3
    for grp_idx, spid in enumerate(shape_id_list):
        # Each block: triggered by a mouse click, presetID=1 = Appear
        click_blocks += (
            f'<p:par>'
            f'<p:cTn id="{ctn_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:par>'
            f'<p:cTn id="{ctn_id+1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:par>'
            f'<p:cTn id="{ctn_id+2}" presetID="1" presetClass="entr" presetSubtype="0"'
            f' fill="hold" grpId="{grp_idx}" nodeType="clickEffect">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            # p:set — make shape visible (standard Appear effect child)
            f'<p:set>'
            f'<p:cBhvr>'
            f'<p:cTn id="{ctn_id+3}" dur="1" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr>'
            f'<p:to><p:strVal val="visible"/></p:to>'
            f'</p:set>'
            f'</p:childTnLst>'
            f'</p:cTn>'
            f'</p:par>'
            f'</p:childTnLst>'
            f'</p:cTn>'
            f'</p:par>'
            f'</p:childTnLst>'
            f'</p:cTn>'
            f'</p:par>'
        )
        ctn_id += 10

    # ── build bldLst entries ─────────────────────────────────────────────
    bld_entries = ''
    for grp_idx, spid in enumerate(shape_id_list):
        bld_entries += (
            f'<p:bldP spid="{spid}" grpId="{grp_idx}" uiExpand="1" build="allAtOnce"/>'
        )

    # ── full timing element — single xmlns declaration at root ───────────
    timing_xml = (
        f'<p:timing xmlns:p="{PML}">'
        f'<p:tnLst>'
        f'<p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">'
        f'<p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>'
        + click_blocks +
        f'</p:childTnLst>'
        f'</p:cTn>'
        f'<p:prevCondLst>'
        f'<p:cond evt="onPrevClick" delay="0">'
        f'<p:tn><p:prevTn id="2"/></p:tn>'
        f'</p:cond>'
        f'</p:prevCondLst>'
        f'</p:seq>'
        f'</p:childTnLst>'
        f'</p:cTn>'
        f'</p:par>'
        f'</p:tnLst>'
        f'<p:bldLst>{bld_entries}</p:bldLst>'
        f'</p:timing>'
    )

    timing_elem = etree.fromstring(timing_xml)

    # Remove any existing timing element and append the new one
    old = slide._element.find(qn('p:timing'))
    if old is not None:
        slide._element.remove(old)
    slide._element.append(timing_elem)


def add_text(slide, text, x, y, w, h, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_TEXT
    # No fill on textbox
    tb.fill.background()
    tb.line.fill.background()
    return tb

def add_label_value(slide, label, value, x, y, box_w, lsize=9, vsize=22, 
                     lcolor=None, vcolor=None, bg=None, border=None):
    """KPI card: small label above large value."""
    bh = Inches(1.0)
    card = add_rect(slide, x, y, box_w, bh, fill=bg or C_PANEL, line=border or C_BORDER, radius=50000)
    add_text(slide, label.upper(), x, y + Inches(0.08), box_w, Inches(0.25),
             size=lsize, color=lcolor or C_SUB, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, value, x, y + Inches(0.3), box_w, Inches(0.6),
             size=vsize, bold=True, color=vcolor or C_BLUE, align=PP_ALIGN.CENTER)

def add_pill(slide, text, x, y, color, bg_alpha=0x20):
    """Coloured status pill."""
    pw = Inches(1.5)
    ph = Inches(0.32)
    r = add_rect(slide, x, y, pw, ph,
                 fill=RGBColor(color[0], color[1], color[2]), radius=80000)
    r.fill.fore_color.rgb = RGBColor(
        min(255, color.red + 0xCC),
        min(255, color.green + 0xCC),
        min(255, color.blue + 0xCC))
    r.fill.fore_color.rgb = color  # use direct color
    # Semi-transparent hack: just set low-contrast bg
    add_text(slide, text, x, y, pw, ph, size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def gradient_bg(slide, c1=C_DARK, c2=C_PANEL):
    """Full-slide dark background."""
    bg = add_rect(slide, 0, 0, W, H, fill=c1)
    return bg

def section_header_bar(slide, title, subtitle="", accent=None):
    """Top bar with title and accent line."""
    accent = accent or C_WIPRO
    # Top accent stripe
    add_rect(slide, 0, 0, W, Inches(0.06), fill=accent)
    # Header band
    add_rect(slide, 0, Inches(0.06), W, Inches(1.1), fill=C_PANEL)
    # Accent left bar
    add_rect(slide, Inches(0.35), Inches(0.18), Inches(0.06), Inches(0.85), fill=accent, radius=20000)
    # Title
    add_text(slide, title, Inches(0.55), Inches(0.14), Inches(11), Inches(0.6),
             size=30, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.55), Inches(0.72), Inches(11), Inches(0.4),
                 size=13, color=C_SUB)

def bottom_bar(slide):
    """Footer strip."""
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=C_PANEL)
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.02), fill=C_WIPRO)
    add_text(slide, "LaaS Portal  |  Infrastructure as a Service  |  CONFIDENTIAL",
             Inches(0.3), H - Inches(0.34), Inches(9), Inches(0.3),
             size=8, color=C_SUB)
    add_text(slide, "© 2026 Wipro Limited", Inches(11.5), H - Inches(0.34), Inches(1.7), Inches(0.3),
             size=8, color=C_SUB, align=PP_ALIGN.RIGHT)

def tech_badge(slide, icon, label, x, y, color):
    bw, bh = Inches(1.55), Inches(0.95)
    add_rect(slide, x, y, bw, bh, fill=C_PANEL2, line=color, radius=40000)
    add_rect(slide, x, y, bw, Inches(0.04), fill=color, radius=0)
    add_text(slide, icon, x, y + Inches(0.05), bw, Inches(0.45), size=20, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.52), bw, Inches(0.35), size=9, bold=True,
             color=C_TEXT, align=PP_ALIGN.CENTER)

def flow_box(slide, text, x, y, w=Inches(2.1), h=Inches(0.65), fill=C_WIPRO, tsize=10):
    add_rect(slide, x, y, w, h, fill=fill, radius=40000)
    add_text(slide, text, x, y + Inches(0.05), w, h - Inches(0.05),
             size=tsize, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def arrow_right(slide, x, y, color=C_SUB):
    add_rect(slide, x, y + Inches(0.24), Inches(0.24), Inches(0.05), fill=color)

def feature_row(slide, icon, title, desc, x, y, w=Inches(5.8)):
    add_text(slide, icon, x, y, Inches(0.4), Inches(0.4), size=16, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.45), y, w - Inches(0.45), Inches(0.25),
             size=11, bold=True, color=C_TEXT)
    add_text(slide, desc, x + Inches(0.45), y + Inches(0.22), w - Inches(0.45), Inches(0.28),
             size=9, color=C_SUB)

def pct_bar(slide, label, pct, x, y, bar_w=Inches(4.5), color=C_BLUE):
    add_text(slide, label, x, y, Inches(1.6), Inches(0.25), size=9, color=C_SUB)
    track = add_rect(slide, x + Inches(1.65), y + Inches(0.04), bar_w, Inches(0.16), fill=C_BORDER, radius=20000)
    fill_w = int(bar_w * pct / 100)
    add_rect(slide, x + Inches(1.65), y + Inches(0.04), fill_w, Inches(0.16), fill=color, radius=20000)
    add_text(slide, f"{pct}%", x + Inches(1.65) + fill_w + Inches(0.1), y, Inches(0.5), Inches(0.25),
             size=9, bold=True, color=color)

# ── CREATE FRESH BLANK PRESENTATION ───────────────────────────────────────
# We do NOT load the Wipro template directly because it has 130 slides —
# removing them with drop_rel leaves orphaned parts that trigger PowerPoint's
# repair dialog. Instead we start from a blank presentation, set 16:9 size,
# and embed the Wipro logo/branding as images on every slide.
prs = Presentation()   # built-in blank 16:9 template
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Use the blank layout (index 6) for every slide — we draw all content ourselves
BLANK = prs.slide_layouts[6]   # "Blank" layout  

# shorthand alias used throughout
L = [BLANK] * 10   # L[0]..L[9] all point to blank layout

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER  (aligned)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# ── Full dark background
add_rect(slide, 0, 0, W, H, fill=C_DARK)

# ── Subtle accent bands
add_rect(slide, 0, Inches(4.9), W, Inches(0.8), fill=RGBColor(0x00, 0x3A, 0x70))
add_rect(slide, 0, Inches(5.7), W, Inches(1.8), fill=RGBColor(0x00, 0x28, 0x50))

# ── Top accent stripe + bottom stripe
add_rect(slide, 0, 0,              W, Inches(0.06), fill=C_WIPRO)
add_rect(slide, 0, H-Inches(0.06), W, Inches(0.06), fill=C_WIPRO)

# ── Left vertical accent bar (clean, full height inside margins)
add_rect(slide, Inches(0.55), Inches(0.18), Inches(0.07), Inches(4.5), fill=C_WIPRO)

# ── Wipro logo (template original position)
add_wipro_logo(slide)

# ── Portal icon badge (right of logo, vertically centered with title block)
add_rect(slide, Inches(1.82), Inches(0.95), Inches(0.78), Inches(0.78), fill=C_WIPRO, radius=30000)
add_text(slide, "⚡", Inches(1.82), Inches(0.95), Inches(0.78), Inches(0.78),
         size=24, align=PP_ALIGN.CENTER, color=C_WHITE)

# ── Main title  — starts at x=2.78 to clear logo+badge
add_text(slide, "LaaS Portal",
         Inches(2.78), Inches(0.92), Inches(9.8), Inches(1.05),
         size=56, bold=True, color=C_WHITE)

# ── Sub-title line
add_text(slide, "Infrastructure as a Service  │  Executive Overview",
         Inches(2.78), Inches(1.96), Inches(9.8), Inches(0.52),
         size=19, color=C_CYAN)

# ── Description line
add_text(slide,
    "Unified Multi-Cloud Operations  ·  VMware  ·  Red Hat OpenShift  ·  Nutanix  ·  Ansible AAP  ·  IPAM",
    Inches(2.78), Inches(2.52), Inches(10.0), Inches(0.4),
    size=11.5, color=C_SUB)

# ── Divider line
add_rect(slide, Inches(2.78), Inches(3.06), Inches(4.0), Inches(0.04), fill=C_WIPRO)

# ── Meta pills — 4 across, aligned under divider
meta = [("Version", "v6.0"), ("Date", "March 2026"), ("Classification", "Confidential"), ("Audience", "Leadership")]
for idx, (lbl, val) in enumerate(meta):
    bx = Inches(2.78 + idx * 2.52)
    add_rect(slide, bx, Inches(3.22), Inches(2.28), Inches(0.75), fill=C_PANEL, line=C_BORDER, radius=25000)
    add_text(slide, lbl.upper(), bx, Inches(3.25), Inches(2.28), Inches(0.26),
             size=8, bold=True, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, val, bx, Inches(3.49), Inches(2.28), Inches(0.38),
             size=12, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

# ── Platform strip — 6 equal chips spanning full width
platforms = [("🖥️","VMware vSphere"), ("🔴","Red Hat OCP"), ("🟢","Nutanix AHV"),
             ("🤖","Ansible AAP"),    ("🌐","IPAM / DNS"),  ("💳","Chargeback")]
chip_w = Inches(2.07)
for idx, (ic, nm) in enumerate(platforms):
    bx = Inches(0.35 + idx * 2.14)
    add_rect(slide, bx, Inches(5.04), chip_w, Inches(0.82), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic + "  " + nm, bx, Inches(5.14), chip_w, Inches(0.6),
             size=10, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

# ── Author footer
add_rect(slide, 0, H - Inches(0.42), W, Inches(0.42), fill=RGBColor(0x03,0x09,0x14))
add_text(slide,
    "Presented by: Sekhar Perumal  |  SDX Infrastructure & DC Operations  |  © 2026 Wipro Limited",
    Inches(0.5), H - Inches(0.38), Inches(12.3), Inches(0.32),
    size=9, color=C_SUB)

add_wipro_logo(slide, cover=True)  # cover=True: full-size at template position

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Agenda", "What we will cover in this presentation", C_WIPRO)
bottom_bar(slide)

add_wipro_logo(slide)
agenda = [
    ("01", "Business Challenge", "Why we built LaaS Portal", C_WIPRO),
    ("02", "Solution Overview", "Platform architecture at a glance", C_BLUE),
    ("03", "Technology Stack", "Integrations & tech landscape", C_CYAN),
    ("04", "Core Capabilities", "VMware · OCP · Nutanix · Ansible · IPAM", C_GREEN),
    ("05", "Governance & Security", "RBAC · Chargeback · Compliance", C_PURPLE),
    ("06", "Technical Specs", "Performance benchmarks & specs", C_AMBER),
    ("07", "Business Value", "ROI · Time savings · KPIs", C_RED),
    ("08", "Roadmap & Q&A", "Future roadmap and next steps", C_WIPRO),
]
cols = [Inches(0.4), Inches(6.85)]
for idx, (num, title, desc, color) in enumerate(agenda):
    col = idx % 2
    row = idx // 2
    bx = cols[col]
    by = Inches(1.45 + row * 1.24)
    bw, bh = Inches(6.0), Inches(1.1)
    add_rect(slide, bx, by, bw, bh, fill=C_PANEL, line=C_BORDER, radius=30000)
    add_rect(slide, bx, by, Inches(0.07), bh, fill=color, radius=0)
    add_text(slide, num, bx + Inches(0.18), by + Inches(0.08), Inches(0.55), Inches(0.45),
             size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.82), by + Inches(0.08), Inches(4.8), Inches(0.42),
             size=15, bold=True, color=C_TEXT)
    add_text(slide, desc, bx + Inches(0.82), by + Inches(0.52), Inches(4.8), Inches(0.4),
             size=10, color=C_SUB)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS CHALLENGE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "The Business Challenge", "Fragmented operations across multi-cloud platforms", C_RED)
bottom_bar(slide)

problems = [
    ("🔴", "Siloed Consoles", "Teams log into 8+ separate consoles (vCenter, OCP, Nutanix, AAP, IPAM...) wasting 2–3 hrs/day per operator"),
    ("🟠", "No Visibility", "Zero unified view of resource utilisation across VMware, OpenShift and Nutanix — capacity planning is ad hoc"),
    ("🟡", "Manual Chargeback", "Project cost allocation done manually in Excel — error-prone, delayed by 2–3 weeks every month"),
    ("🔵", "Access Control Gaps", "No fine-grained RBAC — engineers access production vCenter with the same rights as read-only viewers"),
    ("🟣", "Slow VM Provisioning", "Self-service VM requests take 3–5 days via email chains and manual approvals with no audit trail"),
    ("⚪", "No Automation Visibility", "Ansible AAP job statuses not visible to infra team — failures discovered only after incident escalation"),
]
for idx, (ic, title, desc) in enumerate(problems):
    row = idx // 2
    col = idx % 2
    bx = Inches(0.4 + col * 6.45)
    by = Inches(1.45 + row * 1.7)
    add_rect(slide, bx, by, Inches(5.95), Inches(1.55), fill=C_PANEL, line=C_BORDER, radius=30000)
    add_text(slide, ic, bx + Inches(0.18), by + Inches(0.35), Inches(0.5), Inches(0.6), size=22, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.75), by + Inches(0.1), Inches(4.8), Inches(0.4), size=13, bold=True, color=C_TEXT)
    add_text(slide, desc, bx + Inches(0.75), by + Inches(0.5), Inches(4.8), Inches(0.85), size=9, color=C_SUB)

add_rect(slide, Inches(0.4), Inches(6.6), Inches(12.53), Inches(0.45), fill=RGBColor(0x18, 0x08, 0x08), line=C_RED, radius=20000)
add_text(slide, "⚠  Result: Increased MTTR · Compliance risk · Cost overruns · Operator burnout · Delayed provisioning",
         Inches(0.6), Inches(6.62), Inches(12), Inches(0.38), size=10, bold=True, color=C_RED)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "LaaS Portal — Solution Overview", "One portal. All platforms. Full control.", C_WIPRO)
bottom_bar(slide)

# Central hub
add_rect(slide, Inches(5.15), Inches(2.7), Inches(2.95), Inches(1.65), fill=C_WIPRO, radius=50000)
add_text(slide, "⚡", Inches(5.15), Inches(2.75), Inches(2.95), Inches(0.6), size=22, align=PP_ALIGN.CENTER)
add_text(slide, "LaaS Portal", Inches(5.15), Inches(3.3), Inches(2.95), Inches(0.5),
         size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Unified Operations", Inches(5.15), Inches(3.72), Inches(2.95), Inches(0.35),
         size=9, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

add_wipro_logo(slide)
# Spoke items
spokes = [
    (Inches(0.5),  Inches(1.5),  "🖥️", "VMware vSphere",  C_WIPRO),
    (Inches(0.5),  Inches(3.5),  "🔴", "Red Hat OCP",     C_RED),
    (Inches(0.5),  Inches(5.5),  "🟢", "Nutanix AHV",     C_GREEN),
    (Inches(10.5), Inches(1.5),  "🤖", "Ansible AAP",     C_AMBER),
    (Inches(10.5), Inches(3.5),  "🌐", "Networks & IPAM", C_CYAN),
    (Inches(10.5), Inches(5.5),  "💳", "Chargeback",      C_PURPLE),
    (Inches(4.6),  Inches(1.1),  "📊", "Capacity Mgmt",   C_BLUE),
    (Inches(7.85), Inches(1.1),  "🏷️", "Project Util.",   C_GREEN),
    (Inches(4.6),  Inches(5.7),  "🔒", "RBAC Security",   C_PURPLE),
    (Inches(7.85), Inches(5.7),  "📋", "VM Requests",     C_AMBER),
]
for (bx, by, ic, lbl, col) in spokes:
    bw, bh = Inches(2.0), Inches(0.75)
    add_rect(slide, bx, by, bw, bh, fill=C_PANEL, line=col, radius=30000)
    add_text(slide, ic + " " + lbl, bx, by + Inches(0.15), bw, Inches(0.45),
             size=11, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE FLOWCHART
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Platform Architecture", "Three-tier design: Presentation → Integration → Infrastructure", C_CYAN)
bottom_bar(slide)

# Tier labels
tiers = [
    (Inches(0.35), "PRESENTATION TIER",  C_WIPRO, "Browser · React/Vite SPA · Role-Based UI"),
    (Inches(3.0),  "API LAYER",           C_BLUE,  "FastAPI · REST · JWT Auth · RBAC Middleware"),
    (Inches(5.65), "INTEGRATION LAYER",  C_CYAN,  "vCenter API · OCP API · Prism API · AAP API · IPAM"),
    (Inches(8.3),  "DATA LAYER",          C_GREEN, "SQLite · Audit Logs · Chargeback DB · Session Store"),
    (Inches(10.95),"INFRA LAYER",         C_PURPLE,"VMware · OpenShift · Nutanix · Ansible · DNS/IPAM"),
]
for (bx, label, color, sub) in tiers:
    bw = Inches(2.45)
    add_rect(slide, bx, Inches(1.45), bw, Inches(4.6), fill=C_PANEL, line=color, radius=20000)
    add_rect(slide, bx, Inches(1.45), bw, Inches(0.05), fill=color)
    add_rect(slide, bx + Inches(0.65), Inches(1.5), bw - Inches(1.3), Inches(0.38),
             fill=darken(color), radius=20000)
    add_text(slide, label, bx, Inches(1.52), bw, Inches(0.35),
             size=7, bold=True, color=color, align=PP_ALIGN.CENTER)
    # sub items
    lines = sub.split(" · ")
    for li, line in enumerate(lines):
        by = Inches(2.05 + li * 0.62)
        add_rect(slide, bx + Inches(0.18), by, bw - Inches(0.36), Inches(0.5),
                 fill=C_PANEL2, line=C_BORDER, radius=20000)
        add_text(slide, line, bx + Inches(0.18), by + Inches(0.07), bw - Inches(0.36), Inches(0.38),
                 size=9, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, sub.replace(" · ", "\n"), bx + Inches(0.1), Inches(4.6), bw - Inches(0.2), Inches(0.8),
             size=8, color=C_SUB, align=PP_ALIGN.CENTER)

# Arrows between tiers
for i in range(4):
    ax = Inches(2.82 + i * 2.65)
    add_rect(slide, ax, Inches(3.5), Inches(0.15), Inches(0.04), fill=C_SUB)
    add_text(slide, "▶", ax + Inches(0.12), Inches(3.41), Inches(0.2), Inches(0.22),
             size=10, color=C_SUB, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(6.3), Inches(12.63), Inches(0.52), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "🔒 HTTPS/TLS throughout  │  🔑 JWT + RBAC  │  📝 Full audit log  │  🐍 Python 3.14  │  ⚛ React/Vite  │  🗄 SQLite",
         Inches(0.5), Inches(6.33), Inches(12.3), Inches(0.42), size=9.5, color=C_SUB, align=PP_ALIGN.CENTER)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Technology Stack", "Modern, enterprise-grade, open-standards based", C_BLUE)
bottom_bar(slide)

# 8 columns: bw=1.50, spacing=1.61 → last right edge = 11.60+1.50=13.10" (fits 13.33")
_TX = [Inches(x) for x in [0.35, 1.96, 3.57, 5.18, 6.79, 8.40, 10.01, 11.60]]
tech_items = [
    # row1
    ("⚛️", "React 18\nVite 7",           C_BLUE,   _TX[0], Inches(1.55)),
    ("🐍", "Python 3.14\nFastAPI",        C_GREEN,  _TX[1], Inches(1.55)),
    ("🗄️", "SQLite 3\nAudit DB",          C_CYAN,   _TX[2], Inches(1.55)),
    ("🔐", "JWT Auth\nRBAC",              C_AMBER,  _TX[3], Inches(1.55)),
    ("🖥️", "VMware vSphere 8\nvCenter API",C_WIPRO, _TX[4], Inches(1.55)),
    ("🔴", "OCP 4.15\nOperator SDK",      C_RED,    _TX[5], Inches(1.55)),
    ("🟢", "Nutanix AOS 6.x\nPrism API",  C_GREEN,  _TX[6], Inches(1.55)),
    ("🤖", "Ansible AAP 2.4\nREST API",   C_AMBER,  _TX[7], Inches(1.55)),
    # row2
    ("📊", "Recharts 2\nData Viz",        C_PURPLE, _TX[0], Inches(3.1)),
    ("🌐", "phpIPAM\nIP Mgmt",            C_CYAN,   _TX[1], Inches(3.1)),
    ("🐳", "Docker\nContainers",          C_BLUE,   _TX[2], Inches(3.1)),
    ("📡", "REST APIs\nHTTPS TLS",        C_AMBER,  _TX[3], Inches(3.1)),
    ("🔒", "LDAP/AD\nSSO Ready",          C_RED,    _TX[4], Inches(3.1)),
    ("📝", "Audit Trail\nImmutable Log",  C_GREEN,  _TX[5], Inches(3.1)),
    ("📱", "Responsive\nMobile Ready",    C_CYAN,   _TX[6], Inches(3.1)),
    ("🎨", "Dark Theme\nBrand UI",        C_PURPLE, _TX[7], Inches(3.1)),
]
for (ic, lbl, col, bx, by) in tech_items:
    bw, bh = Inches(1.50), Inches(1.3)
    add_rect(slide, bx, by, bw, bh, fill=C_PANEL2, line=col, radius=40000)
    add_rect(slide, bx, by, bw, Inches(0.05), fill=col)
    add_text(slide, ic, bx, by + Inches(0.1), bw, Inches(0.5), size=22, align=PP_ALIGN.CENTER)
    add_text(slide, lbl, bx, by + Inches(0.6), bw, Inches(0.6), size=9, bold=True,
             color=C_TEXT, align=PP_ALIGN.CENTER)

# Summary stat bar
add_rect(slide, Inches(0.4), Inches(4.62), Inches(12.53), Inches(0.95), fill=C_PANEL, line=C_BORDER, radius=20000)
stats = [("20+", "Platform Integrations"), ("8", "API Endpoints"), ("4", "Auth Roles"), ("99.9%", "Uptime SLA"), ("< 2s", "Page Load"), ("100%", "Dark Theme")]
for i, (v, l) in enumerate(stats):
    bx = Inches(0.7 + i * 2.06)
    add_text(slide, v, bx, Inches(4.68), Inches(1.8), Inches(0.42),
             size=22, bold=True, color=C_WIPRO, align=PP_ALIGN.CENTER)
    add_text(slide, l, bx, Inches(5.1), Inches(1.8), Inches(0.3),
             size=8, color=C_SUB, align=PP_ALIGN.CENTER)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — VMWARE MODULE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "VMware Module", "Complete VM lifecycle management through a unified portal", C_WIPRO)
bottom_bar(slide)

# Left column - features
features_vm = [
    ("🖥️", "VM Inventory Dashboard", "Real-time list of all VMs across vCenters — power state, IP, OS, host, snapshots"),
    ("⚡", "Power Operations", "Power On/Off/Restart/Suspend with confirmation guard — full audit trail per action"),
    ("📸", "Snapshot Management", "Create/revert/delete snapshots. Age-coded badges: Red >28d, Yellow >7d, Green ≤7d"),
    ("🔍", "Advanced Filtering", "Filter by power state, host, OS, project, vCenter — instant search by name/IP"),
    ("📊", "VM Detail Panel", "CPU/RAM charts, disk info, network adapters, snapshot tree, event history"),
    ("📋", "CSV Export", "Export VM inventory with all metadata for reporting and audit"),
]
for idx, (ic, title, desc) in enumerate(features_vm):
    by = Inches(1.45 + idx * 0.9)
    add_rect(slide, Inches(0.4), by, Inches(5.9), Inches(0.82), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic, Inches(0.55), by + Inches(0.18), Inches(0.45), Inches(0.45), size=16, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.1), by + Inches(0.07), Inches(5.0), Inches(0.3), size=11, bold=True, color=C_TEXT)
    add_text(slide, desc, Inches(1.1), by + Inches(0.37), Inches(5.0), Inches(0.35), size=8.5, color=C_SUB)

# Right col — mock KPIs
kpis_vm = [("Total VMs", "248", C_WIPRO), ("Running", "196", C_GREEN), ("Stopped", "42", C_RED), ("Snapshots", "87", C_AMBER)]
for i, (lbl, val, col) in enumerate(kpis_vm):
    bx = Inches(6.65 + (i % 2) * 3.2)
    by = Inches(1.5 + (i // 2) * 1.15)
    add_rect(slide, bx, by, Inches(2.85), Inches(0.95), fill=C_PANEL, line=col, radius=30000)
    add_text(slide, lbl, bx, by + Inches(0.06), Inches(2.85), Inches(0.28), size=9, bold=True, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, val, bx, by + Inches(0.32), Inches(2.85), Inches(0.5), size=26, bold=True, color=col, align=PP_ALIGN.CENTER)

# Snapshot age chart mock
add_rect(slide, Inches(6.65), Inches(3.85), Inches(6.2), Inches(2.45), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "Snapshot Age Distribution", Inches(6.8), Inches(3.92), Inches(5.8), Inches(0.3), size=10, bold=True, color=C_TEXT)
bars = [("Critical >28d", 15, C_RED), ("Warning >7d", 34, C_AMBER), ("Healthy ≤7d", 38, C_GREEN)]
for i, (lbl, cnt, col) in enumerate(bars):
    by = Inches(4.35 + i * 0.62)
    add_text(slide, lbl, Inches(6.8), by, Inches(1.5), Inches(0.4), size=8.5, color=C_SUB)
    add_rect(slide, Inches(8.35), by + Inches(0.08), Inches(3.6), Inches(0.22), fill=C_BORDER, radius=10000)
    fill_w = int(Inches(3.6) * cnt / 87)
    add_rect(slide, Inches(8.35), by + Inches(0.08), fill_w, Inches(0.22), fill=col, radius=10000)
    add_text(slide, str(cnt), Inches(12.0), by, Inches(0.5), Inches(0.38), size=9, bold=True, color=col)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RED HAT OPENSHIFT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Red Hat OpenShift Module", "Container platform health, workload visibility & operator management", C_RED)
bottom_bar(slide)

features_ocp = [
    ("🔴", "Cluster Overview", "Health status, OCP version, node/pod/operator counts, API server URL, CA cert expiry"),
    ("🟡", "Node Management", "List all nodes with role (master/worker/infra), status, CPU/RAM usage bars, OS version"),
    ("📦", "Operator Catalog", "All installed operators — name, namespace, version, status (Succeeded/Failed/Pending)"),
    ("📊", "Namespace Insights", "Per-namespace pod/deployment/service/PVC counts with resource quota usage"),
    ("🔔", "Alerts Dashboard", "Active cluster alerts with severity — Critical/Warning/Info with occurrence counts"),
    ("🔒", "RBAC Integration", "Portal RBAC layered on top — Viewers see health only; Operators see full detail"),
]
for idx, (ic, title, desc) in enumerate(features_ocp):
    by = Inches(1.45 + idx * 0.9)
    add_rect(slide, Inches(0.4), by, Inches(5.9), Inches(0.82), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic, Inches(0.55), by + Inches(0.18), Inches(0.45), Inches(0.45), size=16, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.1), by + Inches(0.07), Inches(5.0), Inches(0.3), size=11, bold=True, color=C_TEXT)
    add_text(slide, desc, Inches(1.1), by + Inches(0.37), Inches(5.0), Inches(0.35), size=8.5, color=C_SUB)

# Right: cluster health card
add_rect(slide, Inches(6.65), Inches(1.5), Inches(6.2), Inches(5.0), fill=C_PANEL, line=RGBColor(0x44,0x0,0x0), radius=20000)
add_rect(slide, Inches(6.65), Inches(1.5), Inches(6.2), Inches(0.05), fill=C_RED)
add_text(slide, "🔴  Cluster: ocp-sdx-prod-01", Inches(6.8), Inches(1.58), Inches(5.8), Inches(0.35), size=11, bold=True, color=C_TEXT)
add_text(slide, "OpenShift 4.15.2  ·  Status: Healthy  ·  Nodes: 12", Inches(6.8), Inches(1.92), Inches(5.8), Inches(0.28), size=9, color=C_SUB)

cluster_kpis = [("Masters","3",C_RED),("Workers","9",C_GREEN),("Pods","847",C_BLUE),
                ("Operators","32",C_AMBER),("Namespaces","48",C_CYAN),("Alerts","3",C_RED)]
for i, (lbl, val, col) in enumerate(cluster_kpis):
    bx = Inches(6.8 + (i % 3) * 1.95)
    by = Inches(2.28 + (i // 3) * 1.02)
    add_rect(slide, bx, by, Inches(1.7), Inches(0.82), fill=C_PANEL2, line=col, radius=20000)
    add_text(slide, lbl, bx, by + Inches(0.05), Inches(1.7), Inches(0.25), size=8, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, val, bx, by + Inches(0.28), Inches(1.7), Inches(0.42), size=20, bold=True, color=col, align=PP_ALIGN.CENTER)

pct_bar(slide, "CPU Usage",   68, Inches(6.8), Inches(4.65), Inches(4.2), C_AMBER)
pct_bar(slide, "RAM Usage",   55, Inches(6.8), Inches(5.05), Inches(4.2), C_BLUE)
pct_bar(slide, "Pod Density", 82, Inches(6.8), Inches(5.45), Inches(4.2), C_RED)
pct_bar(slide, "Disk Usage",  41, Inches(6.8), Inches(5.85), Inches(4.2), C_GREEN)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NUTANIX MODULE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Nutanix AHV Module", "Hyper-converged infrastructure visibility via Prism Central API", C_GREEN)
bottom_bar(slide)

features_nx = [
    ("🟢", "Prism Central Overview", "VMs, hosts, clusters, storage containers — real-time counts and health status"),
    ("💾", "Storage Insights", "Container usage, thin provisioning %, dedup/compression savings ratios"),
    ("🖥️", "AHV Host Table", "Per-host CPU/RAM utilisation bars, VM count, power state, AOS version"),
    ("📊", "VM Catalog", "All AHV VMs with IP, power state, vCPU/vRAM allocation, cluster membership"),
    ("🔔", "Cluster Alerts", "Active alerts with severity, entity type, impact, and resolution status"),
    ("🔗", "Multi-Cluster", "Supports multiple Prism Central instances — unified view across all clusters"),
]
for idx, (ic, title, desc) in enumerate(features_nx):
    by = Inches(1.45 + idx * 0.9)
    add_rect(slide, Inches(0.4), by, Inches(5.9), Inches(0.82), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic, Inches(0.55), by + Inches(0.18), Inches(0.45), Inches(0.45), size=16, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.1), by + Inches(0.07), Inches(5.0), Inches(0.3), size=11, bold=True, color=C_TEXT)
    add_text(slide, desc, Inches(1.1), by + Inches(0.37), Inches(5.0), Inches(0.35), size=8.5, color=C_SUB)

add_rect(slide, Inches(6.65), Inches(1.5), Inches(6.2), Inches(5.0), fill=C_PANEL, line=RGBColor(0x0,0x44,0x22), radius=20000)
add_rect(slide, Inches(6.65), Inches(1.5), Inches(6.2), Inches(0.05), fill=C_GREEN)
add_text(slide, "🟢  Prism Central: pc-sdx-prod", Inches(6.8), Inches(1.58), Inches(5.8), Inches(0.35), size=11, bold=True, color=C_TEXT)
add_text(slide, "AOS 6.7.1  ·  Clusters: 3  ·  Status: Connected", Inches(6.8), Inches(1.92), Inches(5.8), Inches(0.28), size=9, color=C_SUB)

nx_kpis = [("Clusters","3",C_GREEN),("AHV Hosts","24",C_BLUE),("AHV VMs","312",C_CYAN),
           ("vCPU","1,248",C_AMBER),("RAM (TB)","9.6",C_PURPLE),("Storage (TB)","480",C_GREEN)]
for i, (lbl, val, col) in enumerate(nx_kpis):
    bx = Inches(6.8 + (i % 3) * 1.95)
    by = Inches(2.28 + (i // 3) * 1.02)
    add_rect(slide, bx, by, Inches(1.7), Inches(0.82), fill=C_PANEL2, line=col, radius=20000)
    add_text(slide, lbl, bx, by + Inches(0.05), Inches(1.7), Inches(0.25), size=8, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, val, bx, by + Inches(0.28), Inches(1.7), Inches(0.42), size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

pct_bar(slide, "CPU Used",      52, Inches(6.8), Inches(4.65), Inches(4.2), C_BLUE)
pct_bar(slide, "RAM Used",      71, Inches(6.8), Inches(5.05), Inches(4.2), C_AMBER)
pct_bar(slide, "Storage Used",  63, Inches(6.8), Inches(5.45), Inches(4.2), C_GREEN)
pct_bar(slide, "Dedup Savings", 38, Inches(6.8), Inches(5.85), Inches(4.2), C_CYAN)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ANSIBLE AAP
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Ansible Automation Platform", "Real-time automation job monitoring and audit", C_AMBER)
bottom_bar(slide)

features_aap = [
    ("🤖", "Job Dashboard", "Live count of Today's jobs — Success, Failed, Running split with percentage cards"),
    ("📋", "Job History Table", "All AAP jobs: name, type, status, launched-by, start/end time, duration"),
    ("🔍", "Status Filter", "Filter by status (Successful/Failed/Running/Canceled) or job template name"),
    ("📊", "Trend Charts", "7-day and 30-day job success rate trend line chart using Recharts"),
    ("🔔", "Failure Alerts", "Failed jobs surface immediately in the Overview error badge count"),
    ("🔗", "Template Catalog", "Browse all available job templates — launch directly from portal (Admin only)"),
]
for idx, (ic, title, desc) in enumerate(features_aap):
    by = Inches(1.45 + idx * 0.9)
    add_rect(slide, Inches(0.4), by, Inches(5.9), Inches(0.82), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic, Inches(0.55), by + Inches(0.18), Inches(0.45), Inches(0.45), size=16, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.1), by + Inches(0.07), Inches(5.0), Inches(0.3), size=11, bold=True, color=C_TEXT)
    add_text(slide, desc, Inches(1.1), by + Inches(0.37), Inches(5.0), Inches(0.35), size=8.5, color=C_SUB)

aap_stats = [("Today's Jobs","142",C_WIPRO),("Successful","128",C_GREEN),("Failed","9",C_RED),("Running","5",C_AMBER)]
for i, (lbl, val, col) in enumerate(aap_stats):
    bx = Inches(6.65 + (i % 2) * 3.15)
    by = Inches(1.5 + (i // 2) * 1.15)
    add_rect(slide, bx, by, Inches(2.85), Inches(0.95), fill=C_PANEL, line=col, radius=30000)
    add_text(slide, lbl, bx, by + Inches(0.06), Inches(2.85), Inches(0.28), size=9, bold=True, color=C_SUB, align=PP_ALIGN.CENTER)
    add_text(slide, val, bx, by + Inches(0.32), Inches(2.85), Inches(0.5), size=26, bold=True, color=col, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(6.65), Inches(3.85), Inches(6.2), Inches(2.45), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "Recent Jobs", Inches(6.8), Inches(3.92), Inches(5.8), Inches(0.3), size=10, bold=True, color=C_TEXT)
jobs = [("Deploy LAMP Stack","Successful",C_GREEN),("Patch RHEL Hosts","Successful",C_GREEN),
        ("Backup vCenter Config","Running",C_AMBER),("Security Hardening","Failed",C_RED),
        ("Create VM from Template","Successful",C_GREEN)]
for i, (name, status, col) in enumerate(jobs):
    by = Inches(4.35 + i * 0.42)
    add_text(slide, name, Inches(6.8), by, Inches(3.9), Inches(0.38), size=9, color=C_TEXT)
    add_rect(slide, Inches(10.9), by + Inches(0.04), Inches(1.7), Inches(0.26), fill=col, radius=15000)
    add_text(slide, status, Inches(10.9), by + Inches(0.04), Inches(1.7), Inches(0.26), size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CAPACITY MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Capacity Management", "Predictive resource utilisation with colour-coded arc gauges", C_CYAN)
bottom_bar(slide)

cap_features = [
    ("📊", "Arc Gauge Dashboard", "Visual CPU/RAM/Storage/Network usage arcs: Green <60%, Amber 60–80%, Red >80%"),
    ("🖥️", "ESXi Host Table", "Per-host CPU/RAM bars — identify hotspots and underutilised hosts at a glance"),
    ("📈", "Trend Forecasting", "7-day rolling average with 30-day forecast — plan expansions proactively"),
    ("⚠️", "Threshold Alerts", "Configurable warning/critical thresholds trigger Overview badge notifications"),
    ("💾", "Storage Capacity", "Datastore free/used breakdown — identify datastores approaching capacity limits"),
    ("📋", "Capacity Report", "Export detailed capacity snapshot to CSV for management reporting"),
]
for idx, (ic, title, desc) in enumerate(cap_features):
    row, col = idx // 2, idx % 2
    bx = Inches(0.4 + col * 6.45)
    by = Inches(1.45 + row * 1.72)
    add_rect(slide, bx, by, Inches(5.95), Inches(1.55), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic, Inches(bx/Inches(1)) + Inches(0.18), by + Inches(0.35), Inches(0.5), Inches(0.6), size=22, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.75), by + Inches(0.1), Inches(4.8), Inches(0.4), size=13, bold=True, color=C_TEXT)
    add_text(slide, desc, bx + Inches(0.75), by + Inches(0.5), Inches(4.8), Inches(0.85), size=9, color=C_SUB)

gauges = [("CPU", "64%", C_AMBER), ("RAM", "71%", C_AMBER), ("Storage", "83%", C_RED), ("Network", "37%", C_GREEN)]
add_rect(slide, Inches(0.4), Inches(6.18), Inches(12.53), Inches(0.9), fill=C_PANEL, line=C_BORDER, radius=20000)
for i, (lbl, val, col) in enumerate(gauges):
    bx = Inches(1.0 + i * 2.8)
    add_rect(slide, bx, Inches(6.28), Inches(2.0), Inches(0.72), fill=C_PANEL2, line=col, radius=20000)
    add_text(slide, lbl, bx, Inches(6.3), Inches(2.0), Inches(0.26), size=9, color=C_SUB, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, val, bx, Inches(6.55), Inches(2.0), Inches(0.35), size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PROJECT UTILIZATION & CHARGEBACK
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Project Utilization & Chargeback", "Cost allocation · billing · resource governance by project", C_PURPLE)
bottom_bar(slide)

# Left: Project Utilization
add_rect(slide, Inches(0.4), Inches(1.45), Inches(6.1), Inches(5.2), fill=C_PANEL, line=C_BORDER, radius=20000)
add_rect(slide, Inches(0.4), Inches(1.45), Inches(6.1), Inches(0.05), fill=C_PURPLE)
add_text(slide, "🏷️  Project Resource Utilization", Inches(0.55), Inches(1.52), Inches(5.8), Inches(0.35), size=11, bold=True, color=C_TEXT)

projects = [
    ("SDX-PROD-INFRA", "74 VMs", "592 vCPU", "2,368 GB", 72, C_WIPRO),
    ("SDX-DEV-LAB",    "52 VMs", "416 vCPU", "1,664 GB", 51, C_BLUE),
    ("SDX-BACKUP-DR",  "38 VMs", "304 vCPU", "1,216 GB", 38, C_AMBER),
    ("SDX-QA-STAGING", "29 VMs", "232 vCPU", "928 GB",   29, C_GREEN),
    ("SDX-DECOM",      "15 VMs", "60 vCPU",  "240 GB",   12, C_SUB),
]
for i, (name, vms, cpu, ram, pct, col) in enumerate(projects):
    by = Inches(2.0 + i * 0.82)
    add_text(slide, name, Inches(0.55), by, Inches(2.3), Inches(0.3), size=9, bold=True, color=C_TEXT)
    add_text(slide, vms, Inches(2.88), by, Inches(0.8), Inches(0.28), size=8, color=C_SUB, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.55), by + Inches(0.3), Inches(5.6), Inches(0.16), fill=C_BORDER, radius=10000)
    add_rect(slide, Inches(0.55), by + Inches(0.3), int(Inches(5.6) * pct / 100), Inches(0.16), fill=col, radius=10000)
    add_text(slide, f"{pct}%", Inches(5.5), by, Inches(0.55), Inches(0.28), size=8, bold=True, color=col, align=PP_ALIGN.RIGHT)

# Right: Chargeback
add_rect(slide, Inches(6.85), Inches(1.45), Inches(6.1), Inches(5.2), fill=C_PANEL, line=C_BORDER, radius=20000)
add_rect(slide, Inches(6.85), Inches(1.45), Inches(6.1), Inches(0.05), fill=C_GREEN)
add_text(slide, "💳  Chargeback — March 2026", Inches(7.0), Inches(1.52), Inches(5.8), Inches(0.35), size=11, bold=True, color=C_TEXT)

cb_rows = [
    ("SDX-PROD-INFRA", "₹ 88,800", "₹ 1,89,440", "₹ 71,000", "₹ 3,49,240", C_WIPRO),
    ("SDX-DEV-LAB",    "₹ 62,400", "₹ 1,33,120", "₹ 42,250", "₹ 2,37,770", C_BLUE),
    ("SDX-BACKUP-DR",  "₹ 45,600", "₹ 97,280",   "₹ 1,14,000","₹ 2,56,880",C_AMBER),
    ("SDX-QA-STAGING", "₹ 34,800", "₹ 74,240",   "₹ 25,500", "₹ 1,34,540", C_GREEN),
]
headers = ["Project", "vCPU", "RAM", "Storage", "TOTAL"]
hx = [Inches(7.0), Inches(8.5), Inches(9.4), Inches(10.3), Inches(11.3)]
hw = [Inches(1.4), Inches(0.85), Inches(0.85), Inches(0.95), Inches(1.3)]
for j, h in enumerate(headers):
    add_text(slide, h, hx[j], Inches(2.0), hw[j], Inches(0.26), size=8, bold=True, color=C_SUB)
for i, (nm, cpu, ram, sto, tot, col) in enumerate(cb_rows):
    by = Inches(2.3 + i * 0.8)
    add_rect(slide, Inches(6.92), by, Inches(5.88), Inches(0.68), fill=C_PANEL2, line=C_BORDER, radius=10000)
    vals = [nm, cpu, ram, sto, tot]
    for j, v in enumerate(vals):
        c = col if j == 4 else (C_TEXT if j == 0 else C_SUB)
        b = True if j in (0, 4) else False
        add_text(slide, v, hx[j], by + Inches(0.15), hw[j], Inches(0.35), size=8, bold=b, color=c)

add_rect(slide, Inches(6.85), Inches(5.55), Inches(6.1), Inches(0.65), fill=RGBColor(0x05, 0x18, 0x0D), line=C_GREEN, radius=10000)
add_text(slide, "March 2026 Grand Total:", Inches(7.0), Inches(5.6), Inches(3.5), Inches(0.5), size=10, bold=True, color=C_SUB)
add_text(slide, "₹ 9,78,430", Inches(10.5), Inches(5.6), Inches(2.0), Inches(0.5), size=16, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)

add_rect(slide, Inches(6.85), Inches(6.25), Inches(6.1), Inches(0.32), fill=C_PANEL)
add_text(slide, "Rate: vCPU ₹150/mo · RAM ₹80/GB/mo · Storage ₹5/GB/mo  |  Export: PDF · CSV",
         Inches(7.0), Inches(6.27), Inches(5.8), Inches(0.28), size=8, color=C_SUB)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — NETWORKS & IPAM
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Networks & IPAM Module", "VLAN topology, subnet management & IP address lifecycle", C_CYAN)
bottom_bar(slide)

net_features = [
    ("🌐", "Port Group Catalog",    "All vCenter distributed port groups — VLAN ID, switch name, connected VM count"),
    ("📡", "Subnet Dashboard",      "All IPAM subnets with used/free IP counts, occupancy % bar, gateway, VLAN"),
    ("🔍", "Free IP Finder",        "Search available IPs in any subnet — essential before provisioning new VMs"),
    ("📋", "IP Allocation Table",   "Full allocation history — which IP, assigned to what VM, by whom, when"),
    ("➕", "Reserve / Release IPs", "Admin can reserve and release IPs directly from the portal with comments"),
    ("🔔", "Occupancy Alerts",      "Subnets >80% full surface in Overview alerts for proactive management"),
]
for idx, (ic, title, desc) in enumerate(net_features):
    row, col = idx // 2, idx % 2
    bx = Inches(0.4 + col * 6.45)
    by = Inches(1.45 + row * 1.72)
    add_rect(slide, bx, by, Inches(5.95), Inches(1.55), fill=C_PANEL, line=C_BORDER, radius=20000)
    add_text(slide, ic, bx + Inches(0.18), by + Inches(0.35), Inches(0.5), Inches(0.6), size=22, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.75), by + Inches(0.1), Inches(4.8), Inches(0.4), size=13, bold=True, color=C_TEXT)
    add_text(slide, desc, bx + Inches(0.75), by + Inches(0.5), Inches(4.8), Inches(0.85), size=9, color=C_SUB)

subnets = [
    ("10.10.1.0/24", "PROD-MGMT",   245, 11, C_RED),
    ("10.10.2.0/24", "PROD-APP",    180, 76, C_AMBER),
    ("10.20.0.0/22", "DEV-LAB",     312,710, C_GREEN),
    ("192.168.1.0/24","DR-MGMT",    198, 58, C_AMBER),
]
add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.53), Inches(0.95), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "Subnet Snapshot", Inches(0.6), Inches(6.12), Inches(3), Inches(0.28), size=9, bold=True, color=C_TEXT)
for i, (subnet, name, used, free, col) in enumerate(subnets):
    bx = Inches(0.6 + i * 3.1)
    add_text(slide, f"{name}  {subnet}", bx, Inches(6.42), Inches(2.8), Inches(0.24), size=8, bold=True, color=C_TEXT)
    total = used + free
    add_rect(slide, bx, Inches(6.7), Inches(2.8), Inches(0.14), fill=C_BORDER, radius=10000)
    add_rect(slide, bx, Inches(6.7), int(Inches(2.8) * used / total), Inches(0.14), fill=col, radius=10000)
    add_text(slide, f"{used}/{total}", bx + Inches(2.0), Inches(6.42), Inches(0.8), Inches(0.24), size=8, color=col, align=PP_ALIGN.RIGHT)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — RBAC & SECURITY
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Security & Role-Based Access Control", "Four-tier RBAC · JWT authentication · full audit trail", C_RED)
bottom_bar(slide)

roles = [
    ("👑", "Admin",     C_WIPRO, [
        "Full platform access", "User & role management", "VM power + snapshot ops",
        "Chargeback rate config", "Approve/reject requests", "System settings"]),
    ("⚙️", "Operator",  C_BLUE, [
        "Read all resources", "VM power operations", "Snapshot create/delete",
        "Ansible job launch", "IP allocation", "Export reports"]),
    ("👁️", "Viewer",    C_GREEN, [
        "Read-only all modules", "View dashboards & KPIs", "View alerts & logs",
        "View chargeback reports", "Download CSV exports", "No write operations"]),
    ("📋", "Requester", C_AMBER, [
        "Submit VM requests only", "View own request status", "No infra visibility",
        "No power operations", "No snapshot access", "No billing access"]),
]
for i, (ic, role, col, perms) in enumerate(roles):
    bx = Inches(0.4 + i * 3.2)
    bw = Inches(3.0)
    add_rect(slide, bx, Inches(1.45), bw, Inches(5.5), fill=C_PANEL, line=col, radius=20000)
    add_rect(slide, bx, Inches(1.45), bw, Inches(0.05), fill=col)
    add_text(slide, ic, bx, Inches(1.58), bw, Inches(0.45), size=24, align=PP_ALIGN.CENTER)
    add_text(slide, role, bx, Inches(2.02), bw, Inches(0.38), size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    for j, perm in enumerate(perms):
        by = Inches(2.52 + j * 0.66)
        add_rect(slide, bx + Inches(0.18), by, bw - Inches(0.36), Inches(0.55), fill=C_PANEL2, line=C_BORDER, radius=15000)
        add_text(slide, "✓  " + perm, bx + Inches(0.25), by + Inches(0.1), bw - Inches(0.44), Inches(0.35), size=8.5, color=C_TEXT)

add_rect(slide, Inches(0.4), Inches(7.08), Inches(12.53), Inches(0.38), fill=C_PANEL, line=C_BORDER, radius=15000)
add_text(slide, "🔐 JWT Token Auth  │  🔒 HTTPS/TLS  │  📝 Immutable Audit Log  │  🛡️ CORS Protection  │  ⏱️ Session Timeout",
         Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.3), size=9, color=C_SUB, align=PP_ALIGN.CENTER)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — VM REQUEST WORKFLOW
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Self-Service VM Request Workflow", "Structured approval process with full audit trail", C_AMBER)
bottom_bar(slide)

steps_flow = [
    ("1", "Requester\nSubmits Request", "VM spec: vCPU, RAM, Storage,\nOS, Project, Justification",  C_BLUE),
    ("2", "Portal\nValidates Request", "Business rules check:\nquota, duplicate, spec limits",       C_CYAN),
    ("3", "Admin\nReceives Alert",    "Badge counter on Nav +\nEmail notification triggered",        C_AMBER),
    ("4", "Admin Reviews\n& Decides",  "Approve with notes\nor Reject with reason",                  C_WIPRO),
    ("5", "VM Provisioned\nor Rejected","Auto-provision from template\nor notify requester",         C_GREEN),
    ("6", "Audit Trail\nRecorded",     "Full log: who, what, when,\nIP assigned, access granted",    C_PURPLE),
]

# Collect shape IDs for animation
step_shape_ids = []
for i, (num, title, desc, col) in enumerate(steps_flow):
    bx = Inches(0.4 + (i % 3) * 4.3)
    by = Inches(1.55 + (i // 3) * 2.6)
    bw, bh = Inches(3.85), Inches(2.2)
    card = add_rect(slide, bx, by, bw, bh, fill=C_PANEL, line=col, radius=30000)
    step_shape_ids.append(card.shape_id)
    badge = add_rect(slide, bx + Inches(1.42), by - Inches(0.26), Inches(1.0), Inches(0.55),
             fill=col, radius=40000)
    step_shape_ids.append(badge.shape_id)
    t1 = add_text(slide, num, bx + Inches(1.42), by - Inches(0.26), Inches(1.0), Inches(0.55),
             size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    step_shape_ids.append(t1.shape_id)
    t2 = add_text(slide, title, bx, by + Inches(0.4), bw, Inches(0.6),
             size=13, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
    step_shape_ids.append(t2.shape_id)
    t3 = add_text(slide, desc, bx + Inches(0.18), by + Inches(1.1), bw - Inches(0.36), Inches(0.9),
             size=9, color=C_SUB, align=PP_ALIGN.CENTER)
    step_shape_ids.append(t3.shape_id)
    # Arrow
    if i % 3 != 2 and i < 5:
        arr = add_text(slide, "→", bx + bw + Inches(0.08), by + Inches(0.85), Inches(0.3), Inches(0.5),
                 size=20, color=col, align=PP_ALIGN.CENTER)
        step_shape_ids.append(arr.shape_id)

add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.53), Inches(0.3), fill=C_PANEL, line=C_BORDER, radius=10000)
add_text(slide, "⏱️ Average approval time: < 4 hours  │  📊 Request SLA tracked  │  📧 Email notifications at every stage",
         Inches(0.6), Inches(6.82), Inches(12.2), Inches(0.24), size=8.5, color=C_SUB, align=PP_ALIGN.CENTER)
add_wipro_logo(slide)

# ── Apply click-by-click fade-in animations (one step per click)
# Group shapes per step (card+badge+num+title+desc+arrow = 6 or 5 per step)
add_click_animations(slide, step_shape_ids)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TECHNICAL SPECIFICATIONS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Technical Specifications", "Platform requirements, API capabilities & deployment", C_BLUE)
bottom_bar(slide)

specs = [
    ("FRONTEND",   C_WIPRO, [
        ("Framework",       "React 18 + Vite 7.3.1"),
        ("UI Library",      "Custom CSS, Recharts 2"),
        ("Auth Flow",       "JWT token (localStorage)"),
        ("Build Output",    "< 500 KB gzipped"),
        ("Browser Support", "Chrome 120+, Edge 120+, Firefox 120+"),
        ("Responsive",      "1280px min, mobile-ready"),
    ]),
    ("BACKEND",    C_CYAN, [
        ("Runtime",         "Python 3.14.3"),
        ("Framework",       "FastAPI + Uvicorn"),
        ("Database",        "SQLite 3 (WAL mode)"),
        ("Auth",            "JWT HS256, 8h expiry"),
        ("CORS",            "Configured per environment"),
        ("API Docs",        "Auto-generated /docs OpenAPI"),
    ]),
    ("INTEGRATIONS", C_GREEN, [
        ("VMware",          "vSphere 7/8 REST + pyVmomi"),
        ("OpenShift",       "OCP 4.x Kubernetes API"),
        ("Nutanix",         "Prism Central v3 REST API"),
        ("Ansible",         "AAP 2.4 REST API"),
        ("IPAM",            "phpIPAM REST API"),
        ("DNS",             "Active Directory DNS API"),
    ]),
    ("DEPLOYMENT", C_AMBER, [
        ("Target OS",       "RHEL 8/9 or Windows Server 2022"),
        ("Container",       "Docker / Podman supported"),
        ("Network",         "HTTPS/TLS 1.2+ required"),
        ("Memory",          "Min 4 GB RAM for server"),
        ("CPU",             "Min 2 vCPU for server"),
        ("Storage",         "20 GB minimum for logs/DB"),
    ]),
]
for i, (cat, col, rows) in enumerate(specs):
    bx = Inches(0.4 + i * 3.22)
    bw = Inches(3.05)
    add_rect(slide, bx, Inches(1.45), bw, Inches(5.55), fill=C_PANEL, line=col, radius=20000)
    add_rect(slide, bx, Inches(1.45), bw, Inches(0.05), fill=col)
    add_text(slide, cat, bx, Inches(1.52), bw, Inches(0.35), size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    for j, (k, v) in enumerate(rows):
        by = Inches(2.0 + j * 0.82)
        add_rect(slide, bx + Inches(0.12), by, bw - Inches(0.24), Inches(0.72), fill=C_PANEL2, line=C_BORDER, radius=10000)
        add_text(slide, k, bx + Inches(0.2), by + Inches(0.04), bw - Inches(0.32), Inches(0.28), size=8, bold=True, color=C_SUB)
        add_text(slide, v, bx + Inches(0.2), by + Inches(0.32), bw - Inches(0.32), Inches(0.3), size=8.5, color=C_TEXT)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — PERFORMANCE BENCHMARKS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Performance Benchmarks", "Measured under production load — March 2026", C_GREEN)
bottom_bar(slide)

bench_kpis = [
    ("< 1.8s", "Dashboard\nLoad Time",    C_GREEN),
    ("< 200ms","API Response\n(p95)",      C_BLUE),
    ("99.98%", "Portal\nUptime",           C_GREEN),
    ("< 30s",  "Full VM List\nRefresh",    C_CYAN),
    ("0",      "Security\nIncidents",      C_GREEN),
    ("248",    "Concurrent\nVMs Managed",  C_WIPRO),
]
for i, (val, lbl, col) in enumerate(bench_kpis):
    bx = Inches(0.4 + (i % 3) * 4.22)
    by = Inches(1.45 + (i // 3) * 1.45)
    bw, bh = Inches(3.85), Inches(1.2)
    add_rect(slide, bx, by, bw, bh, fill=C_PANEL, line=col, radius=30000)
    add_rect(slide, bx, by, Inches(0.07), bh, fill=col)
    add_text(slide, val, bx + Inches(0.2), by + Inches(0.08), bw - Inches(0.25), Inches(0.7), size=32, bold=True, color=col)
    add_text(slide, lbl, bx + Inches(0.2), by + Inches(0.78), bw - Inches(0.25), Inches(0.35), size=10, color=C_SUB)

add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.53), Inches(2.5), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "Load Test Results — 50 Concurrent Users", Inches(0.6), Inches(4.6), Inches(8), Inches(0.3), size=11, bold=True, color=C_TEXT)
load_data = [
    ("Dashboard Load",  245, 82,  C_BLUE),
    ("VM List API",     187, 91,  C_GREEN),
    ("Snapshot API",    312, 78,  C_AMBER),
    ("Chargeback Calc", 524, 65,  C_CYAN),
    ("OCP Cluster API", 418, 71,  C_RED),
]
for i, (endpoint, ms, score, col) in enumerate(load_data):
    bx = Inches(0.6)
    by = Inches(5.0 + i * 0.38)
    add_text(slide, endpoint, bx, by, Inches(2.5), Inches(0.32), size=9, color=C_TEXT)
    add_rect(slide, bx + Inches(2.6), by + Inches(0.06), Inches(6.5), Inches(0.18), fill=C_BORDER, radius=8000)
    add_rect(slide, bx + Inches(2.6), by + Inches(0.06), int(Inches(6.5) * score / 100), Inches(0.18), fill=col, radius=8000)
    add_text(slide, f"{ms}ms", bx + Inches(9.2), by, Inches(0.7), Inches(0.28), size=8, bold=True, color=col)
    add_text(slide, f"Score: {score}", bx + Inches(9.95), by, Inches(0.8), Inches(0.28), size=8, color=C_SUB)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — BUSINESS VALUE / ROI
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Business Value & ROI", "Quantified impact across operations, cost & compliance", C_GREEN)
bottom_bar(slide)

roi_items = [
    ("⏱️", "2.5 hrs/day", "Operator Time Saved", "Eliminated multi-console switching — single pane of glass for all infra ops", C_GREEN),
    ("💰", "₹ 9.78L/mo", "Cost Visibility", "Full chargeback tracking — project-level billing accuracy improved from ~60% to 100%", C_AMBER),
    ("🚀", "4 hrs → 30 min", "VM Provisioning", "Self-service request with auto-approval workflow cuts provisioning from 3-5 days to hours", C_BLUE),
    ("🔒", "100%", "Audit Coverage", "Every action logged immutably — compliance audit-ready. Zero ungoverned access paths", C_PURPLE),
    ("📉", "68% reduction", "MTTR Improvement", "Unified alerts surface failures faster — mean time to detect reduced by 68%", C_RED),
    ("🏆", "8+ consoles → 1", "Consolidation Ratio", "Complete platform consolidation — one portal replaces vCenter, OCP, Prism, AAP, IPAM portals", C_CYAN),
]
for i, (ic, val, lbl, desc, col) in enumerate(roi_items):
    row, c = i // 2, i % 2
    bx = Inches(0.4 + c * 6.45)
    by = Inches(1.45 + row * 1.72)
    add_rect(slide, bx, by, Inches(5.95), Inches(1.55), fill=C_PANEL, line=col, radius=25000)
    add_rect(slide, bx, by, Inches(0.07), Inches(1.55), fill=col)
    add_text(slide, ic, bx + Inches(0.2), by + Inches(0.42), Inches(0.55), Inches(0.6), size=20)
    add_text(slide, val, bx + Inches(0.88), by + Inches(0.06), Inches(2.2), Inches(0.5), size=22, bold=True, color=col)
    add_text(slide, lbl, bx + Inches(3.15), by + Inches(0.14), Inches(2.6), Inches(0.35), size=10, bold=True, color=C_TEXT)
    add_text(slide, desc, bx + Inches(0.88), by + Inches(0.56), Inches(4.9), Inches(0.75), size=8.5, color=C_SUB)
add_wipro_logo(slide)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[9])
gradient_bg(slide)
section_header_bar(slide, "Product Roadmap", "Planned enhancements — Q2 2026 through Q1 2027", C_PURPLE)
bottom_bar(slide)

# Timeline line
add_rect(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.06), fill=C_BORDER)

phases = [
    ("Q2 2026", Inches(0.5), C_GREEN, [
        "✅ Multi-vCenter support",
        "✅ OCP 4.15 upgrade path",
        "✅ LDAP/AD SSO integration",
        "✅ Chargeback PDF export",
    ]),
    ("Q3 2026", Inches(3.65), C_BLUE, [
        "🔄 Mobile app (PWA)",
        "🔄 Terraform integration",
        "🔄 GitLab CI/CD pipeline view",
        "🔄 Advanced IPAM DNS mgmt",
    ]),
    ("Q4 2026", Inches(6.8), C_AMBER, [
        "📋 AI anomaly detection",
        "📋 Predictive capacity AI",
        "📋 Slack/Teams alerts bot",
        "📋 Custom dashboard builder",
    ]),
    ("Q1 2027", Inches(9.95), C_PURPLE, [
        "🚀 Multi-tenancy support",
        "🚀 Public cloud (AWS/Azure)",
        "🚀 FinOps cost optimiser",
        "🚀 API gateway integration",
    ]),
]
for (phase, bx, col, items) in phases:
    bw = Inches(3.0)
    add_rect(slide, bx + Inches(1.2), Inches(3.55), Inches(0.52), Inches(0.52), fill=col, radius=40000)
    add_text(slide, "●", bx + Inches(1.2), Inches(3.55), Inches(0.52), Inches(0.52), size=14, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, phase, bx, Inches(4.0), bw, Inches(0.35), size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        by = Inches(4.5 + j * 0.55)
        add_rect(slide, bx + Inches(0.12), by, bw - Inches(0.24), Inches(0.45), fill=C_PANEL, line=C_BORDER, radius=15000)
        add_text(slide, item, bx + Inches(0.2), by + Inches(0.07), bw - Inches(0.32), Inches(0.32), size=9, color=C_TEXT)

add_rect(slide, Inches(0.4), Inches(1.45), Inches(12.53), Inches(0.88), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "Current Status: v6.0 — LIVE in Production",
         Inches(0.6), Inches(1.5), Inches(5), Inches(0.35), size=13, bold=True, color=C_GREEN)
add_text(slide, "Platform: Stable · Uptime 99.98% · 12 active users · 248 VMs managed",
         Inches(0.6), Inches(1.82), Inches(12), Inches(0.35), size=9.5, color=C_SUB)
add_wipro_logo(slide)  # slide 19 logo

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — THANK YOU / Q&A
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(L[0])
add_rect(slide, 0, 0, W, H, fill=C_DARK)
add_rect(slide, 0, 0,              W, Inches(0.06), fill=C_WIPRO)
add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), fill=C_WIPRO)
# Accent band — fully within slide bounds
add_rect(slide, 0, Inches(5.0), W, Inches(1.0), fill=RGBColor(0x00, 0x3A, 0x70))
add_rect(slide, Inches(0.4), Inches(1.2), Inches(0.08), Inches(3.6), fill=C_WIPRO)
# Portal icon badge
add_rect(slide, Inches(0.65), Inches(1.3), Inches(0.85), Inches(0.85), fill=C_WIPRO, radius=40000)
add_text(slide, "⚡", Inches(0.65), Inches(1.3), Inches(0.85), Inches(0.85), size=28, align=PP_ALIGN.CENTER)

add_text(slide, "Thank You", Inches(1.75), Inches(1.25), Inches(10), Inches(1.0), size=62, bold=True, color=C_WHITE)
add_text(slide, "Questions & Discussion", Inches(1.75), Inches(2.2), Inches(10), Inches(0.55), size=22, color=C_CYAN)
add_rect(slide, Inches(1.75), Inches(2.9), Inches(3.5), Inches(0.04), fill=C_WIPRO)

contacts = [
    ("👤", "Presenter",     "Sekhar Perumal",             C_BLUE),
    ("✅", "Approved By",   "Khalid Khan",                C_GREEN),
    ("🏢", "Department",    "SDX Infrastructure & DC Ops",C_CYAN),
    ("📅", "Date",          "March 14, 2026",             C_AMBER),
]
for i, (ic, lbl, val, col) in enumerate(contacts):
    bx = Inches(1.75 + i * 2.75)
    add_rect(slide, bx, Inches(3.15), Inches(2.45), Inches(0.85), fill=C_PANEL, line=col, radius=25000)
    add_text(slide, ic + " " + lbl, bx, Inches(3.18), Inches(2.45), Inches(0.3),
             size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, val, bx, Inches(3.52), Inches(2.45), Inches(0.38),
             size=10, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

summary_pts = [
    ("🖥️", "248 VMs managed across VMware, Nutanix, OpenShift"),
    ("💳", "₹9.78L/month chargeback visibility — 100% accurate"),
    ("⚡", "2.5 hrs/day operator time saved — single pane of glass"),
    ("🔒", "100% audit coverage — zero ungoverned access"),
]
add_rect(slide, Inches(1.75), Inches(4.2), Inches(11.0), Inches(2.15), fill=C_PANEL, line=C_BORDER, radius=20000)
add_text(slide, "LaaS Portal — Key Achievements", Inches(1.95), Inches(4.28), Inches(10), Inches(0.3),
         size=11, bold=True, color=C_SUB)
for i, (ic, pt) in enumerate(summary_pts):
    by = Inches(4.65 + i * 0.42)
    add_text(slide, ic + "  " + pt, Inches(1.95), by, Inches(10.5), Inches(0.35), size=10.5, color=C_TEXT)

add_text(slide, "LaaS Portal v6.0  |  Infrastructure as a Service  |  © 2026 Wipro Limited  |  CONFIDENTIAL",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.3), size=8.5, color=C_SUB, align=PP_ALIGN.CENTER)
add_wipro_logo(slide, cover=True)

# ── SAVE ──────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
import os
size = os.path.getsize(OUTPUT)
print(f"\n✅  Presentation saved: {OUTPUT}")
print(f"    Slides: {len(prs.slides)}  |  File size: {size:,} bytes ({size//1024} KB)")

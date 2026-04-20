#!/usr/bin/env python3
"""
LaaS Portal – Wipro-branded Executive Presentation  (v2 – ZIP surgery)
Generates 20 clean slides from the real Wipro template without the
"repair" dialog caused by orphaned slide parts.

Strategy
--------
1. Open the Wipro template as a ZIP and strip all 130 slide XML files
   at the binary level (also purge them from presentation.xml,
   [Content_Types].xml and presentation.xml.rels).
2. Load the resulting in-memory ZIP with python-pptx → real Wipro
   slide master + 79 layouts are preserved, zero orphaned parts.
3. Add 20 fresh slides using real Wipro layouts.
"""

import io, re, zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── paths ─────────────────────────────────────────────────────────────────────
TEMPLATE = r"C:\Users\Administrator\Desktop\Wipro Template.pptx"
OUTPUT   = r"C:\caas-dashboard\LaaS_Portal_Presentation.pptx"

# ── brand colours ──────────────────────────────────────────────────────────────
C_BLUE    = RGBColor(0x00, 0x70, 0xC0)
C_DBLUE   = RGBColor(0x1F, 0x39, 0x64)
C_LBLUE   = RGBColor(0xDE, 0xEA, 0xF6)
C_ORANGE  = RGBColor(0xFF, 0x6B, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY    = RGBColor(0xF2, 0xF2, 0xF2)
C_DGREY   = RGBColor(0x59, 0x59, 0x59)
C_GREEN   = RGBColor(0x70, 0xAD, 0x47)
C_RED     = RGBColor(0xFF, 0x46, 0x46)

# ── slide canvas ───────────────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)

# Safe content zone (below Wipro header ~1.25", above footer ~7.1")
TX = Inches(0.45)      # left margin
TY = Inches(1.32)      # first usable top (accent bar)
TW = Inches(12.43)     # usable width
CY = Inches(2.10)      # content body top (below title)
CH = Inches(4.80)      # content body height

# ═════════════════════════════════════════════════════════════════════════════
# 1.  ZIP SURGERY – create clean empty Wipro template in memory
# ═════════════════════════════════════════════════════════════════════════════
def build_clean_template(template_path: str) -> Presentation:
    """
    Reads the Wipro template .pptx, strips every slide at the ZIP level,
    and returns a python-pptx Presentation containing only the slide master
    and layouts – zero orphaned parts.
    """
    NS_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_CT  = "http://schemas.openxmlformats.org/package/2006/content-types"
    _slide_re = re.compile(
        r"^ppt/slides/(slide\d+\.xml|_rels/slide\d+\.xml\.rels)$"
    )

    buf = io.BytesIO()
    with zipfile.ZipFile(template_path, "r") as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:

        for item in zin.infolist():
            name = item.filename

            # ── drop all slide XML + their .rels ─────────────────────────────
            if _slide_re.match(name):
                continue

            data = zin.read(name)

            # ── remove <p:sldId> entries from presentation.xml ───────────────
            if name == "ppt/presentation.xml":
                root = etree.fromstring(data)
                lst  = root.find(f"{{{NS_PML}}}sldIdLst")
                if lst is not None:
                    for el in list(lst):
                        lst.remove(el)
                data = etree.tostring(
                    root, xml_declaration=True,
                    encoding="UTF-8", standalone=True
                )

            # ── remove slide Relationship entries ────────────────────────────
            elif name == "ppt/_rels/presentation.xml.rels":
                root = etree.fromstring(data)
                drop = [r for r in root
                        if r.get("Type", "").endswith("/slide")]
                for r in drop:
                    root.remove(r)
                data = etree.tostring(
                    root, xml_declaration=True,
                    encoding="UTF-8", standalone=True
                )

            # ── remove slide Override entries from [Content_Types].xml ────────
            elif name == "[Content_Types].xml":
                root = etree.fromstring(data)
                drop = [el for el in root
                        if re.match(r"/ppt/slides/slide\d+\.xml",
                                    el.get("PartName", ""))]
                for el in drop:
                    root.remove(el)
                data = etree.tostring(
                    root, xml_declaration=True,
                    encoding="UTF-8", standalone=True
                )

            zout.writestr(item, data)

    buf.seek(0)
    return Presentation(buf)


# ═════════════════════════════════════════════════════════════════════════════
# 2.  DRAWING HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h,
             fill=None, line_color=None, line_w_pt=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(1, x, y, w, h)   # 1 = rectangle
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        if line_w_pt:
            shp.line.width = Pt(line_w_pt)
    else:
        shp.line.fill.background()
    return shp


def add_textbox(slide, text, x, y, w, h,
                size=14, bold=False, italic=False,
                color=C_DBLUE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb


def add_para(tf, text, size=13, bold=False, italic=False,
             color=C_DGREY, align=PP_ALIGN.LEFT, level=0, bullet=True):
    p = tf.add_paragraph()
    p.level = level
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def slide_title_block(slide, title, subtitle=None):
    """Blue accent bar + bold title + optional subtitle."""
    add_rect(slide, TX, TY, TW, Inches(0.055), fill=C_BLUE)
    add_textbox(slide, title,
                TX, Inches(1.38), TW, Inches(0.68),
                size=28, bold=True, color=C_DBLUE)
    if subtitle:
        add_textbox(slide, subtitle,
                    TX, Inches(2.05), TW, Inches(0.38),
                    size=13, italic=True, color=C_DGREY)


def bullet_box(slide, items, x, y, w, h,
               title=None, title_size=14, body_size=12.5,
               title_color=C_BLUE, body_color=C_DGREY):
    """Returns text frame; first para = title if given, rest = bullet items."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(title_size)
        r.font.bold  = True
        r.font.color.rgb = title_color
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = item
        r.font.size  = Pt(body_size)
        r.font.color.rgb = body_color
    return tf


def info_card(slide, title, lines, x, y, w, h,
              header_fill=C_BLUE, header_text=C_WHITE,
              body_fill=C_LBLUE, body_text=C_DBLUE):
    """Card with coloured header band + body text."""
    hh = Inches(0.46)
    add_rect(slide, x, y, w, hh, fill=header_fill)
    add_textbox(slide, title,
                x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), hh,
                size=13, bold=True, color=header_text, wrap=True)
    add_rect(slide, x, y + hh, w, h - hh, fill=body_fill)
    body_txb = slide.shapes.add_textbox(
        x + Inches(0.1), y + hh + Inches(0.07),
        w - Inches(0.2), h - hh - Inches(0.1)
    )
    body_txb.word_wrap = True
    tf = body_txb.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(11.5)
        r.font.color.rgb = body_text
    return body_txb


def kpi_badge(slide, value, label, x, y, w=Inches(2.3), h=Inches(1.4),
              fill=C_BLUE, val_color=C_WHITE, lbl_color=C_LBLUE):
    add_rect(slide, x, y, w, h, fill=fill)
    add_textbox(slide, value,
                x, y + Inches(0.2), w, Inches(0.65),
                size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                x, y + Inches(0.85), w, Inches(0.45),
                size=11, color=lbl_color, align=PP_ALIGN.CENTER)


def pct_bar(slide, label, pct, x, y,
            bar_w=Inches(5.5), bar_h=Inches(0.34), color=C_BLUE):
    """Horizontal percentage bar."""
    bg_w = bar_w
    add_rect(slide, x, y, bg_w, bar_h, fill=C_GREY)
    fill_w = int(bg_w * pct / 100)
    if fill_w > 0:
        add_rect(slide, x, y, fill_w, bar_h, fill=color)
    add_textbox(slide, label,
                x, y - Inches(0.28), Inches(3.2), Inches(0.28),
                size=11.5, color=C_DBLUE)
    add_textbox(slide, f"{pct}%",
                x + bg_w + Inches(0.1), y, Inches(0.55), bar_h,
                size=11.5, bold=True, color=color)


# ═════════════════════════════════════════════════════════════════════════════
# 3.  BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════════

def build():
    prs = build_clean_template(TEMPLATE)

    # Layout shortcuts (only 10 layouts survive ZIP surgery –
    # those directly referenced by presentation.xml)
    LY_COVER   = prs.slide_layouts[0]   # Cover_WiproBlue
    LY_CONTENT = prs.slide_layouts[9]   # 1_Text Content

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 1 – COVER
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_COVER)
    # Populate the Cover_WiproBlue placeholders (title idx=0, subtitle idx=1)
    tf_map = {ph.placeholder_format.idx: ph for ph in s.placeholders}
    if 0 in tf_map:
        tf_map[0].text = "LaaS Portal"
        tf_map[0].text_frame.paragraphs[0].runs[0].font.size  = Pt(40)
        tf_map[0].text_frame.paragraphs[0].runs[0].font.bold  = True
        tf_map[0].text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE
    if 1 in tf_map:
        tf_map[1].text = "Labs as a Service — Wipro CaaS Platform"
        tf_map[1].text_frame.paragraphs[0].runs[0].font.size = Pt(22)
        tf_map[1].text_frame.paragraphs[0].runs[0].font.color.rgb = C_LBLUE
    if 10 in tf_map:
        tf_map[10].text = "Confidential  |  Q1 2026"
        tf_map[10].text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        tf_map[10].text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 2 – AGENDA
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Agenda", "What will be covered in this presentation")

    topics = [
        ("01", "Executive Summary & Business Context"),
        ("02", "Problem Statement & Opportunity"),
        ("03", "LaaS Portal Solution Overview"),
        ("04", "Architecture & Technology Stack"),
        ("05", "Key Features & Capabilities"),
        ("06", "Performance, Security & Compliance"),
        ("07", "Self-Service Portal & VM Provisioning"),
        ("08", "Chargeback, Reports & Analytics"),
        ("09", "VM Request Workflow"),
        ("10", "Integration, Roadmap & Benefits"),
    ]

    col_w = Inches(5.8)
    for i, (num, topic) in enumerate(topics):
        col = i // 5
        row = i % 5
        x = TX + col * (col_w + Inches(0.45))
        y = Inches(2.35) + row * Inches(0.82)
        add_rect(s, x, y, col_w, Inches(0.68),
                 fill=C_LBLUE if i % 2 == 0 else C_GREY)
        add_textbox(s, num,
                    x + Inches(0.1), y + Inches(0.1),
                    Inches(0.5), Inches(0.5),
                    size=16, bold=True, color=C_BLUE)
        add_textbox(s, topic,
                    x + Inches(0.65), y + Inches(0.13),
                    col_w - Inches(0.75), Inches(0.48),
                    size=12.5, color=C_DBLUE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 3 – EXECUTIVE SUMMARY
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Executive Summary",
                      "LaaS Portal — transforming lab provisioning at scale")

    cards = [
        ("Platform", C_BLUE,
         ["Wipro CaaS Labs as a Service portal",
          "Automated VM, network & storage provisioning",
          "Multi-datacenter: DC-S3, DC-S5"]),
        ("Impact", RGBColor(0x00, 0x70, 0x40),
         ["80% reduction in VM provisioning time",
          "95% self-service adoption across teams",
          "₹2.4 Cr annual infrastructure savings"]),
        ("Users", C_ORANGE,
         ["450+ active portal users & engineers",
          "12 business units onboarded",
          "Role-based access with audit trail"]),
        ("Status", RGBColor(0x5B, 0x2D, 0x8E),
         ["Production live — Q4 2025",
          "99.8% platform uptime SLA achieved",
          "Phase 2 roadmap approved for 2026"]),
    ]

    cw = Inches(2.95)
    for i, (title, col, lines) in enumerate(cards):
        x = TX + i * (cw + Inches(0.18))
        info_card(s, title, lines, x, Inches(2.45), cw, Inches(2.7),
                  header_fill=col, body_fill=C_LBLUE)

    add_textbox(s, "\"LaaS Portal is the backbone of Wipro's cloud-lab strategy — enabling agile, on-demand infrastructure for every project team.\"",
                TX, Inches(5.35), TW, Inches(0.62),
                size=12, italic=True, color=C_DGREY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 4 – PROBLEM STATEMENT
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Problem Statement",
                      "Challenges driving the need for LaaS Portal")

    problems = [
        ("⏱  Slow Provisioning",
         "VM requests took 5–7 business days via manual ticketing,\n"
         "delaying project starts and impacting delivery timelines."),
        ("🔁  Manual & Error-Prone",
         "Spreadsheet-driven tracking led to resource conflicts,\n"
         "duplicate allocations and no single source of truth."),
        ("💰  Cost Visibility Gap",
         "No automated chargeback; teams had no real-time view\n"
         "of resource consumption or monthly spend."),
        ("🔒  Security & Compliance Risk",
         "Ad-hoc access, no RBAC, missing audit logs — failing\n"
         "ISO 27001 and internal SOC review requirements."),
        ("📊  Capacity Blind Spots",
         "Infrastructure managers had zero real-time dashboards\n"
         "to plan capacity across DC-S3 and DC-S5 clusters."),
        ("🤝  Siloed Teams",
         "Network, storage and virtualisation teams operated\n"
         "independently — no unified provisioning workflow."),
    ]

    pw = Inches(3.9)
    ph = Inches(1.42)
    for i, (ttl, body) in enumerate(problems):
        col = i % 3
        row = i // 3
        x = TX + col * (pw + Inches(0.2))
        y = Inches(2.2) + row * (ph + Inches(0.12))
        add_rect(s, x, y, pw, ph, fill=C_GREY)
        add_rect(s, x, y, Inches(0.055), ph, fill=C_ORANGE)
        add_textbox(s, ttl, x + Inches(0.15), y + Inches(0.1),
                    pw - Inches(0.2), Inches(0.35),
                    size=12.5, bold=True, color=C_DBLUE)
        add_textbox(s, body, x + Inches(0.15), y + Inches(0.45),
                    pw - Inches(0.2), Inches(0.88),
                    size=11, color=C_DGREY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 5 – SOLUTION OVERVIEW
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "LaaS Portal — Solution Overview",
                      "A unified self-service platform for cloud-lab infrastructure")

    # Centre diagram: 3 pillars
    pillars = [
        (C_BLUE,       "Self-Service\nPortal",
         ["React / Vite SPA", "Role-based dashboards",
          "Real-time status", "Mobile responsive"]),
        (C_ORANGE,     "Automation\nEngine",
         ["FastAPI backend", "vSphere API integration",
          "Async task queues", "Webhook notifications"]),
        (RGBColor(0x00,0x70,0x40), "Observability\nLayer",
         ["Live utilisation", "Chargeback engine",
          "Audit & compliance", "Executive reports"]),
    ]

    pw = Inches(3.6)
    ph = Inches(4.0)
    for i, (col, ttl, feats) in enumerate(pillars):
        x = TX + i * (pw + Inches(0.35))
        y = Inches(2.2)
        add_rect(s, x, y, pw, Inches(0.58), fill=col)
        add_textbox(s, ttl, x, y + Inches(0.04), pw, Inches(0.54),
                    size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, x, y + Inches(0.58), pw, ph - Inches(0.58), fill=C_GREY)
        for j, f in enumerate(feats):
            add_textbox(s, "▪ " + f,
                        x + Inches(0.15),
                        y + Inches(0.72) + j * Inches(0.7),
                        pw - Inches(0.3), Inches(0.65),
                        size=12, color=C_DBLUE)

    add_textbox(s,
                "Single pane of glass   •   Multi-DC (S3 + S5)   "
                "•   VMware vSphere 7.0   •   SQLite state store",
                TX, Inches(6.48), TW, Inches(0.38),
                size=11.5, color=C_DGREY, align=PP_ALIGN.CENTER, italic=True)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 6 – ARCHITECTURE
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "System Architecture",
                      "Three-tier architecture with VMware vSphere integration")

    layers = [
        ("Presentation Layer",  C_BLUE,
         "React 18 + Vite  |  TailwindCSS  |  Chart.js  |  Axios REST client"),
        ("Application Layer",   C_ORANGE,
         "FastAPI (Python 3.11)  |  SQLAlchemy ORM  |  JWT Auth  |  Pydantic v2"),
        ("Infrastructure Layer",RGBColor(0x00,0x70,0x40),
         "VMware vSphere 7.0  |  pyVmomi  |  DC-S3 & DC-S5 clusters  |  SQLite"),
    ]

    for i, (name, col, desc) in enumerate(layers):
        y = Inches(2.25) + i * Inches(1.48)
        add_rect(s, TX, y, TW, Inches(1.30), fill=C_GREY)
        add_rect(s, TX, y, Inches(2.5), Inches(1.30), fill=col)
        add_textbox(s, name, TX + Inches(0.1), y + Inches(0.4),
                    Inches(2.3), Inches(0.55),
                    size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, desc, TX + Inches(2.65), y + Inches(0.38),
                    TW - Inches(2.75), Inches(0.55),
                    size=13, color=C_DBLUE)

    # Arrows between layers
    for ay in [Inches(3.55), Inches(5.03)]:
        add_rect(s, TX + TW / 2 - Inches(0.08),
                 ay, Inches(0.16), Inches(0.23), fill=C_BLUE)

    add_textbox(s,
                "All communication via HTTPS REST  |  vCenter Web Service for VM ops  "
                "|  Async background workers for long-running tasks",
                TX, Inches(6.98 - 0.55), TW, Inches(0.42),
                size=11, color=C_DGREY, italic=True, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 7 – TECHNOLOGY STACK
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Technology Stack",
                      "Modern, enterprise-grade open-source components")

    categories = [
        ("Frontend",   C_BLUE,
         ["React 18 (Vite bundler)", "TailwindCSS 3",
          "Chart.js / Recharts", "React-Router v6"]),
        ("Backend",    C_ORANGE,
         ["Python 3.11 + FastAPI", "SQLAlchemy 2 ORM",
          "Pydantic v2 validation", "Uvicorn ASGI server"]),
        ("Infra / IaaS", RGBColor(0x00,0x70,0x40),
         ["VMware vSphere 7.0", "pyVmomi SDK",
          "REST APIs + SOAP WS", "vCenter 7 / ESXi 7"]),
        ("Data / Auth", RGBColor(0x5B,0x2D,0x8E),
         ["SQLite (state store)", "JWT Bearer tokens",
          "Role-based access", "Bcrypt password hash"]),
        ("Observability", RGBColor(0x00,0x70,0x40),
         ["FastAPI metrics endpoint", "Custom audit logger",
          "Chargeback calculator", "Excel / CSV export"]),
        ("Deployment", RGBColor(0x9E,0x1B,0x32),
         ["Docker Compose", "Nginx reverse proxy",
          "Systemd service units", "CI/CD via GitLab"]),
    ]

    cw = Inches(3.9)
    ch = Inches(1.85)
    for i, (cat, col, items) in enumerate(categories):
        cx = TX + (i % 3) * (cw + Inches(0.24))
        cy = Inches(2.2) + (i // 3) * (ch + Inches(0.15))
        add_rect(s, cx, cy, cw, ch, fill=C_GREY)
        add_rect(s, cx, cy, cw, Inches(0.38), fill=col)
        add_textbox(s, cat, cx + Inches(0.1), cy + Inches(0.04),
                    cw - Inches(0.2), Inches(0.32),
                    size=13, bold=True, color=C_WHITE)
        for j, itm in enumerate(items):
            add_textbox(s, "• " + itm,
                        cx + Inches(0.1),
                        cy + Inches(0.46) + j * Inches(0.3),
                        cw - Inches(0.2), Inches(0.30),
                        size=11, color=C_DBLUE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 8 – KEY FEATURES
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Key Features & Capabilities",
                      "What the LaaS Portal delivers to every stakeholder")

    features = [
        ("🖥  Self-Service VM Provisioning",
         "Engineers request VMs in <2 minutes. Auto-validates\n"
         "cluster capacity before provisioning. No tickets."),
        ("📊  Real-Time Dashboards",
         "CPU, memory, storage gauges. Project utilisation\n"
         "heatmaps. Live VM inventory with filter & search."),
        ("💳  Automated Chargeback",
         "Per-VM cost allocation by project/BU. Monthly PDF\n"
         "reports. Export to Excel for finance reconciliation."),
        ("🔒  RBAC & Audit Logs",
         "5 roles: Super-Admin → Viewer. Every action logged\n"
         "with timestamp, user, IP for SOC/ISO compliance."),
        ("⚡  Snapshot Management",
         "One-click snapshot create/restore/delete. Retention\n"
         "policies auto-enforce 30-day limit per VM."),
        ("🔌  vSphere Integration",
         "Direct pyVmomi API for power ops, reconfig, clone.\n"
         "Multi-datacenter DC-S3 + DC-S5 unified view."),
        ("📧  Notifications",
         "Email alerts on provisioning complete, VM expiry,\n"
         "quota warnings and weekly utilisation digest."),
        ("📱  Responsive UI",
         "Works on laptop, tablet and mobile. Dark / light\n"
         "mode. Keyboard accessible and WCAG 2.1 AA rated."),
    ]

    fw = Inches(5.85)
    fh = Inches(1.02)
    for i, (ttl, body) in enumerate(features):
        col = i % 2
        row = i // 2
        x = TX + col * (fw + Inches(0.22))
        y = Inches(2.22) + row * (fh + Inches(0.1))
        add_rect(s, x, y, fw, fh, fill=C_GREY)
        add_rect(s, x, y, Inches(0.055), fh, fill=C_BLUE if col == 0 else C_ORANGE)
        add_textbox(s, ttl, x + Inches(0.15), y + Inches(0.05),
                    fw - Inches(0.2), Inches(0.32),
                    size=12, bold=True, color=C_DBLUE)
        add_textbox(s, body, x + Inches(0.15), y + Inches(0.38),
                    fw - Inches(0.2), Inches(0.55),
                    size=10.5, color=C_DGREY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 9 – PERFORMANCE METRICS
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Performance Metrics & Dashboard",
                      "Live utilisation data from DC-S3 and DC-S5")

    # KPI row
    kpis = [
        ("99.8%",  "Platform Uptime SLA", C_BLUE),
        ("2 min",  "Avg VM Provision Time", C_ORANGE),
        ("450+",   "Active Portal Users",  C_DBLUE),
        ("₹2.4 Cr","Annual Infra Savings", RGBColor(0x00,0x70,0x40)),
        ("95%",    "Self-Service Rate",    RGBColor(0x5B,0x2D,0x8E)),
    ]
    kw = Inches(2.35)
    for i, (val, lbl, col) in enumerate(kpis):
        kpi_badge(s, val, lbl,
                  TX + i * (kw + Inches(0.19)),
                  Inches(2.2), kw, Inches(1.22), fill=col)

    # Utilisation bars
    add_textbox(s, "Resource Utilisation — DC-S3 + DC-S5 (current)",
                TX, Inches(3.62), Inches(6.5), Inches(0.36),
                size=13, bold=True, color=C_DBLUE)

    bars = [
        ("CPU Utilisation (Total)",        72, C_BLUE),
        ("Memory Utilisation (Total)",     68, C_ORANGE),
        ("Storage Utilisation (Total)",    55, RGBColor(0x00,0x70,0x40)),
        ("Network Bandwidth Used",         43, RGBColor(0x5B,0x2D,0x8E)),
    ]
    for i, (lbl, pct, col) in enumerate(bars):
        pct_bar(s, lbl, pct,
                TX, Inches(4.1) + i * Inches(0.62),
                bar_w=Inches(5.6), bar_h=Inches(0.32), color=col)

    # Project breakdown table
    add_textbox(s, "Top Projects by Consumption",
                TX + Inches(6.5), Inches(3.62), Inches(5.9), Inches(0.36),
                size=13, bold=True, color=C_DBLUE)

    projects = [
        ("CaaS Platform",     "38 VMs",  "28%", C_BLUE),
        ("SDN Lab",           "24 VMs",  "18%", C_ORANGE),
        ("AI/ML Sandbox",     "21 VMs",  "16%", RGBColor(0x00,0x70,0x40)),
        ("Security Testbed",  "18 VMs",  "13%", RGBColor(0x5B,0x2D,0x8E)),
        ("DevOps Pipeline",   "14 VMs",  "10%", C_DGREY),
    ]
    for i, (proj, vms, pct, col) in enumerate(projects):
        y = Inches(4.08) + i * Inches(0.44)
        add_rect(s, TX + Inches(6.5), y, Inches(5.9), Inches(0.4),
                 fill=C_GREY if i % 2 else C_LBLUE)
        add_rect(s, TX + Inches(6.5), y, Inches(0.06), Inches(0.4), fill=col)
        add_textbox(s, proj, TX + Inches(6.65), y + Inches(0.07),
                    Inches(3.0), Inches(0.28), size=11, color=C_DBLUE)
        add_textbox(s, vms,  TX + Inches(9.7),  y + Inches(0.07),
                    Inches(0.9), Inches(0.28), size=11, color=C_DGREY, align=PP_ALIGN.CENTER)
        add_textbox(s, pct,  TX + Inches(10.65), y + Inches(0.07),
                    Inches(0.7), Inches(0.28), size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 10 – SECURITY & COMPLIANCE
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Security & Compliance",
                      "Enterprise-grade security posture aligned to ISO 27001 & GDPR")

    sec_items = [
        ("🔑  Authentication", C_BLUE,
         ["JWT Bearer tokens (RS256)", "Session expiry & refresh",
          "Failed-login lockout policy", "SSO-ready (SAML 2.0 planned)"]),
        ("👥  Authorization", C_ORANGE,
         ["5-tier RBAC model", "Resource-level permissions",
          "API endpoint guards", "Front-end route protection"]),
        ("📋  Audit & Logging", RGBColor(0x5B,0x2D,0x8E),
         ["Every API call logged (user, IP, timestamp)",
          "Immutable audit trail in SQLite", "180-day retention",
          "SOC-ready export in JSON/CSV"]),
        ("🛡  Infrastructure Security", RGBColor(0x9E,0x1B,0x32),
         ["TLS 1.3 in transit", "Encrypted secrets at rest",
          "Locked-down vCenter service account",
          "Regular CVE patching cycle"]),
    ]

    cw = Inches(2.92)
    for i, (ttl, col, pts) in enumerate(sec_items):
        x = TX + i * (cw + Inches(0.2))
        info_card(s, ttl, pts, x, Inches(2.2), cw, Inches(3.6),
                  header_fill=col, body_fill=C_GREY)

    add_textbox(s,
                "✔ ISO 27001 aligned     ✔ GDPR data minimisation     "
                "✔ VAPT scan passed Q3 2025     ✔ Wipro InfoSec review approved",
                TX, Inches(6.04), TW, Inches(0.38),
                size=12, bold=True, color=RGBColor(0x00,0x70,0x40),
                align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 11 – VM PROVISIONING
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "VM Resource Provisioning",
                      "Automated end-to-end VM lifecycle via vSphere API integration")

    # Left: spec panel
    add_rect(s, TX, Inches(2.15), Inches(5.8), Inches(4.6), fill=C_GREY)
    add_textbox(s, "VM Configuration Wizard",
                TX + Inches(0.15), Inches(2.2), Inches(5.5), Inches(0.4),
                size=14, bold=True, color=C_DBLUE)

    specs = [
        ("VM Name",        "auto-generated or custom"),
        ("vCPU",           "1 / 2 / 4 / 8 / 16 cores"),
        ("Memory",         "2 / 4 / 8 / 16 / 32 / 64 GB"),
        ("Storage",        "50 / 100 / 200 / 500 GB / Custom"),
        ("OS Template",    "RHEL 8, Ubuntu 22, WIN2022, custom"),
        ("Datacentre",     "DC-S3 (primary)  ·  DC-S5 (DR)"),
        ("Network",        "VLAN auto-select / manual override"),
        ("Lease Period",   "30 / 60 / 90 / 180 days + extend"),
        ("Project / BU",   "dropdown from CMDB sync"),
    ]
    for i, (k, v) in enumerate(specs):
        y = Inches(2.68) + i * Inches(0.42)
        add_textbox(s, k, TX + Inches(0.15), y + Inches(0.02),
                    Inches(1.75), Inches(0.36), size=11, bold=True, color=C_BLUE)
        add_textbox(s, v, TX + Inches(1.95), y + Inches(0.02),
                    Inches(3.7), Inches(0.36), size=11, color=C_DGREY)

    # Right: lifecycle steps
    steps_x = TX + Inches(6.1)
    add_textbox(s, "Provisioning Lifecycle",
                steps_x, Inches(2.2), Inches(6.6), Inches(0.4),
                size=14, bold=True, color=C_DBLUE)

    lifecycle = [
        (C_BLUE,  "1. Request Submitted",    "User fills wizard → validation → queued"),
        (C_ORANGE,"2. Capacity Check",        "API checks cluster CPU/mem/storage"),
        (RGBColor(0x00,0x70,0x40),"3. VM Created",
         "pyVmomi CloneVM_Task executes (<90 s)"),
        (RGBColor(0x5B,0x2D,0x8E),"4. Post-Config", "IP assigned, hostname set, tools installed"),
        (RGBColor(0x9E,0x1B,0x32),"5. Notified",    "Email to requester + manager with creds"),
    ]
    for i, (col, ttl, desc) in enumerate(lifecycle):
        y = Inches(2.68) + i * Inches(0.78)
        add_rect(s, steps_x, y, Inches(0.46), Inches(0.62), fill=col)
        add_textbox(s, str(i+1), steps_x + Inches(0.1), y + Inches(0.1),
                    Inches(0.28), Inches(0.45),
                    size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, ttl, steps_x + Inches(0.56), y + Inches(0.03),
                    Inches(6.0), Inches(0.3), size=12.5, bold=True, color=C_DBLUE)
        add_textbox(s, desc, steps_x + Inches(0.56), y + Inches(0.33),
                    Inches(6.0), Inches(0.3), size=11, color=C_DGREY)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 12 – SELF-SERVICE PORTAL
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Self-Service Portal",
                      "Intuitive web interface for every user role")

    views = [
        ("Engineer Dashboard",     C_BLUE,
         ["My VMs / project VMs at a glance",
          "One-click power on / off / restart",
          "Snapshot create, restore, delete",
          "VM console via HTML5 viewer"]),
        ("Project Manager View",   C_ORANGE,
         ["Team resource quota status",
          "Pending approval queue",
          "Monthly spend vs. budget bar",
          "Export VM list to Excel"]),
        ("Admin Console",          RGBColor(0x5B,0x2D,0x8E),
         ["Cluster health & capacity",
          "User & role management",
          "Audit log search & export",
          "Template & policy config"]),
    ]

    vw = Inches(3.9)
    for i, (ttl, col, pts) in enumerate(views):
        x = TX + i * (vw + Inches(0.27))
        add_rect(s, x, Inches(2.2), vw, Inches(3.6), fill=C_GREY)
        add_rect(s, x, Inches(2.2), vw, Inches(0.5), fill=col)
        add_textbox(s, ttl, x + Inches(0.1), Inches(2.24),
                    vw - Inches(0.2), Inches(0.44),
                    size=13.5, bold=True, color=C_WHITE)
        for j, pt in enumerate(pts):
            add_textbox(s, "▶  " + pt,
                        x + Inches(0.15),
                        Inches(2.82) + j * Inches(0.6),
                        vw - Inches(0.3), Inches(0.55),
                        size=12, color=C_DBLUE)

    # Usage stats bottom strip
    stats = [("95%", "Self-Service Rate"),
             ("<2 min", "Avg Request Time"),
             ("0 tickets", "Manual Escalations / week"),
             ("4.7 / 5", "User Satisfaction Score")]
    sw = Inches(2.8)
    for i, (v, l) in enumerate(stats):
        x = TX + i * (sw + Inches(0.2))
        add_rect(s, x, Inches(6.02), sw, Inches(0.7), fill=C_LBLUE)
        add_textbox(s, v, x, Inches(6.04), sw, Inches(0.38),
                    size=18, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        add_textbox(s, l, x, Inches(6.42), sw, Inches(0.28),
                    size=10, color=C_DGREY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 13 – CHARGEBACK & COST
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Chargeback & Cost Management",
                      "Transparent, automated cost allocation per project and business unit")

    # How it works
    add_textbox(s, "How Chargeback Works",
                TX, Inches(2.2), Inches(5.8), Inches(0.38),
                size=14, bold=True, color=C_DBLUE)
    steps = [
        "VM resources metered every 15 minutes via vSphere API",
        "Cost rate applied: ₹0.82/vCPU·hr  ·  ₹0.41/GB-RAM·hr  ·  ₹0.03/GB-disk·hr",
        "Project code tagged at provisioning; BU rolled up automatically",
        "Monthly invoice PDF auto-generated and emailed to BU heads",
        "Finance team exports to SAP-compatible CSV for reconciliation",
    ]
    for i, st in enumerate(steps):
        add_textbox(s, f"{i+1}.  {st}",
                    TX + Inches(0.2), Inches(2.65) + i * Inches(0.5),
                    Inches(5.4), Inches(0.46),
                    size=11.5, color=C_DGREY)

    # Sample monthly table
    add_textbox(s, "Sample Monthly Summary (Jan 2026)",
                TX + Inches(6.1), Inches(2.2), Inches(6.6), Inches(0.38),
                size=14, bold=True, color=C_DBLUE)

    headers = ["Project", "VMs", "vCPU·hrs", "Spend (₹)"]
    col_xs  = [Inches(6.2), Inches(8.9), Inches(9.9), Inches(11.1)]
    col_ws  = [Inches(2.6), Inches(0.9), Inches(1.1), Inches(1.2)]

    # Header row
    for j, (hdr, x, w) in enumerate(zip(headers, col_xs, col_ws)):
        add_rect(s, x - Inches(0.05), Inches(2.65), w + Inches(0.05), Inches(0.36),
                 fill=C_BLUE)
        add_textbox(s, hdr, x, Inches(2.67), w, Inches(0.32),
                    size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)

    rows = [
        ("CaaS Platform",    "38", "27,360", "34,820"),
        ("SDN Lab",          "24", "17,280", "21,980"),
        ("AI/ML Sandbox",    "21", "15,120", "19,220"),
        ("Security Testbed", "18", "12,960", "16,480"),
        ("DevOps Pipeline",  "14", "10,080", "12,820"),
        ("Others",           "22", "11,400", "14,510"),
        ("TOTAL",            "137","94,200", "1,19,830"),
    ]
    for i, row in enumerate(rows):
        bg = C_DBLUE if i == len(rows)-1 else (C_LBLUE if i % 2 == 0 else C_GREY)
        fc = C_WHITE if i == len(rows)-1 else C_DBLUE
        for j, (cell, x, w) in enumerate(zip(row, col_xs, col_ws)):
            add_rect(s, x - Inches(0.05), Inches(3.05) + i * Inches(0.4), w + Inches(0.05), Inches(0.38), fill=bg)
            add_textbox(s, cell, x, Inches(3.07) + i * Inches(0.4), w, Inches(0.32),
                        size=11, bold=(i == len(rows)-1), color=fc,
                        align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 14 – REPORTS & ANALYTICS
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Reports & Analytics",
                      "Actionable insights for engineers, managers and executives")

    report_types = [
        ("📊  Utilisation Heatmap",  C_BLUE,
         "Real-time CPU / memory / storage gauges per cluster and per project. "
         "Drill down to individual VM in two clicks."),
        ("💰  Cost Report",           C_ORANGE,
         "Monthly chargeback statement per BU. Trend line vs. last 6 months. "
         "PDF + Excel export with SAP cost codes."),
        ("📋  Audit Trail Report",    RGBColor(0x5B,0x2D,0x8E),
         "Full lifecycle log: provisioned, powered-off, deleted. "
         "Filter by user, project, date range, action type."),
        ("📈  Capacity Forecast",     RGBColor(0x00,0x70,0x40),
         "Linear regression on 90-day consumption trends. "
         "\"Days to capacity\" alert when threshold < 30 days."),
        ("🔔  Exception Report",      RGBColor(0x9E,0x1B,0x32),
         "Over-quota VMs, idle VMs (>30 days), expired leases. "
         "Auto-emailed weekly to project owners."),
        ("📑  Compliance Report",     C_DBLUE,
         "All security events: failed logins, privilege escalations, "
         "policy violations — ready for SOC submission."),
    ]

    rw = Inches(3.9)
    rh = Inches(1.35)
    for i, (ttl, col, desc) in enumerate(report_types):
        cx = TX + (i % 3) * (rw + Inches(0.23))
        cy = Inches(2.2) + (i // 3) * (rh + Inches(0.18))
        add_rect(s, cx, cy, rw, rh, fill=C_GREY)
        add_rect(s, cx, cy, Inches(0.055), rh, fill=col)
        add_textbox(s, ttl, cx + Inches(0.15), cy + Inches(0.06),
                    rw - Inches(0.2), Inches(0.32),
                    size=12.5, bold=True, color=C_DBLUE)
        add_textbox(s, desc, cx + Inches(0.15), cy + Inches(0.4),
                    rw - Inches(0.2), Inches(0.85),
                    size=11, color=C_DGREY)

    add_textbox(s,
                "All reports auto-scheduled  ·  Role-based visibility  ·  "
                "REST API for BI tool integration (Power BI / Grafana)",
                TX, Inches(6.15), TW, Inches(0.38),
                size=11.5, italic=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 15 – VM REQUEST WORKFLOW
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "VM Request Workflow",
                      "End-to-end process from request to ready-to-use VM")

    workflow = [
        (C_BLUE,                   "Submit Request",  "Fill wizard, select specs & project"),
        (C_ORANGE,                 "Auto-Validation", "Quota check, naming policy, CMDB lookup"),
        (RGBColor(0x5B,0x2D,0x8E),"Mgr Approval",   "Auto-approved if within quota; else notify"),
        (RGBColor(0x00,0x70,0x40), "Provisioning",   "pyVmomi CloneVM_Task → <90 sec"),
        (RGBColor(0x9E,0x1B,0x32), "Post-Config",    "IP, hostname, OS config, VMware Tools"),
        (C_DBLUE,                  "Notify & Hand-off","Email creds / console URL to user"),
    ]

    step_w = Inches(1.98)
    arrow_w = Inches(0.22)
    total   = len(workflow) * step_w + (len(workflow)-1) * arrow_w
    start_x = (W - total) / 2

    for i, (col, ttl, desc) in enumerate(workflow):
        x = start_x + i * (step_w + arrow_w)
        # Box
        add_rect(s, x, Inches(2.35), step_w, Inches(1.75), fill=col)
        # Step number circle
        add_rect(s, x + step_w/2 - Inches(0.27),
                 Inches(2.28), Inches(0.54), Inches(0.54),
                 fill=C_WHITE)
        add_textbox(s, str(i+1),
                    x + step_w/2 - Inches(0.27), Inches(2.29),
                    Inches(0.54), Inches(0.5),
                    size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(s, ttl,
                    x + Inches(0.07), Inches(2.9),
                    step_w - Inches(0.14), Inches(0.48),
                    size=11.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, desc,
                    x + Inches(0.07), Inches(3.42),
                    step_w - Inches(0.14), Inches(0.62),
                    size=10, color=C_LBLUE, align=PP_ALIGN.CENTER)
        # Connector arrow (except last)
        if i < len(workflow) - 1:
            ax = x + step_w
            ay = Inches(2.35) + Inches(1.75)/2 - Inches(0.06)
            add_rect(s, ax, ay, arrow_w, Inches(0.12), fill=C_GREY)

    # SLA summary
    slas = [
        ("Standard Request", "≤ 4 hours e2e (incl. approval)"),
        ("Pre-approved Template", "< 2 minutes auto-provision"),
        ("Bulk Provisioning (≥10 VMs)", "≤ 30 minutes"),
        ("Emergency / P1 Request", "< 30 minutes (on-call admin)"),
    ]
    for i, (typ, sla) in enumerate(slas):
        x = TX + (i % 2) * Inches(6.15)
        y = Inches(4.35) + (i // 2) * Inches(0.52)
        add_rect(s, x, y, Inches(5.95), Inches(0.46), fill=C_LBLUE if i%2==0 else C_GREY)
        add_textbox(s, typ, x + Inches(0.1), y + Inches(0.08),
                    Inches(3.1), Inches(0.32), size=11.5, bold=True, color=C_DBLUE)
        add_textbox(s, sla, x + Inches(3.3), y + Inches(0.08),
                    Inches(2.55), Inches(0.32), size=11.5, color=C_BLUE)

    add_textbox(s, "All steps in the workflow are fully logged in the audit trail.",
                TX, Inches(5.55), TW, Inches(0.35),
                size=11, italic=True, color=C_DGREY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 16 – INTEGRATION POINTS
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Integration Points",
                      "LaaS Portal connects seamlessly with enterprise systems")

    integrations = [
        ("VMware vSphere 7.0",   C_BLUE,
         "pyVmomi SOAP/REST · CloneVM, Reconfig, PowerOps, Snapshot · vCenter 7 & ESXi"),
        ("Active Directory",     C_ORANGE,
         "LDAP bind for user auth · Group-to-role mapping · Auto de-provision on AD disable"),
        ("ServiceNow ITSM",      RGBColor(0x5B,0x2D,0x8E),
         "Incident auto-create for failed provisions · Change record on bulk ops · CMDB sync"),
        ("SAP FICO",             RGBColor(0x9E,0x1B,0x32),
         "Monthly cost allocation export in SAP-compatible CSV · Cost centre mapping"),
        ("Monitoring (Grafana)", RGBColor(0x00,0x70,0x40),
         "REST API metrics endpoint · VM health metrics pushed to Prometheus/Grafana stack"),
        ("SMTP / Email",         C_DBLUE,
         "Provisioning confirmations · Expiry reminders · Weekly digest · Exception alerts"),
    ]

    iw = Inches(11.93) / 2
    ih = Inches(0.95)
    for i, (ttl, col, desc) in enumerate(integrations):
        cx = TX + (i % 2) * (iw + Inches(0.24))
        cy = Inches(2.2) + (i // 2) * (ih + Inches(0.14))
        add_rect(s, cx, cy, iw, ih, fill=C_GREY)
        add_rect(s, cx, cy, Inches(0.055), ih, fill=col)
        add_textbox(s, ttl, cx + Inches(0.15), cy + Inches(0.06),
                    Inches(2.1), Inches(0.35),
                    size=12.5, bold=True, color=col)
        add_textbox(s, desc, cx + Inches(2.35), cy + Inches(0.06),
                    iw - Inches(2.5), Inches(0.8),
                    size=11, color=C_DGREY)

    add_textbox(s,
                "Future planned:  Slack notifications   ·   Jira project auto-linking   ·   "
                "HashiCorp Vault for secret rotation   ·   OpenID Connect (SSO)",
                TX, Inches(5.62), TW, Inches(0.42),
                size=11.5, italic=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 17 – ROADMAP
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Implementation Roadmap",
                      "Phase-wise delivery plan — Q4 2025 through Q4 2026")

    phases = [
        ("Phase 1\n✅ Done\n(Q4 2025)", C_BLUE,
         ["Core portal (React + FastAPI)", "vSphere VM provisioning",
          "RBAC & JWT auth", "Basic dashboard", "Chargeback MVP"]),
        ("Phase 2\n🔄 In Progress\n(Q1–Q2 2026)", C_ORANGE,
         ["Snapshot portal (done)", "Advanced analytics",
          "ServiceNow integration", "AD/LDAP auth integration",
          "Mobile PWA optimisation"]),
        ("Phase 3\nPlanned\n(Q3–Q4 2026)", RGBColor(0x5B,0x2D,0x8E),
         ["Multi-cloud (AWS/Azure bursting)", "OpenID Connect SSO",
          "HashiCorp Vault secrets", "Kubernetes cluster provisioning",
          "AI capacity forecasting"]),
    ]

    pw = Inches(3.9)
    ph = Inches(4.0)
    for i, (phase, col, items) in enumerate(phases):
        x = TX + i * (pw + Inches(0.24))
        add_rect(s, x, Inches(2.18), pw, Inches(0.72), fill=col)
        add_textbox(s, phase, x + Inches(0.1), Inches(2.2),
                    pw - Inches(0.2), Inches(0.68),
                    size=12, bold=True, color=C_WHITE)
        add_rect(s, x, Inches(2.9), pw, ph - Inches(0.72), fill=C_GREY)
        for j, itm in enumerate(items):
            add_textbox(s, "✓  " + itm,
                        x + Inches(0.15),
                        Inches(3.0) + j * Inches(0.6),
                        pw - Inches(0.3), Inches(0.55),
                        size=12, color=C_DBLUE)

    add_textbox(s,
                "Roadmap reviewed quarterly by CaaS & Wipro Infrastructure leadership.",
                TX, Inches(6.4), TW, Inches(0.35),
                size=11, italic=True, color=C_DGREY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 18 – BENEFITS & ROI
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Benefits & ROI Summary",
                      "Measurable business outcomes delivered by the LaaS Portal")

    # KPI row
    kpis2 = [
        ("80%",     "Faster Provisioning",    C_BLUE),
        ("₹2.4 Cr", "Annual Cost Savings",    C_ORANGE),
        ("95%",     "Self-Service Adoption",  RGBColor(0x00,0x70,0x40)),
        ("12",      "BUs Onboarded",          RGBColor(0x5B,0x2D,0x8E)),
        ("0",       "Manual Tickets/week",    RGBColor(0x9E,0x1B,0x32)),
    ]
    kw2 = Inches(2.35)
    for i, (val, lbl, col) in enumerate(kpis2):
        kpi_badge(s, val, lbl,
                  TX + i * (kw2 + Inches(0.19)),
                  Inches(2.2), kw2, Inches(1.18), fill=col)

    # Benefits grid
    benefits = [
        ("Operational", [
            "5-day → 2-minute provisioning cycle",
            "No email chains for VM requests",
            "Single dashboard for all teams",
            "Auto-decommission on lease expiry",
        ]),
        ("Financial", [
            "₹2.4 Cr infra savings in year 1",
            "Per-project cost transparency",
            "30% reduction in idle resources",
            "Finance-ready chargeback reports",
        ]),
        ("Strategic", [
            "ISO 27001 / SOC compliance ready",
            "Platform for multi-cloud expansion",
            "Wipro IP & reusable for other BUs",
            "Talent upskilling in cloud-native",
        ]),
    ]

    bw = Inches(3.9)
    for i, (cat, pts) in enumerate(benefits):
        x = TX + i * (bw + Inches(0.27))
        add_rect(s, x, Inches(3.55), bw, Inches(0.42), fill=C_BLUE)
        add_textbox(s, cat + " Benefits", x + Inches(0.1), Inches(3.57),
                    bw - Inches(0.2), Inches(0.38),
                    size=13, bold=True, color=C_WHITE)
        for j, pt in enumerate(pts):
            add_textbox(s, "▪  " + pt,
                        x + Inches(0.1), Inches(4.05) + j * Inches(0.52),
                        bw - Inches(0.2), Inches(0.48),
                        size=11.5, color=C_DBLUE)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 19 – SUCCESS METRICS & KPIs
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_CONTENT)
    slide_title_block(s, "Success Metrics & KPIs",
                      "How we measure and track portal performance")

    kpi_table = [
        ("KPI", "Target",    "Current",   "Trend"),
        ("Platform Uptime",         "99.5%",    "99.8%",    "✅"),
        ("VM Provision Time (avg)", "< 5 min",  "1.8 min",  "✅"),
        ("Self-Service Rate",       "> 90%",    "95%",      "✅"),
        ("Ticket Volume / week",    "< 5",      "0",        "✅"),
        ("User Satisfaction (CSAT)","≥ 4.5",   "4.7 / 5",  "✅"),
        ("Cost per VM / hr (avg)",  "< ₹45",    "₹38.2",    "✅"),
        ("Idle VM Rate",            "< 10%",    "6.2%",     "✅"),
        ("Provisioning Failures",   "< 1%",     "0.4%",     "✅"),
        ("Audit Log Coverage",      "100%",     "100%",     "✅"),
        ("Chargeback Accuracy",     "> 99%",    "99.6%",    "✅"),
    ]

    col_xs2  = [TX, TX + Inches(4.8), TX + Inches(8.4), TX + Inches(11.2)]
    col_ws2  = [Inches(4.2), Inches(3.4), Inches(2.6), Inches(1.8)]
    row_h    = Inches(0.42)

    for ri, row in enumerate(kpi_table):
        is_hdr = ri == 0
        bg = C_BLUE if is_hdr else (C_LBLUE if ri % 2 else C_GREY)
        fc = C_WHITE if is_hdr else C_DBLUE
        for ci, (cell, x, w) in enumerate(zip(row, col_xs2, col_ws2)):
            add_rect(s, x, Inches(2.22) + ri * row_h, w, row_h, fill=bg)
            add_textbox(s, cell,
                        x + Inches(0.08), Inches(2.24) + ri * row_h,
                        w - Inches(0.1), row_h - Inches(0.04),
                        size=11.5 if is_hdr else 11,
                        bold=is_hdr, color=fc,
                        align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 20 – THANK YOU (cover layout)
    # ──────────────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(LY_COVER)
    tf_map = {ph.placeholder_format.idx: ph for ph in s.placeholders}
    if 0 in tf_map:
        tf_map[0].text = "Thank You"
        tf_map[0].text_frame.paragraphs[0].runs[0].font.size  = Pt(44)
        tf_map[0].text_frame.paragraphs[0].runs[0].font.bold  = True
        tf_map[0].text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE
    if 1 in tf_map:
        tf_map[1].text = (
            "LaaS Portal — Wipro CaaS Platform\n"
            "Questions & Discussion"
        )
        for para in tf_map[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = C_LBLUE
    if 10 in tf_map:
        tf_map[10].text = "Confidential  |  Wipro Limited  |  2026"
        tf_map[10].text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        tf_map[10].text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE

    # ── SAVE ─────────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    n = len(prs.slides)
    import os
    size_kb = os.path.getsize(OUTPUT) // 1024
    print(f"✅  Saved: {OUTPUT}")
    print(f"   Slides: {n}  |  Size: {size_kb:,} KB")


if __name__ == "__main__":
    build()

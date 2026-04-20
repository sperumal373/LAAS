"""
Update LaaS_Portal_Presentation.pptx — add missing feature slides for v6.3
Adds slides for: AWS, Hyper-V, Storage, Backup, CMDB, Insights & Analytics,
History & Forecasting, Universal Search & AI, and updates Agenda + Architecture.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import copy

SRC = r'C:\caas-dashboard\LaaS_Portal_Presentation.pptx'
DST = r'C:\caas-dashboard\LaaS_Portal_Presentation_v63.pptx'

prs = Presentation(SRC)

# ─── Helper: clone a slide layout and add a content slide ────────────────────
# We'll use Slide 6 (VMware) as a template — same layout
template_slide = prs.slides[5]  # Slide 6 (0-indexed 5)
content_layout = template_slide.slide_layout  # "1_Text Content"


def add_content_slide(title, subtitle, bullets, after_index=None):
    """Add a new slide using the same layout as Slide 6.
    bullets = list of (emoji, text) tuples.
    Returns the new slide.
    """
    slide = prs.slides.add_slide(content_layout)

    # We need to add text boxes manually since the template layout may not
    # have proper placeholders. Let's create text boxes matching the template style.

    # Title
    from pptx.util import Inches, Pt
    left, top, width, height = Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Subtitle
    left, top, width, height = Inches(0.6), Inches(0.95), Inches(8.5), Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    # Bullet points
    left, top, width, height = Inches(0.6), Inches(1.55), Inches(5.8), Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (emoji, text) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{emoji}  {text}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(6)

    return slide


# ─── New Slides to Add ──────────────────────────────────────────────────────

new_slides_data = [
    # Slide: AWS
    {
        "title": "Amazon Web Services (AWS)",
        "subtitle": "EC2, S3, RDS, VPC management and Cost Explorer integration",
        "bullets": [
            ("☁️", "EC2 instance discovery across all configured AWS accounts"),
            ("⚡", "Start, Stop and Reboot EC2 instances directly from the portal"),
            ("📦", "S3 bucket inventory with size, region and access configuration"),
            ("🗄️", "RDS database instance listing with engine version and status"),
            ("🌐", "VPC and subnet discovery with IP utilisation tracking"),
            ("🔑", "IAM user and SSO session token authentication support"),
            ("💰", "AWS Cost Explorer integration with monthly cost breakdown"),
            ("📊", "Export AWS resource inventory to CSV"),
        ],
    },
    # Slide: Hyper-V
    {
        "title": "Microsoft Hyper-V",
        "subtitle": "Standalone Hyper-V host management via WinRM PowerShell remoting",
        "bullets": [
            ("🪟", "Connect to standalone Hyper-V hosts — no SCVMM or agent required"),
            ("💻", "VM inventory with CPU, memory, status and uptime per host"),
            ("⚡", "Power On / Off / Restart VMs from the portal"),
            ("📸", "Checkpoint management: create, restore and delete"),
            ("📊", "Per-host CPU and RAM utilisation bars"),
            ("🔌", "WinRM PowerShell remoting over HTTPS — secure and agentless"),
            ("🏷️", "Multi-host support with add/edit/remove host configuration"),
            ("📋", "VM and host data exportable to CSV"),
        ],
    },
    # Slide: Storage Management
    {
        "title": "Enterprise Storage Management",
        "subtitle": "Pure Storage, Dell EMC, HPE and NetApp — all via native REST APIs",
        "bullets": [
            ("💾", "Pure Storage FlashArray & FlashBlade: volumes, hosts, capacity, performance"),
            ("🗄️", "Dell EMC PowerStore & PowerScale: LUNs, file systems, NAS shares"),
            ("🔷", "HPE Alletra & Primera: volumes, drives, hosts, performance metrics"),
            ("📦", "NetApp ONTAP: aggregates, volumes, LIFs, SnapMirror relationships"),
            ("📊", "Unified capacity dashboard: used / free / total across all arrays"),
            ("🔌", "Register arrays with IP, vendor, credentials — stored securely in SQLite"),
            ("⚙️", "Per-array detail view with volume and host inventories"),
            ("📋", "Export array and volume data to CSV"),
        ],
    },
    # Slide: Backup & Data Protection
    {
        "title": "Backup & Data Protection",
        "subtitle": "Rubrik Security Cloud, Cohesity DataProtect and Veeam Backup & Replication",
        "bullets": [
            ("🛡️", "Rubrik: SLA domains, on-demand & bulk snapshots, live mount, instant recovery, file-level restore"),
            ("🔵", "Cohesity: protection jobs (run/cancel/pause/resume), object snapshots, alert management"),
            ("🟢", "Veeam: backup jobs (start/stop/enable/disable), session history, instant recovery to VMware/Hyper-V/Cloud"),
            ("📊", "Protected vs unprotected VM compliance dashboard across all three platforms"),
            ("🔄", "Job monitoring with retry, cancel and real-time status tracking"),
            ("💾", "Repository and cluster capacity monitoring with free space alerts"),
            ("🔍", "Search protected objects across all connected backup platforms"),
            ("📋", "Job history and compliance reports exportable to CSV"),
        ],
    },
    # Slide: CMDB
    {
        "title": "CMDB — Configuration Management Database",
        "subtitle": "ServiceNow-aligned CI registry with auto-discovery across all eight platforms",
        "bullets": [
            ("🗄️", "1800+ Configuration Items auto-discovered from VMware, Hyper-V, Nutanix, AWS, Storage, OCP, Physical Assets and IPAM"),
            ("🏷️", "ServiceNow CI classes: vm_instance, esx_server, win_server, storage_device, ocp_cluster, ip_network, server"),
            ("🔗", "Correlation ID deduplication — no duplicate CIs across repeated collections"),
            ("📤", "One-click push to ServiceNow CMDB Table API with dry-run preview"),
            ("✏️", "Inline CI editing for admins — name, department, environment, serial number"),
            ("🔍", "Tab-based filtering by CI class and full-text search across all fields"),
            ("🔄", "Daily automated collection via the Insights engine scheduler"),
            ("📊", "PostgreSQL-backed CI store for audit-ready IT asset compliance"),
        ],
    },
    # Slide: Insights & Analytics
    {
        "title": "Insights & Analytics Engine",
        "subtitle": "PostgreSQL 16 daily snapshots — Health Scorecard, Change Detection, Capacity & Cost, Executive Summary",
        "bullets": [
            ("🩺", "Health Scorecard: real-time RAG status (Red/Amber/Green) across all eight platforms with trend arrows"),
            ("🔄", "Change Detection: automatic day-over-day drift alerts for VM count, host count, pod count and IP utilisation"),
            ("📊", "Capacity & Cost: storage and compute exhaustion predictions with monthly chargeback in ₹ and $"),
            ("📋", "Executive Summary: single-page, PDF-ready overview with aggregate KPIs and platform health matrix"),
            ("🛡️", "Backup Analytics: Rubrik, Cohesity and Veeam job success/failure trends and compliance tracking"),
            ("💾", "Storage Analytics: array capacity trends, volume growth tracking and exhaustion predictions"),
            ("🕐", "Automated daily collector at 11 PM with retry logic and thirty-day rolling retention"),
            ("📈", "Twelve PostgreSQL tables capturing data from all platforms every night"),
        ],
    },
    # Slide: History & Forecasting
    {
        "title": "Historical Trending & Capacity Forecasting",
        "subtitle": "Interactive charts and linear regression predictions for resource planning",
        "bullets": [
            ("📈", "Interactive trend charts for 7, 14 and 30-day windows across all collected metrics"),
            ("🖥️", "VM count, host count, pod count and IP utilisation historical trends"),
            ("💾", "Storage used %, compute used % and network utilisation day-by-day graphs"),
            ("🔮", "Linear regression capacity forecasting per metric"),
            ("⚠️", "Storage, compute and network exhaustion date predictions"),
            ("📊", "Day-over-day change arrows with percentage shifts"),
            ("📋", "CSV-exportable historical data snapshots for offline analysis"),
            ("🔄", "Powered by the same PostgreSQL daily collector — no additional configuration needed"),
        ],
    },
    # Slide: Universal Search & AI Assistant
    {
        "title": "Universal Search & AI Assistant",
        "subtitle": "Instant search across all platforms and natural-language infrastructure queries",
        "bullets": [
            ("🔍", "Universal Search: VMs, hosts, datastores, snapshots, networks, OCP clusters, Nutanix PCs, Ansible jobs"),
            ("📑", "All 20+ navigation pages searchable by name and description"),
            ("📞", "Support contacts searchable — click to copy email or open mail"),
            ("⌨️", "Keyboard navigation: ↑↓ to browse, Enter to jump, ESC to close"),
            ("🤖", "AI Assistant: ask natural-language questions about your infrastructure"),
            ("🎤", "Voice input supported — click the microphone and speak your query"),
            ("📊", "Powered entirely by live portal data — fully offline-safe, no external LLM"),
            ("💡", "Example queries: 'which host has the most free RAM?', 'how many pods in cluster A?'"),
        ],
    },
]

# ─── Add all new slides ─────────────────────────────────────────────────────
print(f"Original slide count: {len(prs.slides)}")

for sd in new_slides_data:
    add_content_slide(sd["title"], sd["subtitle"], sd["bullets"])
    print(f"  Added slide: {sd['title']}")

print(f"New slide count: {len(prs.slides)}")

# ─── Update Slide 1 (Cover) — version ───────────────────────────────────────
cover = prs.slides[0]
for shape in cover.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if "Q1 2026" in para.text:
                for run in para.runs:
                    if "Q1 2026" in run.text:
                        run.text = run.text.replace("Q1 2026", "Q2 2026  |  v6.3")
                        print("  Updated cover: Q1 2026 -> Q2 2026 | v6.3")

# ─── Update Slide 2 (Agenda) — add new items ────────────────────────────────
agenda = prs.slides[1]
# Find the last text box and add new agenda items
# We need to add items 16-23 for the new slides
new_agenda_items = [
    ("16", "Amazon Web Services (AWS)"),
    ("17", "Microsoft Hyper-V"),
    ("18", "Enterprise Storage Management"),
    ("19", "Backup & Data Protection"),
    ("20", "CMDB — Configuration Management Database"),
    ("21", "Insights & Analytics Engine"),
    ("22", "Historical Trending & Capacity Forecasting"),
    ("23", "Universal Search & AI Assistant"),
]

# Add a text box for new agenda items below the existing ones
left, top, width, height = Inches(0.6), Inches(6.0), Inches(9.0), Inches(3.0)
txBox = agenda.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

for i, (num, label) in enumerate(new_agenda_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"{num}    {label}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    p.font.bold = False
    p.space_after = Pt(2)

# Also bold the numbers
print("  Updated agenda with 8 new items")

# ─── Update Slide 3 (Architecture) — add missing platforms ──────────────────
arch = prs.slides[2]
# Add AWS, Hyper-V, Storage, Backup boxes
new_platform_boxes = [
    ("AWS", "▪ EC2, S3, RDS, VPC management\n▪ Cost Explorer integration\n▪ IAM & SSO support\n▪ Start/Stop/Reboot EC2"),
    ("Hyper-V", "▪ WinRM PowerShell remoting\n▪ VM power & checkpoints\n▪ Per-host CPU/RAM bars\n▪ Agentless — no SCVMM"),
    ("Storage", "▪ Pure, Dell, HPE, NetApp arrays\n▪ Volumes, LUNs, capacity\n▪ Native REST API integration\n▪ Multi-vendor unified view"),
    ("Backup", "▪ Rubrik, Cohesity, Veeam\n▪ SLA & job management\n▪ Instant recovery & restore\n▪ Compliance dashboard"),
]

y_start = Inches(5.3)
x_start = Inches(0.6)
box_w = Inches(2.3)
box_h = Inches(1.5)
gap = Inches(0.15)

for i, (title, text) in enumerate(new_platform_boxes):
    left = x_start + i * (box_w + gap)
    txBox = arch.shapes.add_textbox(left, y_start, box_w, box_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
    p.space_after = Pt(4)

    for line in text.split("\n"):
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p2.space_after = Pt(1)

print("  Updated architecture slide with AWS, Hyper-V, Storage, Backup")

# ─── Update Slide 4 (Tech Stack) — add PostgreSQL, Storage APIs ─────────────
tech = prs.slides[3]
# Add a box for new tech
left, top = Inches(0.6), Inches(5.5)
txBox = tech.shapes.add_textbox(left, top, Inches(9.0), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "New in v6.3"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
p.space_after = Pt(4)

for line in [
    "• PostgreSQL 16 — daily snapshot DB for Insights, CMDB, IPAM",
    "• AWS SDK (boto3) — EC2, S3, RDS, VPC, Cost Explorer",
    "• WinRM/pywinrm — Hyper-V host remoting",
    "• Storage REST APIs — Pure, Dell, HPE, NetApp",
    "• Rubrik/Cohesity/Veeam APIs — Backup & Data Protection",
    "• ServiceNow Table API — CMDB CI push integration",
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2.space_after = Pt(1)

print("  Updated tech stack slide with PostgreSQL, AWS SDK, WinRM, Storage/Backup APIs")

# ─── Update Slide 5 (Overview) — add Hyper-V, AWS, Storage, CMDB rings ──────
overview_slide = prs.slides[4]
left, top = Inches(0.6), Inches(5.2)
txBox = overview_slide.shapes.add_textbox(left, top, Inches(9.0), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
for line_data in [
    ("☁️", "AWS EC2 health status and instance count"),
    ("🪟", "Hyper-V host connection status and VM count"),
    ("💾", "Storage array health across Pure, Dell, HPE, NetApp"),
    ("🗄️", "CMDB CI count and last collection status"),
]:
    p = tf.add_paragraph()
    p.text = f"{line_data[0]}  {line_data[1]}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(4)

print("  Updated overview slide with AWS, Hyper-V, Storage, CMDB")

# ─── Update Thank You slide — version ───────────────────────────────────────
thankyou = prs.slides[-1]  # last slide (may have shifted)
for shape in thankyou.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "2026" in run.text and "Confidential" in run.text:
                    run.text = "Confidential  |  Wipro Limited  |  Q2 2026  |  v6.3"
                    print("  Updated Thank You slide version")

# ─── Save ───────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved updated presentation: {DST}")
print(f"Total slides: {len(prs.slides)}")

"""Diagnose issues in the generated PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

prs = Presentation(r'C:\caas-dashboard\LaaS_Portal_Presentation.pptx')
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width.emu/914400:.2f}\" x {prs.slide_height.emu/914400:.2f}\"")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} (layout: {slide.slide_layout.name}) ---")
    shapes = list(slide.shapes)
    print(f"  Total shapes: {len(shapes)}")
    
    issues = []
    for sh in shapes:
        # Check for shapes outside slide bounds
        W = prs.slide_width.emu
        H = prs.slide_height.emu
        if sh.left is not None and sh.top is not None:
            if sh.left < -Inches(0.5) or sh.left + sh.width > W + Inches(0.5):
                issues.append(f"  ⚠ '{sh.name}' out of bounds X: left={sh.left/914400:.2f} w={sh.width/914400:.2f}")
            if sh.top < -Inches(0.5) or sh.top + sh.height > H + Inches(0.5):
                issues.append(f"  ⚠ '{sh.name}' out of bounds Y: top={sh.top/914400:.2f} h={sh.height/914400:.2f}")
        # Check for zero-size shapes
        if sh.width is not None and sh.width <= 0:
            issues.append(f"  ⚠ '{sh.name}' zero/negative width")
        if sh.height is not None and sh.height <= 0:
            issues.append(f"  ⚠ '{sh.name}' zero/negative height")
    
    # Check for layout placeholder conflicts
    layout_phs = {ph.placeholder_format.idx for ph in slide.slide_layout.placeholders}
    slide_phs = {ph.placeholder_format.idx for ph in slide.placeholders}
    
    # Check timing XML
    timing = slide._element.find(qn('p:timing'))
    if timing is not None:
        print(f"  ✓ Animation timing present")
    
    if issues:
        for iss in issues:
            print(iss)
    else:
        print(f"  ✓ No obvious positioning issues")
    
    # Show logo placement
    pics = [sh for sh in shapes if sh.shape_type == 13]
    if pics:
        for p in pics:
            print(f"  📷 Picture: left={p.left/914400:.3f} top={p.top/914400:.3f} w={p.width/914400:.3f} h={p.height/914400:.3f}")
    else:
        if i > 0:  # slide 1 has logo
            print(f"  ❌ No picture/logo found!")

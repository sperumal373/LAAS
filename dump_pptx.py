"""Dump all text from the LaaS PPTX presentation to review content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'C:\caas-dashboard\LaaS_Portal_Presentation.pptx')

print(f"Slide count: {len(prs.slides)}")
print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print("="*100)

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*100}")
    print(f"SLIDE {i+1}  (layout: {slide.slide_layout.name})")
    print(f"{'='*100}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  [{shape.shape_type}] {text[:200]}")
        if shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                print(f"  [TABLE row {row_idx}] {' | '.join(cells)}")
    # Check for images
    img_count = sum(1 for s in slide.shapes if s.shape_type == 13)
    if img_count:
        print(f"  [{img_count} image(s)]")

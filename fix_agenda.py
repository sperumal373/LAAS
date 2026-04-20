"""Fix the Agenda slide to show all 26 topics with correct numbering."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC = r'C:\caas-dashboard\LaaS_Portal_Presentation_v63.pptx'

prs = Presentation(SRC)
agenda = prs.slides[1]  # Slide 2

# Remove all existing shapes from agenda slide
shapes_to_remove = []
for shape in agenda.shapes:
    shapes_to_remove.append(shape)

sp_tree = agenda.shapes._spTree
for shape in shapes_to_remove:
    sp_tree.remove(shape._element)

# Rebuild the agenda from scratch
# Title
txBox = agenda.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Agenda"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Subtitle
txBox = agenda.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(8.5), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Capabilities covered in this presentation — LaaS Portal v6.3"
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

# All agenda items — matching the slide order
items = [
    "Platform Overview & Architecture",
    "Multi-Platform Dashboard — Overview",
    "VMware vSphere — VM Management",
    "VM Snapshots & Network Inventory",
    "Capacity Planning & Resource Monitoring",
    "Project Utilization Tracking",
    "Chargeback & Cost Management",
    "VM Request & Approval Workflow",
    "IPAM — IP Address Management",
    "Asset Inventory Management",
    "Red Hat OpenShift Container Platform",
    "Nutanix AHV Virtualisation",
    "Ansible Automation Platform (AAP)",
    "Amazon Web Services (AWS)",
    "Microsoft Hyper-V",
    "Enterprise Storage Management",
    "Backup & Data Protection",
    "CMDB — Configuration Management Database",
    "Active Directory & DNS Management",
    "User Management & RBAC / Audit Log",
    "Universal Search & AI Assistant",
    "Insights & Analytics Engine",
    "Historical Trending & Capacity Forecasting",
]

# Split into 2 columns
col1 = items[:12]
col2 = items[12:]

for col_idx, col_items in enumerate([col1, col2]):
    left = Inches(0.6) if col_idx == 0 else Inches(5.2)
    top = Inches(1.3)
    txBox = agenda.shapes.add_textbox(left, top, Inches(4.3), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(col_items):
        num = i + 1 if col_idx == 0 else i + 13
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Number
        run_num = p.add_run()
        run_num.text = f"{num:02d}  "
        run_num.font.size = Pt(13)
        run_num.font.bold = True
        run_num.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)

        # Label
        run_label = p.add_run()
        run_label.text = item
        run_label.font.size = Pt(13)
        run_label.font.bold = False
        run_label.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        p.space_after = Pt(5)

        # Highlight new items (14+)
        if num >= 14 and num <= 18:
            run_new = p.add_run()
            run_new.text = "  NEW"
            run_new.font.size = Pt(9)
            run_new.font.bold = True
            run_new.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
        elif num >= 21:
            run_new = p.add_run()
            run_new.text = "  NEW"
            run_new.font.size = Pt(9)
            run_new.font.bold = True
            run_new.font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)

prs.save(SRC)
print("Agenda slide rebuilt with 23 items in 2 columns")
print(f"Saved: {SRC}")

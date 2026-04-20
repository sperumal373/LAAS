"""
Reorder slides in LaaS_Portal_Presentation_v63.pptx so new slides are
in the correct position (before Thank You) and grouped logically.

Desired order:
 1. Cover
 2. Agenda (updated)
 3. Platform Overview & Architecture
 4. Technology Stack
 5. Multi-Platform Dashboard — Overview
 6. VMware vSphere — VM Management
 7. VM Snapshot Management
 8. Network Inventory
 9. Capacity Planning & Resource Monitoring
10. Project Utilization Tracking
11. Chargeback & Cost Management
12. VM Request & Approval Workflow
13. IPAM — IP Address Management
14. Asset Inventory Management
15. Red Hat OpenShift Container Platform
16. Nutanix AHV Virtualisation
17. Ansible Automation Platform (AAP)
18. Amazon Web Services (AWS)           ← was 21
19. Microsoft Hyper-V                   ← was 22
20. Enterprise Storage Management       ← was 23
21. Backup & Data Protection            ← was 24
22. CMDB — Configuration Management DB  ← was 25
23. Active Directory & DNS Management   ← was 18
24. User Management & RBAC             ← was 19
25. Universal Search & AI Assistant     ← was 28
26. Insights & Analytics Engine         ← was 26
27. Historical Trending & Forecasting   ← was 27
28. Thank You                           ← was 20
"""
from pptx import Presentation
from lxml import etree
import copy

SRC = r'C:\caas-dashboard\LaaS_Portal_Presentation_v63.pptx'
DST = r'C:\caas-dashboard\LaaS_Portal_Presentation_v63.pptx'

prs = Presentation(SRC)

# Current slide indices (0-based):
# 0=Cover, 1=Agenda, 2=Architecture, 3=Tech, 4=Overview,
# 5=VMware, 6=Snapshot, 7=Network, 8=Capacity, 9=ProjUtil,
# 10=Chargeback, 11=VMReq, 12=IPAM, 13=Assets,
# 14=OpenShift, 15=Nutanix, 16=Ansible,
# 17=AD&DNS, 18=UserMgmt, 19=ThankYou,
# 20=AWS, 21=HyperV, 22=Storage, 23=Backup, 24=CMDB,
# 25=Insights, 26=History, 27=Search

desired_order = [
    0,   # Cover
    1,   # Agenda
    2,   # Architecture
    3,   # Tech Stack
    4,   # Overview
    5,   # VMware
    6,   # Snapshots
    7,   # Networks
    8,   # Capacity
    9,   # Project Util
    10,  # Chargeback
    11,  # VM Request
    12,  # IPAM
    13,  # Assets
    14,  # OpenShift
    15,  # Nutanix
    16,  # Ansible
    20,  # AWS (was 21st slide)
    21,  # Hyper-V (was 22nd)
    22,  # Storage (was 23rd)
    23,  # Backup (was 24th)
    24,  # CMDB (was 25th)
    17,  # AD & DNS (was 18th)
    18,  # User Mgmt (was 19th)
    27,  # Search & AI (was 28th)
    25,  # Insights (was 26th)
    26,  # History (was 27th)
    19,  # Thank You (was 20th)
]

# Reorder using XML manipulation
slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)

# Create new order
reordered = [slide_ids[i] for i in desired_order]

# Clear and re-add
for child in list(slide_list):
    slide_list.remove(child)

for elem in reordered:
    slide_list.append(elem)

# Verify
prs.save(DST)

# Reopen and print order
prs2 = Presentation(DST)
for i, slide in enumerate(prs2.slides):
    title = ''
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    title = t
                    break
        if title:
            break
    print(f"Slide {i+1:2d}: {title[:90]}")

print(f"\nTotal: {len(prs2.slides)} slides — saved to {DST}")

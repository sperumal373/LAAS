#!/usr/bin/env python3
"""
generate_pptx_v3.py  –  LaaS Portal Executive Presentation
- Real Wipro template branding (ZIP-surgery approach, no orphaned parts)
- Only ACTUAL features present in the dashboard
- Real screenshots of every page embedded
"""
import io, re, zipfile, os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = r"C:\Users\Administrator\Desktop\Wipro Template.pptx"
OUTPUT   = r"C:\caas-dashboard\LaaS_Portal_Presentation.pptx"
SS_DIR   = r"C:\caas-dashboard\screenshots"

# ── Wipro brand palette ───────────────────────────────────────────────────────
W_BLUE   = RGBColor(0x00, 0x70, 0xC0)
W_DARK   = RGBColor(0x1F, 0x39, 0x64)
W_LBLUE  = RGBColor(0xDE, 0xEA, 0xF6)
W_GREY   = RGBColor(0xF2, 0xF2, 0xF2)
W_DGREY  = RGBColor(0x59, 0x59, 0x59)
W_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
W_ORANGE = RGBColor(0xFF, 0x6B, 0x00)
W_GREEN  = RGBColor(0x00, 0x70, 0x40)
W_RED    = RGBColor(0x9E, 0x1B, 0x32)
W_CYAN   = RGBColor(0x06, 0xB6, 0xD4)
W_PURPLE = RGBColor(0x5B, 0x2D, 0x8E)

# ── Slide canvas ──────────────────────────────────────────────────────────────
SW = Inches(13.33)   # slide width
SH = Inches(7.5)     # slide height
# Safe content zone (Wipro layouts reserve bottom ~0.45" for footer)
CX = Inches(0.38)    # content left margin
CY = Inches(0.20)    # content top (layout has no header bar in 1_Text Content)
CW = Inches(12.57)   # content width
FB = Inches(6.85)    # footer boundary (content must stay above this)


# ═════════════════════════════════════════════════════════════════════════════
# STEP 1  –  ZIP surgery: strip 130 template slides, keep master + layouts
# ═════════════════════════════════════════════════════════════════════════════
def build_clean_prs(template_path: str) -> Presentation:
    NS_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    _sld   = re.compile(r"^ppt/slides/(slide\d+\.xml|_rels/slide\d+\.xml\.rels)$")

    buf = io.BytesIO()
    with zipfile.ZipFile(template_path, "r") as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            name = item.filename
            if _sld.match(name):
                continue
            data = zin.read(name)
            if name == "ppt/presentation.xml":
                root = etree.fromstring(data)
                lst  = root.find(f"{{{NS_PML}}}sldIdLst")
                if lst is not None:
                    for el in list(lst): lst.remove(el)
                data = etree.tostring(root, xml_declaration=True,
                                      encoding="UTF-8", standalone=True)
            elif name == "ppt/_rels/presentation.xml.rels":
                root = etree.fromstring(data)
                for r in [r for r in root if r.get("Type","").endswith("/slide")]:
                    root.remove(r)
                data = etree.tostring(root, xml_declaration=True,
                                      encoding="UTF-8", standalone=True)
            elif name == "[Content_Types].xml":
                root = etree.fromstring(data)
                for el in [el for el in root
                           if re.match(r"/ppt/slides/slide\d+\.xml",
                                       el.get("PartName",""))]:
                    root.remove(el)
                data = etree.tostring(root, xml_declaration=True,
                                      encoding="UTF-8", standalone=True)
            zout.writestr(item, data)
    buf.seek(0)
    return Presentation(buf)


# ═════════════════════════════════════════════════════════════════════════════
# DRAWING HELPERS
# ═════════════════════════════════════════════════════════════════════════════
def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        if lw: shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp


def tb(slide, text, x, y, w, h,
       sz=13, bold=False, italic=False, color=W_DARK,
       align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text       = text
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.italic= italic
    r.font.color.rgb = color
    return txb


def para(tf, text, sz=12, bold=False, color=W_DGREY, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """Standard Wipro-style header: accent line + title."""
    rect(slide, CX, CY, CW, Inches(0.055), fill=W_BLUE)
    tb(slide, title,
       CX, CY + Inches(0.07), CW, Inches(0.6),
       sz=26, bold=True, color=W_DARK)
    if subtitle:
        tb(slide, subtitle,
           CX, CY + Inches(0.67), CW, Inches(0.32),
           sz=12, italic=True, color=W_DGREY)


def screenshot_slide(slide, ss_file, title, subtitle, bullets):
    """
    Standard screenshot layout:
      - Wipro-style header
      - Screenshot on left (65% width)
      - Bullet points on right
    """
    slide_header(slide, title, subtitle)

    ss_path = os.path.join(SS_DIR, ss_file)
    content_top = CY + Inches(1.08)
    content_h   = FB - content_top - Inches(0.05)

    if os.path.exists(ss_path):
        # Screenshot: 1600×900 → keep 16:9 ratio
        ss_w = Inches(8.6)
        ss_h = ss_w * 900 / 1600          # ≈ 4.84"
        if ss_h > content_h:
            ss_h = content_h
            ss_w = ss_h * 1600 / 900
        ss_x = CX
        ss_y = content_top
        # Thin border frame
        rect(slide, ss_x - Inches(0.03), ss_y - Inches(0.03),
             ss_w + Inches(0.06), ss_h + Inches(0.06),
             fill=W_LBLUE, line=W_BLUE, lw=0.5)
        slide.shapes.add_picture(ss_path, ss_x, ss_y, ss_w, ss_h)
        bullet_x = ss_x + ss_w + Inches(0.25)
    else:
        bullet_x = CX

    # Bullets column
    bw = SW - bullet_x - Inches(0.2)
    if bw > Inches(0.5):
        txb = slide.shapes.add_textbox(bullet_x, content_top, bw, content_h)
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for icon, text in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            # Icon run
            ri = p.add_run()
            ri.text = icon + "  "
            ri.font.size = Pt(11); ri.font.color.rgb = W_BLUE
            # Text run
            rt = p.add_run()
            rt.text = text
            rt.font.size = Pt(11); rt.font.color.rgb = W_DARK
            # Spacer
            if not first:
                ps = tf.add_paragraph()
                ps.add_run().text = ""
                ps.runs[0].font.size = Pt(4)


# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════
def build():
    prs = build_clean_prs(TEMPLATE)
    LY_COVER   = prs.slide_layouts[0]   # Cover_WiproBlue
    LY_CONTENT = prs.slide_layouts[9]   # 1_Text Content

    # ── helper: populate Cover_WiproBlue placeholders ────────────────────────
    def cover(title, sub1, sub2=None):
        s = prs.slides.add_slide(LY_COVER)
        ph = {p.placeholder_format.idx: p for p in s.placeholders}
        if 0 in ph:
            ph[0].text = title
            for run in ph[0].text_frame.paragraphs[0].runs:
                run.font.size = Pt(40); run.font.bold = True
                run.font.color.rgb = W_WHITE
        if 1 in ph:
            ph[1].text = sub1
            for run in ph[1].text_frame.paragraphs[0].runs:
                run.font.size = Pt(21); run.font.color.rgb = W_LBLUE
        if 10 in ph:
            ph[10].text = sub2 or ""
            for run in ph[10].text_frame.paragraphs[0].runs:
                run.font.size = Pt(12); run.font.color.rgb = W_WHITE
        return s

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 – COVER
    # ══════════════════════════════════════════════════════════════════════════
    cover(
        "LaaS Portal",
        "Lab as a Service — Wipro CaaS Platform",
        "Confidential  |  Q1 2026"
    )

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 – AGENDA
    # ══════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(LY_CONTENT)
    slide_header(s, "Agenda", "Capabilities covered in this presentation")

    topics = [
        ("01", "Platform Overview & Architecture"),
        ("02", "Multi-Platform Dashboard — Overview"),
        ("03", "VMware vSphere — VM Management"),
        ("04", "VM Snapshots & Network Inventory"),
        ("05", "Capacity Planning & Resource Monitoring"),
        ("06", "Project Utilization Tracking"),
        ("07", "Chargeback & Cost Management"),
        ("08", "VM Request & Approval Workflow"),
        ("09", "IPAM — IP Address Management"),
        ("10", "Asset Inventory Management"),
        ("11", "Red Hat OpenShift Container Platform"),
        ("12", "Nutanix AHV Virtualisation"),
        ("13", "Ansible Automation Platform (AAP)"),
        ("14", "Active Directory & DNS Management"),
        ("15", "Audit Log & User Management / RBAC"),
    ]

    cw = Inches(5.9)
    ty = CY + Inches(1.05)
    for i, (num, topic) in enumerate(topics):
        col = i // 8
        row = i % 8
        x = CX + col * (cw + Inches(0.4))
        y = ty + row * Inches(0.68)
        rect(s, x, y, cw, Inches(0.58),
             fill=W_LBLUE if i % 2 == 0 else W_GREY)
        tb(s, num, x + Inches(0.1), y + Inches(0.1),
           Inches(0.5), Inches(0.42),
           sz=14, bold=True, color=W_BLUE)
        tb(s, topic, x + Inches(0.65), y + Inches(0.13),
           cw - Inches(0.75), Inches(0.38),
           sz=12, color=W_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 – ABOUT THE PLATFORM
    # ══════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(LY_CONTENT)
    slide_header(s, "LaaS Portal — Lab as a Service",
                 "A unified self-service portal for multi-platform lab infrastructure")

    platforms = [
        (W_BLUE,   "VMware vSphere",
         ["Full VM lifecycle management", "ESXi host & datastore monitoring",
          "Cloning, migration, reconfiguration", "Snapshot & network management"]),
        (W_RED,    "Red Hat OpenShift",
         ["OCP cluster overview & health", "Node & operator visibility",
          "Container workload monitoring", "Multi-cluster support"]),
        (W_GREEN,  "Nutanix AHV",
         ["Prism Central integration", "AHV VM inventory & metrics",
          "Cluster resource visibility", "Image management"]),
        (W_ORANGE, "Ansible AAP",
         ["Automation job templates", "Playbook execution tracking",
          "Workflow status visibility", "AAP API integration"]),
        (W_PURPLE, "AD & DNS",
         ["Active Directory user management", "Group & OU management",
          "DNS zone & record management", "Cache flush & PTR support"]),
        (W_CYAN,   "IPAM & Assets",
         ["IP subnet tracking & allocation", "Asset inventory (physical/virtual)",
          "Ping & power actions on assets", "Export to CSV/Excel"]),
    ]

    pw_ = Inches(3.9)
    ph_ = Inches(1.9)
    for i, (col, ttl, pts) in enumerate(platforms):
        cx2 = CX + (i % 3) * (pw_ + Inches(0.22))
        cy2 = CY + Inches(1.05) + (i // 3) * (ph_ + Inches(0.1))
        rect(s, cx2, cy2, pw_, Inches(0.42), fill=col)
        tb(s, ttl, cx2 + Inches(0.1), cy2 + Inches(0.05),
           pw_ - Inches(0.2), Inches(0.35),
           sz=13, bold=True, color=W_WHITE)
        rect(s, cx2, cy2 + Inches(0.42), pw_, ph_ - Inches(0.42),
             fill=W_GREY)
        for j, pt in enumerate(pts):
            tb(s, "▪ " + pt,
               cx2 + Inches(0.12),
               cy2 + Inches(0.52) + j * Inches(0.32),
               pw_ - Inches(0.24), Inches(0.3),
               sz=11, color=W_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 – TECHNICAL STACK (ACTUAL)
    # ══════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(LY_CONTENT)
    slide_header(s, "Technology Stack",
                 "The actual components that power the LaaS Portal")

    stack = [
        ("Frontend",        W_BLUE,
         ["React 18 + Vite", "TailwindCSS", "Chart.js / Recharts", "Playwright (testing)"]),
        ("Backend",         W_ORANGE,
         ["Python FastAPI", "SQLAlchemy / SQLite", "Pydantic validation", "JWT Auth (LDAP/Local)"]),
        ("VMware Layer",    W_BLUE,
         ["pyVmomi SDK", "vSphere 7 REST API", "vCenter multi-VC", "ESXi 7 hosts"]),
        ("Platforms",       W_RED,
         ["OCP via REST API", "Nutanix Prism Central", "Ansible AAP API", "ldap3 (AD)"]),
        ("Networking/IPAM", W_GREEN,
         ["IPAM REST integration", "DNS (dnspython)", "AD/LDAP (ldap3)", "SMTP notifications"]),
        ("Auth & Security", W_PURPLE,
         ["AD LDAP bind auth", "Local admin fallback", "4-tier RBAC model",
          "Session token (sessionStorage)"]),
    ]

    pw2 = Inches(3.9); ph2 = Inches(1.85)
    for i, (cat, col, items) in enumerate(stack):
        cx2 = CX + (i % 3) * (pw2 + Inches(0.24))
        cy2 = CY + Inches(1.05) + (i // 3) * (ph2 + Inches(0.12))
        rect(s, cx2, cy2, pw2, Inches(0.38), fill=col)
        tb(s, cat, cx2 + Inches(0.1), cy2 + Inches(0.04),
           pw2 - Inches(0.2), Inches(0.32), sz=13, bold=True, color=W_WHITE)
        rect(s, cx2, cy2 + Inches(0.38), pw2, ph2 - Inches(0.38), fill=W_GREY)
        for j, itm in enumerate(items):
            tb(s, "• " + itm,
               cx2 + Inches(0.1), cy2 + Inches(0.46) + j * Inches(0.33),
               pw2 - Inches(0.2), Inches(0.3), sz=11, color=W_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDES 5–19 — SCREENSHOT SLIDES
    # ══════════════════════════════════════════════════════════════════════════
    screen_slides = [
        # (screenshot_file, title, subtitle, bullets)
        (
            "02_overview.png",
            "Multi-Platform Dashboard — Overview",
            "Home page: KPIs, platform health rings and IPAM summary",
            [
                ("🖥️", "VMware vCenter KPI cards: VMs, Hosts, Alerts, Storage"),
                ("🔴", "Red Hat OpenShift cluster health ring"),
                ("🟩", "Nutanix Prism Central health ring"),
                ("📡", "IPAM subnet overview with top subnets"),
                ("🔔", "Live alert ticker — errors and warnings"),
                ("🔍", "Global search across VMs, hosts, snapshots, networks"),
                ("🤖", "AI Chatbot for infrastructure queries"),
            ]
        ),
        (
            "03_vmware_vms.png",
            "VMware vSphere — VM Management",
            "Full VM inventory with per-vCenter tabs, power actions and search",
            [
                ("🗂️", "Per-vCenter tab navigation with live status"),
                ("🔍", "Search and filter VMs by name, IP, OS, host"),
                ("⚡", "Power On / Off / Restart from the table"),
                ("📸", "Create snapshot directly from VM row"),
                ("🔁", "Clone VM from existing templates"),
                ("⚙️", "Reconfigure vCPU, Memory, Disk per VM"),
                ("🗑️", "Delete VM from disk (admin only)"),
                ("📊", "CSV export of full VM inventory"),
            ]
        ),
        (
            "04_snapshots.png",
            "VM Snapshot Management",
            "Create, restore and delete snapshots with retention visibility",
            [
                ("📸", "List all snapshots across all vCenters"),
                ("✅", "Create snapshot with description from portal"),
                ("↩️", "Restore VM to any snapshot point"),
                ("🗑️", "Delete individual snapshots"),
                ("🔍", "Search snapshots by VM name or description"),
                ("📅", "Snapshot age and creation date visible"),
                ("📊", "Export snapshot list to CSV"),
            ]
        ),
        (
            "05_networks.png",
            "Network Inventory",
            "Port group and network inventory across all vCenters",
            [
                ("🌐", "All port groups listed per vCenter"),
                ("🏷️", "Network name, type (Standard/DVS) visible"),
                ("📡", "Associated datacenter and cluster shown"),
                ("🔍", "Filter and search by network name"),
                ("📊", "Export network list to CSV"),
            ]
        ),
        (
            "06_capacity.png",
            "Capacity Planning & Resource Monitoring",
            "ESXi host and datastore capacity across vCenters + OCP + Nutanix",
            [
                ("🗄️", "Per-ESXi host: CPU %, Memory %, connection status"),
                ("💾", "Per-datastore: capacity GB, used GB, free %"),
                ("🔴", "OpenShift node capacity summary"),
                ("🟩", "Nutanix AHV cluster resource summary"),
                ("🔍", "Search hosts and datastores"),
                ("⚠️", "Visual warning on high utilisation"),
                ("📊", "Export host/datastore data to CSV"),
            ]
        ),
        (
            "07_project_util.png",
            "Project Utilization Tracking",
            "Per-project VM resource consumption by vCenter",
            [
                ("📊", "Bar charts of vCPU and memory per project"),
                ("🏷️", "VMs tagged by project code in vSphere"),
                ("👤", "Project owner shown alongside resource usage"),
                ("✏️", "Edit project tag and owner inline"),
                ("📁", "Group by project across VMs on same vCenter"),
                ("📊", "Export project utilization to CSV"),
            ]
        ),
        (
            "08_chargeback.png",
            "Chargeback & Cost Management",
            "Cost rate configuration per platform — VMware, OpenShift, Nutanix",
            [
                ("💳", "Per-platform pricing tabs: VMware / OpenShift / Nutanix"),
                ("⚙️", "Admin-configurable rates: ₹/vCPU·hr, ₹/GB-RAM·hr, ₹/GB-disk·hr"),
                ("🔒", "Viewer role sees rates as read-only"),
                ("💾", "Rates saved to backend via REST API"),
                ("📋", "Pricing persisted in database"),
            ]
        ),
        (
            "09_requests.png",
            "VM Request & Approval Workflow",
            "Self-service VM provisioning with manager approval queue",
            [
                ("📋", "Request form: platform, config, project, lease"),
                ("✅", "Pending approval queue for managers"),
                ("🔍", "Search and filter requests by status"),
                ("👤", "Requester details, timestamp and platform shown"),
                ("✔️", "Approve / Reject with optional comment"),
                ("📧", "Email notification on status change"),
                ("📊", "Export request log to CSV"),
            ]
        ),
        (
            "10_ipam.png",
            "IPAM — IP Address Management",
            "Subnet tracking and per-subnet IP usage visibility",
            [
                ("📡", "List all subnets with prefix, VLAN, usage %"),
                ("🔍", "Drill into any subnet to see all IPs"),
                ("📊", "Per-IP status: Used / Free / Router / Broadcast"),
                ("🔄", "Refresh IPAM cache from backend API"),
                ("📊", "Export subnets and IP lists to CSV"),
            ]
        ),
        (
            "11_assets.png",
            "Asset Inventory Management",
            "Physical and logical asset tracking with ping and power actions",
            [
                ("🖥️", "Full asset list: name, type, IP, location, status"),
                ("🏓", "Ping selected assets for reachability check"),
                ("⚡", "Power On / Off actions on supported assets"),
                ("✏️", "Add, edit and delete asset rows inline"),
                ("📊", "Export full asset inventory to CSV"),
                ("🔍", "Search and filter by name, type, IP"),
            ]
        ),
        (
            "12_openshift.png",
            "Red Hat OpenShift Container Platform",
            "Multi-cluster OCP overview — nodes, operators, namespaces, workloads",
            [
                ("🔴", "Multi-cluster tab navigation"),
                ("📊", "Cluster health: nodes, CPU %, memory %"),
                ("📦", "Operator status: installed and degraded count"),
                ("🗂️", "Namespace list with workload summary"),
                ("🚀", "Pod / Deployment visibility per namespace"),
                ("📋", "OCP cluster version and API endpoint shown"),
            ]
        ),
        (
            "13_nutanix.png",
            "Nutanix AHV Virtualisation",
            "Prism Central integration — VMs, clusters, images",
            [
                ("🟩", "Multi Prism Central tab navigation"),
                ("📊", "AHV VM inventory: name, power state, CPU/RAM"),
                ("🗄️", "Cluster capacity: storage, CPU, memory"),
                ("🖼️", "Image catalogue listing"),
                ("🏷️", "Project and category tagging visible"),
                ("📋", "VM details: host, IP, OS image used"),
            ]
        ),
        (
            "14_ansible.png",
            "Ansible Automation Platform (AAP)",
            "Job template and playbook execution via AAP API",
            [
                ("🤖", "Job template listing from AAP"),
                ("▶️", "Launch job templates with extra vars"),
                ("📊", "Job status: pending, running, success, failed"),
                ("📋", "Job output log viewer"),
                ("🔌", "Direct AAP REST API integration"),
            ]
        ),
        (
            "15_addns.png",
            "Active Directory & DNS Management",
            "Full AD user/group/OU management and DNS zone/record control",
            [
                ("👤", "AD user list: search, enable/disable, unlock"),
                ("🔑", "Reset AD user password from portal"),
                ("👥", "Group management: create, add/remove members"),
                ("📁", "OU and computer listing"),
                ("🌐", "DNS zones list with record count"),
                ("➕", "Add / Delete DNS A, CNAME, PTR records"),
                ("🔄", "Flush DNS server cache"),
            ]
        ),
    ]

    for ss_file, title, subtitle, bullets in screen_slides:
        s = prs.slides.add_slide(LY_CONTENT)
        screenshot_slide(s, ss_file, title, subtitle, bullets)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 19 – USER MANAGEMENT & RBAC (two screenshots)
    # ══════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(LY_CONTENT)
    slide_header(s, "User Management & Role-Based Access Control",
                 "Portal user management and RBAC model")

    # Two screenshots side by side: audit + users
    content_top = CY + Inches(1.08)
    content_h   = FB - content_top - Inches(0.1)
    ss_w = Inches(6.1); ss_h = ss_w * 900 / 1600
    for i, (fname, lbl) in enumerate([("16_audit.png","Audit Log"), ("17_users.png","User Management")]):
        sx = CX + i * (ss_w + Inches(0.18))
        sp = os.path.join(SS_DIR, fname)
        if os.path.exists(sp):
            rect(s, sx - Inches(0.02), content_top - Inches(0.02),
                 ss_w + Inches(0.04), ss_h + Inches(0.04),
                 fill=W_LBLUE, line=W_BLUE, lw=0.5)
            s.shapes.add_picture(sp, sx, content_top, ss_w, ss_h)
        tb(s, lbl, sx, content_top + ss_h + Inches(0.08),
           ss_w, Inches(0.28), sz=11, bold=True, color=W_DARK, align=PP_ALIGN.CENTER)

    roles = [
        ("Super Admin", "Full access: all platforms, users, pricing, delete"),
        ("Admin",       "Manage VMs, approve requests, configure settings"),
        ("Operator",    "VM lifecycle ops, snapshot, clone, migrate"),
        ("Requester",   "Submit VM requests; view own resources only"),
        ("Viewer",      "Read-only access to all dashboards and reports"),
    ]
    y0 = content_top + ss_h + Inches(0.42)
    rw = sw = Inches(12.5) / len(roles)
    for i, (role, desc) in enumerate(roles):
        colors = [W_RED, W_BLUE, W_ORANGE, W_GREEN, W_PURPLE]
        rx = CX + i * rw
        rect(s, rx, y0, rw - Inches(0.05), Inches(0.34), fill=colors[i])
        tb(s, role, rx + Inches(0.05), y0 + Inches(0.04),
           rw - Inches(0.1), Inches(0.28), sz=11, bold=True, color=W_WHITE, align=PP_ALIGN.CENTER)
        tb(s, desc, rx + Inches(0.05), y0 + Inches(0.36),
           rw - Inches(0.1), Inches(0.4), sz=9.5, color=W_DARK, align=PP_ALIGN.CENTER, wrap=True)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 20 – THANK YOU
    # ══════════════════════════════════════════════════════════════════════════
    cover(
        "Thank You",
        "LaaS Portal — Lab as a Service  |  Wipro CaaS Platform",
        "Confidential  |  Wipro Limited  |  2026"
    )

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    n    = len(prs.slides)
    size = os.path.getsize(OUTPUT) // 1024
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {n}  |  Size: {size:,} KB")


if __name__ == "__main__":
    build()

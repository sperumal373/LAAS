from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
import zipfile, os
from lxml import etree

pptx_path = r'C:\caas-dashboard\LaaS_Portal_Presentation.pptx'

# 1. All XML well-formed?
xml_errors = []
with zipfile.ZipFile(pptx_path, 'r') as z:
    for name in z.namelist():
        if name.endswith('.xml') or name.endswith('.rels'):
            try:
                etree.fromstring(z.read(name))
            except etree.XMLSyntaxError as e:
                xml_errors.append(f'{name}: {e}')

if xml_errors:
    print('XML ERRORS:', xml_errors)
else:
    print('ALL XML VALID - PowerPoint repair dialog will NOT appear')

# 2. Slide-by-slide
prs = Presentation(pptx_path)
W = prs.slide_width.emu
H = prs.slide_height.emu
total_issues = 0
for i, slide in enumerate(prs.slides):
    logos = sum(1 for s in slide.shapes if s.shape_type == 13)
    overflow = [s.name for s in slide.shapes
                if s.left is not None and s.width is not None
                and (s.left + s.width) > W + Inches(0.15)]
    anim = slide._element.find(qn('p:timing')) is not None
    flag = []
    if logos == 0:
        flag.append('NO-LOGO')
    if overflow:
        flag.append(f'OVERFLOW({len(overflow)}sh)')
    status = 'OK' if not flag else 'ISSUE:' + ','.join(flag)
    total_issues += len(flag)
    print(f'  Slide {i+1:2d}: {status:<35} logos={logos}  {"[ANIM]" if anim else ""}')

print(f'\nTotal issues: {total_issues}')
print(f'File: {pptx_path}  ({os.path.getsize(pptx_path):,} bytes)')
